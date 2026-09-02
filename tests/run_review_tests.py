#!/usr/bin/env python3
"""Tests for the mlos_review presentation builder.

    python3 tests/run_review_tests.py [--prefix NAME] [--only CHECK] [--quiet]

--prefix selects fixtures by leading name, --only selects checks: any
check function whose name contains the given text, so `--only care_days`
runs check_care_days against every fixture and nothing else. A check
whose result later checks consume still runs under --only, muted, since
skipping it would disable the checks that were asked for.

Each check is run isolated. One that raises is recorded as a failure and
the run continues, because a suite built on running every check against
every fixture should not lose twenty-seven fixtures to one KeyError.

Same philosophy as run_tests.R, and deliberately the same shape: a minimal
assertion harness, no test framework dependency, and every check run against
every fixture rather than against one convenient case. The fixtures are the
committed golden bundles in tests/golden/, which between them cover eleven
levels, an unreached median, a single constant-LOS period, an eleven-way tie,
selectively disabled outputs, and stratifiers that are absent altogether. Those
are exactly the shapes that break a deck builder, and none of them are shapes
anyone would think to write a bespoke test for.

Checks come in two kinds. Fixture checks assert invariants over real bundles
and dispatch on what each bundle happens to contain. Synthetic checks build a
small frame by hand and pin down behavior that no fixture happens to exercise
at the boundary, which is the only honest way to test "a selector that cannot
tell the levels apart must abstain".

There are no golden files here. The deck's content spec is still moving, and a
golden of a table that is about to be split in two would generate churn rather
than catch bugs. Add them once the block set settles.
"""

from __future__ import annotations

import itertools
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
from pathlib import Path

import pandas as pd

REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_ROOT))

from mlos_review import blocks, names, salience  # noqa: E402
from mlos_review.blocks import (  # noqa: E402
    FORMATS,
    FULL_TABLE_MEASURES,
    HIGHLIGHT_SHOW_ALL_MAX,
    Format,
    _unified_stays,
    findings_for_stratifier,
    flag_extremes,
    full_table,
    highlights_table,
    requires_full_table,
)
from mlos_review import output, settings as settings_mod  # noqa: E402
from mlos_review.bundle import Bundle  # noqa: E402
from mlos_review.deck import build  # noqa: E402
from mlos_review.names import Vocabulary, capitalize_first  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

# Words allowed to start a label with a capital: acronyms the reader is assumed
# to know, and proper nouns. Everything else must begin lowercase so a label can
# land mid-sentence in a bullet.
ALLOWED_CAPITALIZED = set(names.ACRONYMS.values()) | {
    "Kaplan-Meier",
    "Aalen-Johansen",
}

_state = {"pass": 0, "fail": 0, "skip": 0, "error": 0, "quiet": False,
          "only": None, "matched": False, "muted": False}


# ---------------------------------------------------------------------------
# Harness
# ---------------------------------------------------------------------------


def section(title: str) -> None:
    if _state["muted"]:
        return
    print(f"\n=== {title} ===")


def expect(label: str, ok: bool, detail: str = "") -> None:
    if _state["muted"]:
        return
    if ok:
        _state["pass"] += 1
        if not _state["quiet"]:
            print(f"  [PASS] {label}")
    else:
        _state["fail"] += 1
        print(f"  [FAIL] {label}{': ' + detail if detail else ''}")


def skipped(label: str, reason: str) -> None:
    """A check that could not run, said out loud.

    Counted and reported in the summary rather than passed over, because a
    tool that is missing on this machine is the one case where a green run
    would mean nothing was checked.
    """
    if _state["muted"]:
        return
    _state["skip"] += 1
    print(f"  [SKIP] {label}: {reason}")


def expect_equal(label: str, actual, expected) -> None:
    expect(label, actual == expected, f"expected {expected!r}, got {actual!r}")


def errored(label: str, exc: BaseException) -> None:
    """A check that raised, recorded as a failure rather than ending the run.

    An escaping exception is a bug either in the check or in what it checks,
    and either way it is one result among thousands. Letting it propagate
    would abandon every fixture after this one and print a traceback where
    the summary belongs, so a run that found one broken check would report
    nothing about the parts that were fine. That is the opposite of what a
    suite whose whole design is "every check against every fixture" is for.
    """
    _state["fail"] += 1
    _state["error"] += 1
    reason = "".join(traceback.format_exception_only(type(exc), exc)).strip()
    print(f"  [ERROR] {label}: {reason}")
    for frame in traceback.format_tb(exc.__traceback__):
        for line in frame.rstrip().splitlines():
            print(f"    {line}")


def run_check(fn, *args, required: bool = False):
    """Run one check, isolated from the rest of the suite and subject to --only.

    Returns whatever the check returns, or None if it raised or was filtered
    out. A caller that feeds one check's result to the next therefore has to
    treat None as "no result to check", which is the same thing it means
    after a crash.

    `required` marks a check whose return value later checks need. Under
    --only it still runs, because skipping it would silently disable the
    checks the user actually asked for, but its own output and counters are
    muted so the run reports only what was selected.
    """
    only = _state["only"]
    filtered = only is not None and only not in fn.__name__
    if filtered and not required:
        return None
    if not filtered:
        _state["matched"] = True

    label = fn.__name__
    if args and isinstance(args[0], str):
        label = f"{args[0]}: {label}"

    _state["muted"] = filtered
    try:
        return fn(*args)
    except Exception as exc:  # noqa: BLE001 - any escape is this check failing
        _state["muted"] = False
        errored(label, exc)
        return None
    finally:
        _state["muted"] = False


# ---------------------------------------------------------------------------
# Synthetic checks: behavior at boundaries no fixture pins down
# ---------------------------------------------------------------------------


def check_flag_extremes() -> None:
    section("flags (synthetic)")

    df = pd.DataFrame(
        {
            "spread": [1.0, 5.0, 3.0],
            "all_tied": [2.0, 2.0, 2.0],
            "tied_max": [7.0, 7.0, 1.0],
            "empty": [float("nan")] * 3,
        },
        index=["a", "b", "c"],
    )
    flags = flag_extremes(df)

    expect_equal("highest gets H", flags.loc["b", "spread"], "H")
    expect_equal("lowest gets L", flags.loc["a", "spread"], "L")
    expect_equal("middle gets nothing", flags.loc["c", "spread"], "")

    expect("an all-tied column is F throughout",
           set(flags["all_tied"]) == {"F"}, str(list(flags["all_tied"])))
    expect("an all-tied column has no H or L",
           not {"H", "L"} & set(flags["all_tied"]))

    expect("a shared maximum is shared",
           list(flags["tied_max"]) == ["H", "H", "L"], str(list(flags["tied_max"])))

    expect("an all-NaN column is unflagged",
           set(flags["empty"]) == {""}, str(list(flags["empty"])))

    # The property that makes a combined flag unnecessary: F and the extreme
    # marks are mutually exclusive within a column, so a level is never both
    # "highest" and "everything here is flat".
    for column in flags.columns:
        marks = set(flags[column]) - {""}
        expect(f"F never shares a column with H or L ({column})",
               not ("F" in marks and {"H", "L"} & marks), f"marks={sorted(marks)}")


def check_selector_abstains() -> None:
    """A superlative nothing distinguishes must not put a row on the slide."""
    section("highlight selection (synthetic)")

    # Six levels, arranged so that the three real superlatives all land in the
    # middle of the range and NOTHING legitimately selects G01. The tied column
    # would select it, because idxmax on an all-tied column returns whichever
    # level sorts first, so G01 appearing in the result is proof the selector
    # failed to abstain.
    levels = [f"G{i:02d}" for i in range(1, 7)]
    frame = pd.DataFrame(
        {
            "mean_daily_intakes": [1.0] * 6,                     # all tied
            "expected_census": [2.0, 3.0, 9.0, 4.0, 5.0, 6.0],       # max G03
            "km_restricted_mean": [5.0, 6.0, 7.0, 1.0, 9.0, 4.0],     # min G04, max G05
        },
        index=levels,
    )
    for measure, _ in FULL_TABLE_MEASURES:
        if measure not in frame:
            frame[measure] = [float(i) for i in range(len(levels))]
    table = blocks.Table(df=frame, title="synthetic", formats=dict(FORMATS),
                         flags=flag_extremes(frame))

    # highlights_table takes its levels from the table it is handed, so the
    # bundle is genuinely unused on this path.
    chosen = list(highlights_table(None, "group", table).df.index)

    expect("an all-tied selector contributes no row",
           "G01" not in chosen, f"chosen={chosen}")
    expect_equal("the distinguishing selectors each pick their level",
                 chosen, ["G03", "G04", "G05"])
    expect("selection never exceeds the level count", len(chosen) <= len(levels))

    # Flags must be INHERITED from the full table, never recomputed on the rows
    # that survived selection. No fixture can prove this on its own: the levels
    # a superlative picks usually are the extremes, so recomputing agrees by
    # accident. It only shows up when a shown row is mid-range, which is
    # exactly the case that would mislead a reader.
    highlights = highlights_table(None, "group", table)
    filler = [m for m, _ in FULL_TABLE_MEASURES
              if m not in ("mean_daily_intakes", "expected_census",
                           "km_restricted_mean")][0]
    expect(f"a mid-range shown row carries no flag ({filler})",
           set(highlights.flags[filler]) == {""},
           f"got {list(highlights.flags[filler])}, "
           f"recomputation would have marked one H and one L")
    expect("highlight flags equal the full-table flags on the shown rows",
           highlights.flags.equals(table.flags.loc[chosen]))


def check_tied_extremes() -> None:
    """Ties in prose: named together up to two, silent past that.

    Synthetic because no fixture reaches it and none should have to. These are
    floating-point means over hundreds of stays, so two levels equal to the
    last digit is a thing the data will not produce and the code must still
    handle: the flags share ties in ink, and a sentence beside them naming one
    of two flagged levels would contradict the table it sits under.
    """
    section("tied extremes (synthetic)")
    from mlos_review.blocks import (MAX_TIED_LEVELS, findings_for_aj,
                                    findings_for_stratifier,
                                    extreme_levels, name_levels)

    values = pd.Series([30.0, 12.0, 30.0, 12.0, 20.0],
                       index=["ALPHA", "BRAVO", "ZULU", "YANKEE", "MIKE"])
    expect_equal("both holders of a shared maximum are named",
                 extreme_levels(values, want_max=True), ["ALPHA", "ZULU"])
    expect_equal("holders come back in the frame's level order",
                 extreme_levels(values, want_max=False), ["BRAVO", "YANKEE"])
    expect_equal("a sole holder is still a list of one",
                 extreme_levels(pd.Series([1.0, 2.0], index=["A", "B"]),
                                want_max=True), ["B"])

    crowded = pd.Series([9.0, 9.0, 9.0, 1.0], index=["A", "B", "C", "D"])
    expect_equal(f"past {MAX_TIED_LEVELS} holders the extreme is not named",
                 extreme_levels(crowded, want_max=True), [])
    expect_equal("the other end of the same column still is",
                 extreme_levels(crowded, want_max=False), ["D"])

    flat = pd.Series([5.0, 5.0], index=["A", "B"])
    expect_equal("a flat column names nobody", extreme_levels(flat, True), [])
    expect_equal("an empty column names nobody",
                 extreme_levels(pd.Series([float("nan")], index=["A"]), True), [])

    expect_equal("one level reads as itself", name_levels(["LRG"]), "LRG")
    expect_equal("two read as a pair", name_levels(["LRG", "XL"]), "LRG and XL")

    # The sentences the levels land in, and the verb agreeing with them.
    vocab = Vocabulary({"palette": {"outcome_labels": {
        "type": "named_vector", "names": ["L", "T"],
        "values": ["L community live", "T other live"]}}})

    def spoken(frame):
        table = blocks.Table(df=frame, title="synthetic",
                             flags=flag_extremes(frame))
        return findings_for_aj(table, "animal group", vocab)

    tied = pd.DataFrame(
        {"aj_restricted_mean_L": [30.0, 12.0, 30.0],
         "aj_restricted_mean_T": [40.0, 15.0, 22.0]},
        index=["ALPHA", "BRAVO", "ZULU"])
    said = spoken(tied)[-1]
    expect("tied levels are named together, with a plural verb",
           "ALPHA and ZULU are slowest (30.0 days)" in said, said)
    expect("a sole level keeps its singular verb",
           "BRAVO is fastest (12.0 days)" in said, said)
    # The table flags BOTH tied levels; the sentence must not name one.
    flags = flag_extremes(tied)["aj_restricted_mean_L"]
    expect_equal("the flags and the sentence agree on who is highest",
                 sorted(flags.index[flags == "H"]), ["ALPHA", "ZULU"])
    expect("no sweep while an outcome is tied at that end",
           all("slowest in every outcome type" not in line for line in spoken(tied)),
           f"{spoken(tied)}")

    crowded_end = tied.copy()
    crowded_end.loc["YANKEE"] = [30.0, 19.0]
    said = spoken(crowded_end)[-1]
    expect("an end with too many holders is dropped from its outcome",
           "slowest" not in said.split("For other live")[0], said)
    expect("the other end of that outcome survives",
           "for community live outcomes, BRAVO is fastest" in said, said)

    # The LOS finding agrees with the same rule, verb and all. Its census
    # clause reads "and" when the same levels hold both extremes and "while"
    # when they do not, which is a comparison of the two sets, so a tied pair
    # holding both still reads as one group.
    class _Pooled:
        """Just enough bundle for the share-of-residents clause."""

        def __init__(self, total):
            self.total = total

        def has(self, *path):
            return path == ("strata", "all", "census")

        def stratum(self, stratifier, block):
            return pd.DataFrame({"expected_census": [self.total]}, index=["All"])

    def found(frame):
        return findings_for_stratifier(
            _Pooled(100.0), "group", blocks.Table(df=frame, title="synthetic"),
            "animal group")

    both = pd.DataFrame(
        {"km_restricted_mean": [31.0, 12.0, 31.0],
         "expected_census": [40.0, 20.0, 40.0]}, index=["LRG", "TOY", "XL"])
    expect_equal("tied levels hold the LOS clause together", found(both),
                 ["Among animal group levels, LRG and XL have the longest "
                  "average stay at 31 days, and the largest share of residents "
                  "at 40%."])

    split = both.copy()
    split["expected_census"] = [10.0, 80.0, 10.0]
    expect("a different level holding the census reads as a contrast",
           "while TOY holds the largest share of residents at 80%"
           in found(split)[0], found(split))

    crowded_los = both.copy()
    crowded_los.loc["MED"] = [31.0, 5.0]
    expect_equal("three levels sharing the longest stay conclude nothing",
                 found(crowded_los), [])


def check_cell_formatting() -> None:
    """The number the reader sees, which the format spec alone does not pin."""
    section("cell formatting (synthetic)")
    from mlos_review.render_pptx import _format_cell

    expect_equal("a fraction renders as a percent",
                 _format_cell(0.008, Format(1, percent=True)), "0.8%")
    expect_equal("a percent at one is a hundred",
                 _format_cell(1.0, Format(0, percent=True)), "100%")
    expect_equal("a plain number is not scaled",
                 _format_cell(30.79, Format(1)), "30.8")
    expect_equal("thousands are grouped",
                 _format_cell(9051.8, Format(0)), "9,052")
    expect_equal("a missing value says so",
                 _format_cell(float("nan"), Format(1)), "n/a")
    # The two cases the summary tables added: a column that only some rows
    # have any business filling, and a column whose values were never numbers.
    expect_equal("a structurally empty cell stays empty",
                 _format_cell(float("nan"), Format(1, blank=True)), "")
    expect_equal("text passes through untouched",
                 _format_cell("2019-01-01", Format(0)), "2019-01-01")


def check_settings() -> None:
    """Defaults, coercions, and the refusals that matter."""
    section("settings (synthetic)")
    from mlos_review.settings import Emphasis, SettingsError, from_mapping

    d = from_mapping({})
    expect_equal("default flag style", d.high_low_flag, "MARK")
    expect_equal("default output path", str(d.output_path), "reports/mlos_deck.pptx")
    expect_equal("default emphasis is AUTO", d.emphasis_for("group"), Emphasis())

    picked = from_mapping({"emphasis": {"animal_group": ["LRG", "XL"]}})
    expect_equal("naming levels pins them",
                 picked.emphasis_for("group").levels, ("LRG", "XL"))
    expect("naming levels also forces the stratifier in",
           picked.emphasis_for("group").forced)
    expect("settings keys map to bundle stratifier ids",
           picked.emphasis_for("intake") == Emphasis())

    expect("NEVER suppresses",
           from_mapping({"emphasis": {"period": "NEVER"}}).emphasis_for("period").suppressed)
    expect("YAML booleans are accepted as synonyms",
           from_mapping({"emphasis": {"period": True}}).emphasis_for("period").forced)

    # A settings file is a promise about the output. Silently ignoring part of
    # it is worse than refusing it, so every one of these must raise.
    refusals = [
        ({"aj_covrage": "FULL"}, "a mistyped key"),
        ({"tables": {"high_low_flag": "RAINBOW"}}, "a value outside its enumeration"),
        ({"aj_coverage": "SOME"}, "a bad aj_coverage"),
        ({"emphasis": {"dog_size": "AUTO"}}, "an unknown stratifier"),
        ({"emphasis": {"period": []}}, "an empty level list"),
        ({"output": {"folder": "x"}}, "a mistyped output key"),
        ({"emphasis": {"period": 3}}, "a number where a keyword belongs"),
        ("not a mapping", "a file that is not a mapping"),
        ({"emphasis": "NEVER"}, "a keyword where the emphasis mapping belongs"),
        ({"output": "reports"}, "a bare value where the output mapping belongs"),
    ]
    for data, label in refusals:
        try:
            from_mapping(data)
            expect(f"refuses {label}", False, "accepted it")
        except SettingsError:
            expect(f"refuses {label}", True)

    expect("case and whitespace are forgiven",
           from_mapping({"tables": {"high_low_flag": " color "}}).high_low_flag == "COLOR")


def check_bullet_pagination() -> None:
    """The closing section's page break, on bullets built to order."""
    from mlos_review.render_pptx import (MARGIN, SLIDE_HEIGHT, TITLE_HEIGHT,
                                         bullet_height, bullet_pages)

    section("bullet pagination (synthetic)")
    budget = SLIDE_HEIGHT - (MARGIN + TITLE_HEIGHT) - MARGIN
    expect_equal("no bullets, no pages", bullet_pages([]), [])
    expect_equal("a few short bullets share one page",
                 bullet_pages(["one", "two", "three"]),
                 [["one", "two", "three"]])

    # Enough one-line bullets to overflow a page twice over. Nothing may be
    # dropped, reordered, or duplicated: the closing section is the only place
    # a finding reaches the reader.
    many = [f"finding {i}" for i in range(40)]
    pages = bullet_pages(many)
    expect("a long list breaks into pages", len(pages) > 1, f"got {len(pages)}")
    expect_equal("every bullet survives, in order",
                 [line for page in pages for line in page], many)
    for index, page in enumerate(pages):
        used = sum(bullet_height(line) for line in page)
        expect(f"page {index + 1} fits the slide", used <= budget,
               f"used={used / 914400:.2f}in of {budget / 914400:.2f}in")

    # One bullet longer than a whole page still gets a page rather than being
    # dropped or cut in half.
    giant = "word " * 400
    expect(f"an oversized bullet is oversized", bullet_height(giant) > budget)
    expect_equal("an oversized bullet gets its own page",
                 bullet_pages(["short", giant]), [["short"], [giant]])

    # A lead line eats into the first page and no other, which is the whole
    # reason `reserved` exists: charging every page for it would break each one
    # a bullet early, and charging none of them would run the first page's last
    # line off the bottom under the caveat.
    from mlos_review.deck import AUTOMATION_CAVEAT
    from mlos_review.render_pptx import lead_height

    expect_equal("no lead costs nothing", lead_height(""), 0)
    reserve = lead_height(AUTOMATION_CAVEAT)
    expect("the caveat costs something", reserve > 0)
    led = bullet_pages(many, reserve)
    expect("a lead costs the first page bullets",
           len(led[0]) < len(pages[0]), f"{len(led[0])} vs {len(pages[0])}")
    expect("and costs the pages after it none",
           all(len(page) == len(pages[1]) for page in led[1:-1]),
           f"{[len(page) for page in led]} vs {[len(page) for page in pages]}")
    for index, page in enumerate(led):
        used = sum(bullet_height(line) for line in page)
        used += reserve if index == 0 else 0
        expect(f"led page {index + 1} fits the slide", used <= budget,
               f"used={used / 914400:.2f}in of {budget / 914400:.2f}in")
    expect_equal("and nothing is lost to it",
                 [line for page in led for line in page], many)


def _slide_title(slide) -> str:
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    return texts[0] if texts else ""


def _title_base(slide) -> str:
    """A slide's title up to any colon, which is what a flip-book run shares.

    The competing-risk slides name their outcome after a colon, so the titles
    differ by design across a run that is otherwise held still. Everything
    before the colon is the section.
    """
    return _slide_title(slide).split(":")[0].strip()


def _slide_geometry(slide):
    """Where a rendered slide put its pictures and its table."""
    return ([(s.left, s.top, s.width, s.height)
             for s in slide.shapes if s.shape_type == 13],
            [(s.left, s.top, s.width) for s in slide.shapes if s.has_table])


def _expect_no_table_overlap(case: str, number: int, slide) -> None:
    """No text box may start inside the space a table will actually occupy.

    The bug this pins: a table is added at whatever height the layout asked
    for, but pptx grows every row to fit its text, so the rendered table is as
    deep as its content whatever the shape claims. A footnote positioned from
    an underestimate lands ON the last row of the table. It was found by eye,
    in PowerPoint, on a slide of three tables whose headers had been squeezed
    into wrapping, and nothing in the suite could see it.

    The real depth is counted here the way pptx will lay it out: one row height
    per row, plus an extra for every line a header wraps to. Overlap is judged
    only against boxes that share the table's columns, since a table in the
    left half of a slide has nothing to say about a caption in the right.
    """
    from mlos_review.render_pptx import (TABLE_ROW_HEIGHT, VALUE_PT,
                                         _word_width, _wrapped_lines)

    tables, boxes = [], []
    for shape in slide.shapes:
        if shape.has_table:
            tables.append(shape)
        elif shape.has_text_frame and shape.text_frame.text.strip():
            boxes.append(shape)

    for shape in tables:
        grid = shape.table
        # No header may break inside a word. A column narrower than its
        # longest word has nowhere legal to wrap, so PowerPoint breaks
        # mid-word, one letter to a line, which is what a squeezed "Pct"
        # column looked like. The renderer floors columns at their longest
        # word; this is that floor checked on the file that was written.
        for cell, column in zip(grid.rows[0].cells, grid.columns):
            expect(f"{case}: slide {number} gives every header room for its "
                   f"longest word",
                   not cell.text.strip()
                   or _word_width(cell.text, VALUE_PT) <= column.width,
                   f"{cell.text!r} needs {_word_width(cell.text, VALUE_PT)} "
                   f"in a column of {column.width}")
        header_lines = max(
            [_wrapped_lines(cell.text, VALUE_PT, column.width)
             for cell, column in zip(grid.rows[0].cells, grid.columns)] or [1])
        depth = TABLE_ROW_HEIGHT * (len(grid.rows) - 1 + header_lines)
        for box in boxes:
            # Meaningful horizontal overlap, not any overlap at all. A row of
            # tables puts each footnote at its own table's allotted width,
            # which can reach a few hundredths of an inch into the neighbor's
            # span; that is a sliver of white, not a footnote written over a
            # table, and treating it as one made this check cry wolf on a
            # slide that was correct.
            shared = (min(box.left + box.width, shape.left + shape.width)
                      - max(box.left, shape.left))
            overlaps_columns = shared > min(box.width, shape.width) / 4
            if overlaps_columns and shape.top < box.top < shape.top + depth:
                expect(f"{case}: slide {number} keeps its text clear of its "
                       f"tables", False,
                       f"a box at {box.top} inside a table running "
                       f"{shape.top} to {shape.top + depth}: "
                       f"{box.text_frame.text[:60]!r}")
                return
        expect(f"{case}: slide {number} keeps its text clear of its tables", True)


def _stage_interval_figures(staged: Path, bundle: Bundle) -> list[str]:
    """Add the two interval stacks to a staged fixture's manifest.

    No golden run sets `probability_mass_width`, so no fixture names those two
    figures and the educational slide would return None on every one of them:
    a rule that never builds is a rule nothing checks. Adding the entries to
    the COPY gives it the manifest a run with the setting on would have, and
    the placeholder PNGs beside it, without regenerating twenty-eight goldens
    to exercise one slide.

    Only the manifest is touched. The slide reads its numbers from the tables
    every other whole-sample slide uses, so nothing here fakes a result.
    """
    if not bundle.data.get("aj", {}).get("has_analysis"):
        return []
    results = staged / "results.json"
    data = json.loads(results.read_text())
    added = []
    for kind in ("aj_mass", "aj_fraction"):
        plot = f"{kind}_unified_stack.png"
        data.setdefault("outputs", []).append({
            "csv": None, "plot": plot, "kind": kind, "stratifier": "all",
            "outcome": None, "variant": "stack",
            "description": "staged for the deck geometry checks",
        })
        added.append(plot)
    results.write_text(json.dumps(data))
    return added


def check_deck_with_figures(case: str, bundle: Bundle, directory: Path) -> None:
    """The half of the deck the golden bundles cannot reach on their own.

    Golden fixtures ship results.json and the companion CSVs but no PNGs, which
    is right: images are large, and the R suite already checks that every plot
    named in the manifest was written. The cost is that `bundle.figure` returns
    None for every lookup here, so until this check existed, every rule ran
    down its no-figures fallback and no layout was ever exercised against a
    real bundle.

    So the figures are stood up: a copy of the fixture, plus a 3:2 placeholder
    PNG at every path the manifest names. The images carry no data and are not
    meant to; what is under test is the geometry the deck computes around them,
    and a 300x200 white rectangle has the same aspect ratio and the same header
    the renderer reads as the figure R would have drawn.
    """
    from mlos_review.names import Vocabulary as _Vocab
    try:
        from PIL import Image
    except ImportError:  # pragma: no cover - Pillow ships with python-pptx
        return

    vocab = _Vocab(bundle.data)
    plots = [e["plot"] for e in bundle.data.get("outputs", []) if e.get("plot")]
    if not plots:
        return

    with tempfile.TemporaryDirectory() as tmp:
        staged = Path(tmp) / "fixture"
        shutil.copytree(directory, staged)
        for name in plots + _stage_interval_figures(staged, bundle):
            Image.new("RGB", (300, 200), "white").save(staged / name)

        # Anything that depends on a figure EXISTING has to be asserted here
        # rather than against the bare fixture, where it would pass vacuously.
        # Which column a competing-risk slide highlights is the case in point:
        # there is one slide per drawn panel, so with no panels there is
        # nothing to check and a broken highlight goes unnoticed.
        from mlos_review.blocks import requires_aj
        from mlos_review.deck import _aj_figures, aj_by_stratifier

        staged_bundle = Bundle.load(staged)
        for stratifier in staged_bundle.stratifiers():
            if stratifier == "all" or not requires_aj(staged_bundle, stratifier):
                continue
            panels = _aj_figures(staged_bundle, stratifier)
            section = aj_by_stratifier(staged_bundle, stratifier, vocab)
            # A stratifier whose plots this run did not draw still gets one
            # slide, carrying the table alone and highlighting nothing.
            expect_equal(f"{case}: a slide per drawn panel ({stratifier})",
                         len(section), len(panels) or 1)
            expect_equal(f"{case}: each slide highlights its own outcome "
                         f"({stratifier})",
                         [s.table.highlight for s in section],
                         [code for code, _ in panels] or [""])

        out = Path(tmp) / "deck.pptx"
        build(staged, out)
        deck = Presentation(str(out))

        # Where the outlook slide sits, checked on the rendered file because
        # position is the requirement: it carries the previous slide's rows and
        # its tenure figure forward, and a reader meets it having just seen
        # them. It needs figures to exist, so this is the only place it can be
        # asserted end to end.
        from mlos_review.settings import Settings as _Settings
        emphasis = _Settings().for_dataset(
            [s for s in staged_bundle.stratifiers() if s != "all"])
        titles = [_slide_title(slide) for slide in deck.slides]
        for stratifier in staged_bundle.stratifiers():
            if stratifier == "all" or not emphasis.emphasis_for(stratifier).forced:
                continue
            label = vocab.stratifier(stratifier).label
            outlook = f"What an average day's residents still owe, by {label}"
            if outlook not in titles:
                continue
            expect(f"{case}: the {stratifier} outlook slide follows its LOS "
                   f"slide", titles.index(outlook) == titles.index(f"LOS by {label}") + 1,
                   f"{titles}")

        # The educational slide, on every fixture whose staged manifest names
        # its two figures. Its shape is what is asserted: two figures side by
        # side over a row of two tables is the one place STACKED draws a row
        # rather than a single table, and the geometry loop below is what
        # catches the row colliding with the figures above it.
        if staged_bundle.figure("aj_mass", "all", variant="stack") is not None:
            interval = "Working with Probabilities in Time Intervals"
            expect(f"{case}: the educational slide is built", interval in titles,
                   f"{titles[-4:]}")
            if interval in titles:
                at = titles.index(interval)
                expect(f"{case}: the educational section sits at the back",
                       at == len(titles) - 1, f"{titles[at:]}")
                expect_equal(f"{case}: the section runs divider, opener, "
                             f"intervals",
                             titles[-3:],
                             ["Educational", "Looking at Length of Stay (LOS)",
                              interval])
                page = list(deck.slides)[at]
                pictures = [sh for sh in page.shapes if sh.shape_type == 13]
                grids = [sh for sh in page.shapes if sh.has_table]
                expect_equal(f"{case}: it carries both figures", len(pictures), 2)
                expect_equal(f"{case}: it carries both tables", len(grids), 2)
                expect(f"{case}: its figures are side by side",
                       len({sh.top for sh in pictures}) == 1
                       and len({sh.left for sh in pictures}) == 2,
                       f"{[(sh.left, sh.top) for sh in pictures]}")
                expect(f"{case}: its tables sit below its figures",
                       min(sh.top for sh in grids)
                       >= max(sh.top + sh.height for sh in pictures),
                       f"figures to {max(sh.top + sh.height for sh in pictures)}, "
                       f"tables from {min(sh.top for sh in grids)}")

        drawn = sum(1 for slide in deck.slides
                    for shape in slide.shapes if shape.shape_type == 13)
        expect(f"{case}: the deck draws figures once they exist", drawn > 0,
               f"{len(plots)} plots available, {drawn} placed")
        for index, slide in enumerate(deck.slides):
            for shape in slide.shapes:
                expect(f"{case}: slide {index + 1} keeps every shape on the page",
                       shape.left >= 0 and shape.top >= 0
                       and shape.left + shape.width <= Inches(13.333)
                       and shape.top + shape.height <= Inches(7.5),
                       f"{shape.shape_type} at {shape.left}, {shape.top}, "
                       f"{shape.width}x{shape.height}")
            _expect_no_table_overlap(case, index + 1, slide)

        # The flip-book property, checked on the rendered file because that is
        # where it is either true or not: consecutive slides sharing a title
        # must put their picture and their table in exactly the same place, or
        # clicking through them jumps instead of animating.
        run = []
        for slide in list(deck.slides) + [None]:
            if run and (slide is None or _title_base(slide) != _title_base(run[0])):
                if len(run) > 1:
                    name = _title_base(run[0])[:30]
                    pictures = [_slide_geometry(s)[0] for s in run]
                    tables = [_slide_geometry(s)[1] for s in run]
                    expect(f"{case}: '{name}' holds its figure still across "
                           f"{len(run)} slides",
                           all(p == pictures[0] for p in pictures), f"{pictures}")
                    expect(f"{case}: '{name}' holds its table still across "
                           f"{len(run)} slides",
                           all(t == tables[0] for t in tables), f"{tables}")
                run = []
            if slide is not None:
                run.append(slide)


# What each layout is for, and so what it must draw when it is handed one of
# everything: how many figures, and how many tables. The layouts that hold a
# table beside figures take it in `table` and place exactly one; TABLES takes
# `tables` and places all of them; TITLE draws no figure and puts its one table
# under the bullets.
#
# A layout added to LAYOUT_FUNCTIONS and not listed here fails the coverage
# assertion below, which is the point: an untested layout is the failure mode
# QUADRANTS was already one release away from.
LAYOUT_CONTRACTS = {
    "STACKED": (True, 1),
    "SPLIT": (True, 1),
    "QUADRANTS": (True, 1),
    "TABLES": (False, 3),
    "TITLE": (False, 1),
}


def check_educational_section() -> None:
    """The section's standing rules, which no single slide rule can hold.

    Two of them. The order, because the opener asks the question the second
    slide answers and a section that led with the answer would be a different
    section. And the guideline: an educational slide contributes no finding and
    no recommendation, which `educational_section` enforces rather than trusts,
    so a rule added later cannot leak a sentence about method into a summary of
    sentences about the shelter.
    """
    from mlos_review.deck import educational_section

    section("educational section (synthetic)")
    bundle = Bundle.load(sorted((REPO_ROOT / "tests" / "golden").glob("*/results.json"))[0].parent)
    vocab = Vocabulary(bundle.data)

    slides = educational_section(bundle, vocab)
    expect(f"the section is built", bool(slides))
    titles = [slide.title for slide in slides]
    expect_equal("it opens with its divider", titles[0], "Educational")
    expect_equal("the LOS opener comes first",
                 titles[1], "Looking at Length of Stay (LOS)")
    # Carried on the golden bundles, which name no interval figure, so the
    # opener is the whole section there and the rule below is asserted over
    # whatever the section did build.
    for slide in slides:
        expect_equal(f"{slide.title!r} carries no findings", slide.findings, [])
        expect_equal(f"{slide.title!r} carries no recommendations",
                     slide.recommendations, [])

    # A rule that returns findings has them dropped, which is the guarantee
    # rather than every rule remembering. Asserted by making one do it.
    import mlos_review.deck as deck_module
    original = deck_module.los_overview_slide

    def _noisy(bundle, vocab):
        slide = original(bundle, vocab)
        slide.findings = ["a finding a teaching slide should not carry"]
        slide.recommendations = ["nor this"]
        return slide

    deck_module.los_overview_slide = _noisy
    try:
        noisy = deck_module.educational_section(bundle, vocab)
    finally:
        deck_module.los_overview_slide = original
    expect("a rule's findings are dropped rather than gathered",
           all(not slide.findings and not slide.recommendations
               for slide in noisy),
           f"{[(s.title, s.findings, s.recommendations) for s in noisy]}")


def check_slide_citations() -> None:
    """The two DOIs the opener prints, against the files that own them.

    A slide is the one place a citation is read aloud and the one place nothing
    resolves it, so both are held against their source: the software DOI
    against CITATION.cff, which is what Zenodo and GitHub read, and the paper
    against the user guide's reference list. A DOI that drifts here names
    something that does not exist and nothing else in the suite would notice.
    """
    from mlos_review.deck import PAPER_CITATION, SOFTWARE_DOI

    section("slide citations (synthetic)")
    citation = (REPO_ROOT / "CITATION.cff").read_text(encoding="utf-8")
    expect(f"the software DOI is the one in CITATION.cff",
           f'doi: "{SOFTWARE_DOI}"' in citation, SOFTWARE_DOI)

    paper_doi = PAPER_CITATION.split("doi:")[-1].strip()
    guide = (REPO_ROOT / "mlos_user_guide.md").read_text(encoding="utf-8")
    expect("the paper DOI is in the user guide's references",
           f"doi:{paper_doi}" in guide, paper_doi)
    # The rest of the line, so a citation cannot keep a live DOI while naming
    # the wrong journal or year.
    for part in ("Mavrovouniotis ML.", "PLOS ONE. 2026;21(1):e0342102"):
        expect(f"the guide carries {part!r}", part in guide)


def check_layouts_render() -> None:
    """Every layout in LAYOUT_FUNCTIONS draws, including the one no rule uses.

    QUADRANTS is kept for a four-figure slide nobody has written yet. Kept
    code that nothing calls and nothing tests is kept code that has quietly
    stopped working, so it is exercised here: a real figure from the results
    directory if there is one, and geometry assertions either way.
    """
    from mlos_review.render_pptx import LAYOUT_FUNCTIONS, Slide, render

    section("layouts (synthetic)")
    expect_equal("every layout is under contract",
                 sorted(LAYOUT_CONTRACTS), sorted(LAYOUT_FUNCTIONS))

    figures = sorted((REPO_ROOT / "results").glob("*.png"))[:4]
    frame = pd.DataFrame({"km_median_los": [4.0, 9.0], "km_p90_los": [12.0, 30.0]},
                         index=["one", "two"])
    table = blocks.Table(df=frame, title="synthetic", formats=dict(FORMATS),
                         flags=flag_extremes(frame), footnotes=["a footnote"])
    vocab = Vocabulary({})

    with tempfile.TemporaryDirectory() as tmp:
        for layout, (takes_figures, table_count) in LAYOUT_CONTRACTS.items():
            wanted = list(figures) if takes_figures else []
            out = Path(tmp) / f"{layout}.pptx"
            render([Slide(title=f"{layout} layout", figures=wanted,
                          table=table if table_count == 1 else None,
                          tables=[table] * table_count if table_count > 1 else [],
                          bullets=["a bullet that says something"],
                          footnote="drawn at a time", notes=["note"],
                          layout=layout)],
                   out, vocab)
            expect(f"{layout} renders", out.exists())
            deck = Presentation(str(out))
            shapes = deck.slides[0].shapes
            pictures = [s for s in shapes if s.shape_type == 13]
            expect_equal(f"{layout} draws every figure it was given",
                         len(pictures), len(wanted))
            table_shapes = [s for s in shapes if s.has_table]
            expect_equal(f"{layout} draws every table it was given",
                         len(table_shapes), table_count)
            # Bands, not rules, and never on the header. The stacked
            # competing-risk table leans on this to tell its share rows from
            # its days rows, so a change that dropped it would be silent.
            for grid in (s.table for s in table_shapes):
                shaded = [i for i, row in enumerate(grid.rows)
                          if row.cells[0].fill.type == 1]
                expect_equal(f"{layout} washes every other data row",
                             shaded, list(range(1, len(grid.rows), 2)))
            expect(f"{layout} draws the slide footnote",
                   any(s.has_text_frame and "drawn at a time" in s.text_frame.text
                       for s in shapes))
            expect(f"{layout} keeps every shape on the slide",
                   all(s.left >= 0 and s.top >= 0
                       and s.left + s.width <= Inches(13.333)
                       and s.top + s.height <= Inches(7.5) for s in shapes),
                   f"{[(s.left, s.top, s.width, s.height) for s in shapes]}")
            _expect_no_table_overlap(layout, 1, deck.slides[0])

    # A row of tables that cannot fit, which is the shape that broke: three
    # tables of long headings on a slide with room for about ten columns. What
    # is under test is the FLOOR, that no column is drawn narrower than its
    # longest word, and that the footnote still lands under the table rather
    # than on it. Synthetic because no fixture is wide enough to force it: the
    # golden bundles carry two or three levels, and the squeeze needs a real
    # dataset's worth of them.
    crowded = pd.DataFrame(
        {"km_median_los": [4.0, 9.0, 14.0], "km_p90_los": [12.0, 30.0, 44.0],
         "expected_census": [120.0, 33.0, 12.0]},
        index=["_UNKNOWN_", "SOMETHING_LONG", "ANOTHER_ONE"])
    wide = blocks.Table(df=crowded, title="a title of some length",
                        formats=dict(FORMATS), footnotes=["a footnote"])
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "crowded.pptx"
        render([Slide(title="three crowded tables", tables=[wide] * 3,
                      layout="TABLES")], out, vocab)
        deck = Presentation(str(out))
        _expect_no_table_overlap("crowded", 1, deck.slides[0])

    # And the contract when a row genuinely cannot fit: the columns keep their
    # floors and the table hangs off the slide, where the geometry check and
    # the eye can both see it. Silently crushing the columns instead is the
    # failure this pins, and it is the one that shipped.
    from mlos_review.render_pptx import _column_minimums, _fitted_widths
    impossible = blocks.Table(
        df=crowded.rename(columns={"km_median_los": "km_restricted_mean"}),
        title="t", formats=dict(FORMATS))
    floors = _column_minimums(impossible, vocab)
    expect_equal("an impossible row keeps its column floors",
                 _fitted_widths(impossible, vocab, "MARK", int(sum(floors) / 2)),
                 floors)


def check_capital_widths() -> None:
    """Capitals are charged more, because level names are written in them.

    The bug this pins: "OWNER" was budgeted the same width as five lowercase
    characters and wrapped a letter onto a second line. Measured against Arial,
    whose capitals are wider than the Calibri a deck asks for, so the estimate
    errs the safe way.
    """
    from mlos_review.render_pptx import (CELL_MARGIN, CELL_PADDING, VALUE_PT,
                                         _text_width)

    section("column widths (synthetic)")
    expect("capitals cost more than lowercase",
           _text_width("OWNER", VALUE_PT) > _text_width("owner", VALUE_PT))
    expect("a digit is not a capital",
           _text_width("80.8%", VALUE_PT) == _text_width("aa.a%", VALUE_PT))

    # The real font, for the labels that actually wrapped.
    try:
        from fontTools.ttLib import TTFont
        font = TTFont("/System/Library/Fonts/Supplemental/Arial.ttf")
    except Exception:  # pragma: no cover - fontTools or Arial absent
        return
    upm, hmtx, cmap = font["head"].unitsPerEm, font["hmtx"], font.getBestCmap()
    for label in ("OWNER", "PUPPY", "SMALL", "STRAY", "OTHER", "Early-C"):
        needed = sum(hmtx[cmap[ord(c)]][0] / upm for c in label) * VALUE_PT.pt / 72
        needed = needed * 914400 + 2 * CELL_MARGIN
        expect(f"{label} fits the column it is budgeted",
               _text_width(label, VALUE_PT) + CELL_PADDING >= needed,
               f"needs {needed / 914400:.3f}in, budgeted "
               f"{(_text_width(label, VALUE_PT) + CELL_PADDING) / 914400:.3f}in")


def check_report_archiving() -> None:
    """A report is never overwritten, and never deleted."""
    section("report archiving (synthetic)")
    from datetime import date

    with tempfile.TemporaryDirectory() as tmp:
        target = Path(tmp) / "reports" / "mlos_deck.pptx"
        day = date(2026, 7, 26)

        path, archived = output.prepare_output(target, today=day)
        expect("the report directory is created", path.parent.is_dir())
        expect_equal("nothing to archive on a first run", archived, None)

        path.write_text("first")
        _, archived = output.prepare_output(target, today=day)
        expect_equal("the previous deck is renamed with date and ordinal",
                     archived.name, "mlos_deck_20260726_001.pptx")
        expect("the previous deck still exists", archived.exists())
        expect_equal("and keeps its contents", archived.read_text(), "first")
        expect("the target name is free again", not target.exists())

        path.write_text("second")
        _, archived = output.prepare_output(target, today=day)
        expect_equal("the ordinal counts within the date",
                     archived.name, "mlos_deck_20260726_002.pptx")

        path.write_text("third")
        _, archived = output.prepare_output(target, today=date(2026, 7, 27))
        expect_equal("a new date restarts the ordinal",
                     archived.name, "mlos_deck_20260727_001.pptx")

        kept = sorted(q.name for q in target.parent.iterdir())
        expect_equal("every earlier deck is still there", len(kept), 3)


def check_aj_teaser(case: str, bundle: Bundle) -> None:
    """The teaser slide, where the bundle has competing risks to teach with."""
    from mlos_review.blocks import (aj_outcome_table, aj_teaser_table,
                                    requires_aj_teaser)
    from mlos_review.deck import aj_teaser
    from mlos_review.names import Vocabulary as _Vocab

    vocab = _Vocab(bundle.data)
    if not requires_aj_teaser(bundle):
        expect(f"{case}: no teaser without a competing-risks analysis",
               aj_teaser(bundle, vocab) is None)
        return

    upright = aj_outcome_table(bundle, vocab)
    expect(f"{case}: outcome rows are the real outcomes, not Any",
           all("any" not in str(i).lower() for i in upright.df.index),
           f"rows={list(upright.df.index)}")
    expect(f"{case}: outcome rows are labeled, not coded",
           len(upright.df) > 0
           and all(str(i) not in ("L", "T", "N") for i in upright.df.index),
           f"rows={list(upright.df.index)}")
    # A cumulative incidence that reaches every stay comes back as 1 plus a
    # rounding error, so the bound is 1 + epsilon rather than 1. Anything
    # meaningfully above that would mean the column is already a percentage.
    expect(f"{case}: shares are fractions rendered as percents",
           bool((upright.df["aj_final_cif"].dropna() <= 1.0 + 1e-9).all())
           and upright.formats["aj_final_cif"].percent,
           f"max={upright.df['aj_final_cif'].max()!r}")
    expect(f"{case}: the upright table carries no flags", upright.flags is None)

    # The flat wrapper is the same numbers rearranged, so it holds every value
    # the upright one does and nothing else.
    table = aj_teaser_table(bundle, vocab)
    expect_equal(f"{case}: the teaser table is one row", len(table.df), 1)
    expect_equal(f"{case}: one column per outcome per measure",
                 len(table.df.columns), 2 * len(upright.df))
    expect(f"{case}: shares lead, mean days follow",
           [c.rsplit("_", 1)[0] for c in table.df.columns]
           == ["aj_final_cif"] * len(upright.df)
           + ["aj_restricted_mean"] * len(upright.df),
           f"columns={list(table.df.columns)}")
    # Each header names its outcome AND which of the two measures it holds,
    # which is what lets six columns be read without a footnote about ordering.
    expect(f"{case}: every column is headed by its outcome and its measure",
           all(table.headers[c]
               == f"{c.rsplit('_', 1)[1]} "
                  f"{vocab.metric(c.rsplit('_', 1)[0]).text(short=True)}"
               for c in table.df.columns),
           f"headers={table.headers}")
    expect(f"{case}: the outcome order repeats across the halves",
           [table.headers[c].split(" ")[0] for c in table.df.columns[:len(upright.df)]]
           == [table.headers[c].split(" ")[0]
               for c in table.df.columns[len(upright.df):]])
    expect(f"{case}: the two halves are told apart by their headers, not a note",
           len({table.headers[c].split(" ", 1)[1] for c in table.df.columns}) == 2,
           f"headers={table.headers}")
    flat = sorted(table.df.iloc[0].dropna().round(9))
    upright_values = sorted(upright.df.to_numpy().ravel().round(9))
    expect(f"{case}: the flat table holds the upright one's numbers",
           flat == [v for v in upright_values if v == v],
           f"flat={flat} upright={upright_values}")
    expect(f"{case}: the teaser carries no flags", table.flags is None)

    # The two figures are two DIFFERENT drawings of two different analyses;
    # without the manifest variant, asking for two stacks could silently return
    # the same file twice.
    slide = aj_teaser(bundle, vocab)
    expect(f"{case}: teaser slide built", slide is not None)
    if slide is not None and len(slide.figures) == 2:
        expect(f"{case}: the teaser's two figures are distinct files",
               slide.figures[0] != slide.figures[1],
               f"figures={[f.name for f in slide.figures]}")
        expect(f"{case}: both teaser figures are stacked",
               all(f.name.endswith("_stack.png") for f in slide.figures),
               f"figures={[f.name for f in slide.figures]}")


def check_aj_by_stratifier(case: str, bundle: Bundle, stratifier: str) -> None:
    """The full competing-risks block for one stratifier."""
    from mlos_review.blocks import (AJ_SHOW_ALL_MAX, aj_highlights_table,
                                    aj_levels_table, aj_outcome_codes,
                                    aj_teaser_table, findings_for_aj,
                                    flag_extremes, requires_aj)
    from mlos_review.deck import _aj_figures, aj_by_stratifier
    from mlos_review.names import Vocabulary as _Vocab

    vocab = _Vocab(bundle.data)
    if not requires_aj(bundle, stratifier):
        expect(f"{case}: no AJ slides for {stratifier} without AJ numbers",
               aj_by_stratifier(bundle, stratifier, vocab) == [])
        return

    codes = aj_outcome_codes(bundle, stratifier)
    levels = aj_levels_table(bundle, stratifier, vocab)
    expect_equal(f"{case}: one row per level ({stratifier})",
                 list(levels.df.index), bundle.levels(stratifier))
    expect_equal(f"{case}: one column per outcome per measure ({stratifier})",
                 len(levels.df.columns), 2 * len(codes))
    expect(f"{case}: shares lead, mean days follow ({stratifier})",
           [c.rsplit("_", 1)[0] for c in levels.df.columns]
           == ["aj_final_cif"] * len(codes) + ["aj_restricted_mean"] * len(codes),
           f"columns={list(levels.df.columns)}")

    # The same block the teaser uses, so the two cannot drift. The teaser is
    # this table asked for the single-level stratifier.
    if stratifier == "all":
        expect_equal(f"{case}: the teaser table is this block on `all`",
                     list(aj_teaser_table(bundle, vocab).df.columns),
                     list(levels.df.columns))

    multi = len(levels.df) > 1
    expect_equal(f"{case}: flags appear only with levels to compare ({stratifier})",
                 levels.flags is not None, multi)
    if multi:
        expect(f"{case}: flags are the extremes over every level ({stratifier})",
               levels.flags.equals(flag_extremes(levels.df)))

    highlights = aj_highlights_table(bundle, stratifier, vocab, levels)
    shown = list(highlights.df.index)
    expect(f"{case}: highlights are a subset in canonical order ({stratifier})",
           shown == [lvl for lvl in levels.df.index if lvl in shown],
           f"shown={shown}")
    if len(levels.df) <= AJ_SHOW_ALL_MAX:
        expect_equal(f"{case}: {AJ_SHOW_ALL_MAX} levels or fewer show all "
                     f"({stratifier})", shown, list(levels.df.index))
    else:
        # Every selected row has to be there because it holds the longest or
        # the shortest mean time to some outcome. Nothing else selects.
        times = [c for c in levels.df.columns
                 if c.startswith("aj_restricted_mean_")]
        extremes = set()
        for column in times:
            values = levels.df[column].dropna()
            if values.empty or values.min() == values.max():
                continue
            extremes |= {values.idxmax(), values.idxmin()}
        expect_equal(f"{case}: selection is by mean time to event ({stratifier})",
                     set(shown), extremes)
    if levels.flags is not None:
        expect(f"{case}: highlight flags are the full table's ({stratifier})",
               highlights.flags.equals(levels.flags.loc[shown]))

    # Findings are judged over every level, so they can name a level the
    # highlights table did not show; what they must never do is disagree with
    # the flags, which mark the same extremes of the same columns.
    findings = findings_for_aj(levels, stratifier, vocab)
    expect(f"{case}: at most two AJ findings per slide ({stratifier})",
           len(findings) <= 2, f"got {len(findings)}")

    # Reconstructed from the frame rather than from the block: who is slowest
    # and fastest to each outcome, and whether any outcome had a tie at either
    # end, which is what the sweep is not allowed to talk over.
    said, ties = [], False
    for column in levels.df.columns:
        if not column.startswith("aj_restricted_mean_"):
            continue
        values = levels.df[column].dropna()
        if len(levels.df) < 2 or values.empty or values.min() == values.max():
            ties = True   # an outcome with no extreme to hold blocks the sweep
            continue
        said.append((values.idxmax(), values.idxmin()))
        ties = ties or (values == values.max()).sum() > 1 \
            or (values == values.min()).sum() > 1

    detail = findings[-1] if findings else ""
    expect_equal(f"{case}: a flat or single level finds nothing "
                 f"({stratifier})", bool(findings), bool(said))
    for slowest, fastest in said:
        expect(f"{case}: the finding names the slowest and fastest levels "
               f"({stratifier})",
               str(slowest) in detail and str(fastest) in detail,
               f"want {slowest}/{fastest} in {detail!r}")

    # The sweep: a level at the same end of EVERY outcome, which is the only
    # thing that may lead the stratifier's findings and must abstain wherever
    # an outcome ties or there is only one of them.
    slowest = {pair[0] for pair in said}
    fastest = {pair[1] for pair in said}
    sweeps = (not ties and len(said) > 1
              and (len(slowest) == 1 or len(fastest) == 1))
    expect_equal(f"{case}: a sweep leads only when one level holds every "
                 f"extreme ({stratifier})", len(findings) == 2, bool(sweeps))
    if sweeps:
        lead = findings[0]
        expect(f"{case}: the sweep names the stratifier ({stratifier})",
               str(stratifier) in lead, lead)
        if len(slowest) == 1:
            expect(f"{case}: the sweep names the level slowest to every "
                   f"outcome ({stratifier})",
                   f"{next(iter(slowest))} is the slowest in every outcome type"
                   in lead, lead)
        if len(fastest) == 1:
            expect(f"{case}: the sweep names the level fastest to every "
                   f"outcome ({stratifier})",
                   f"{next(iter(fastest))} is the fastest in every outcome type"
                   in lead, lead)

    # The section is a flip-book: one slide per outcome that HAS a figure, and
    # everything except the figure identical across them, which is what makes
    # the flip read as the panel changing rather than the slide changing.
    section = aj_by_stratifier(bundle, stratifier, vocab)
    expect(f"{case}: AJ section built for {stratifier}", bool(section))
    if not section:
        return
    drawn = [c for c in codes if bundle.figure("aj_cif", stratifier, outcome=c)]
    expect_equal(f"{case}: one slide per drawn outcome ({stratifier})",
                 len(section), len(drawn) or 1)
    expect_equal(f"{case}: every slide in the section shares one title base "
                 f"({stratifier})",
                 len({s.title.split(":")[0] for s in section}), 1)
    # The table is the stacked arrangement, so its rows are (level, measure)
    # pairs; the levels are what must match the selection, in order.
    expect(f"{case}: every slide shows the same levels ({stratifier})",
           all(list(dict.fromkeys(row[0] for row in s.table.df.index)) == shown
               for s in section), f"want {shown}")
    expect(f"{case}: outcomes are the columns ({stratifier})",
           all(list(s.table.df.columns) == codes for s in section),
           f"want {codes}")
    expect(f"{case}: each slide highlights its own outcome ({stratifier})",
           all(s.table.highlight == code
               for (code, _), s in zip(_aj_figures(bundle, stratifier), section)),
           f"got {[s.table.highlight for s in section]}")
    expect(f"{case}: the highlight is a column, not a mark in the header "
           f"({stratifier})",
           all(set(s.table.headers.values()) == set(codes) for s in section),
           f"got {[s.table.headers for s in section]}")
    expect(f"{case}: every slide has notes ({stratifier})",
           all(s.notes for s in section))
    expect(f"{case}: no slide displays its own finding ({stratifier})",
           all(not s.bullets for s in section))

    # Findings about the stratifier, not one per outcome, so they ride on the
    # first slide alone or the closing section would print them three times.
    # The spread finding follows the timings it summarizes.
    label = vocab.stratifier(stratifier).label
    expect_equal(f"{case}: the section finds once ({stratifier})",
                 [line for s in section for line in s.findings],
                 findings_for_aj(levels, label, vocab)
                 + blocks.findings_for_outcome_spread(bundle, stratifier,
                                                      label, vocab)
                 + blocks.findings_for_outcome_contrast(bundle, stratifier,
                                                        label, vocab))
    expect(f"{case}: the finding is on the first slide ({stratifier})",
           all(not s.findings for s in section[1:]))

    for index, slide in enumerate(section):
        expect_equal(f"{case}: slide {index + 1} holds one figure ({stratifier})",
                     len(slide.figures), 1 if drawn else 0)
        expect_equal(f"{case}: figures put the slide beside its table "
                     f"({stratifier})", slide.layout,
                     "SPLIT" if drawn else "STACKED")
    if drawn:
        expect_equal(f"{case}: the panels are the outcomes in order "
                     f"({stratifier})",
                     [s.figures[0] for s in section],
                     [bundle.figure("aj_cif", stratifier, outcome=c) for c in drawn])


def check_stacked_table(case: str, bundle: Bundle, stratifier: str) -> None:
    """The reshape that makes the split table narrow, and what it must preserve."""
    from mlos_review.blocks import (aj_highlights_table, aj_outcome_codes,
                                    requires_aj, stacked_by_measure)
    from mlos_review.names import Vocabulary as _Vocab

    if not requires_aj(bundle, stratifier):
        return
    vocab = _Vocab(bundle.data)
    flat = aj_highlights_table(bundle, stratifier, vocab)
    codes = aj_outcome_codes(bundle, stratifier)
    stacked = stacked_by_measure(flat, vocab)

    expect_equal(f"{case}: one column per outcome ({stratifier})",
                 list(stacked.df.columns), codes)
    expect_equal(f"{case}: two rows per level ({stratifier})",
                 len(stacked.df), 2 * len(flat.df))
    expect_equal(f"{case}: rows are level then measure ({stratifier})",
                 stacked.df.index.nlevels, 2)

    # Every number survives the turn, in the right cell. Reshaping is the one
    # thing this block does, so a value landing in the wrong row is the only
    # bug it can have. The columns are reconstructed from measure and code
    # rather than parsed back out of the column name, so this check cannot
    # share a parsing bug with the block under test.
    for level in flat.df.index:
        for measure in blocks.AJ_OUTCOME_MEASURES:
            for code in codes:
                column = f"{measure}_{code}"
                if column not in flat.df.columns:
                    continue
                label = vocab.metric(measure).text(short=True)
                before = flat.df.loc[level, column]
                after = stacked.df.loc[(level, label), code]
                # A missing value arrives from JSON as None and comes out of a
                # float column as NaN. Both print as "n/a"; neither is a loss.
                missing = (before is None or before != before,
                           after is None or after != after)
                expect(f"{case}: {level} {label} {code} survives the reshape "
                       f"({stratifier})",
                       all(missing) or (not any(missing) and before == after),
                       f"{before!r} became {after!r}")
                if flat.flags is not None:
                    expect_equal(f"{case}: its flag comes with it ({stratifier})",
                                 stacked.flags.loc[(level, label), code],
                                 flat.flags.loc[level, column])

    # A column now holds a share on one line and a day count on the next, so
    # the format has to follow the row. Getting this wrong prints 20.5 days as
    # 2,054.8%.
    expect(f"{case}: shares are still percents ({stratifier})",
           stacked.row_formats[vocab.metric("aj_final_cif").text(short=True)].percent)
    expect(f"{case}: mean days are not ({stratifier})",
           not stacked.row_formats[
               vocab.metric("aj_restricted_mean").text(short=True)].percent)
    expect(f"{case}: no column format survives to contradict them ({stratifier})",
           not stacked.formats)
    expect(f"{case}: the flat arrangement's reading note is dropped "
           f"({stratifier})",
           all("then mean days" not in note for note in stacked.footnotes),
           f"{stacked.footnotes}")


def check_aj_teaser_pick(case: str, bundle: Bundle) -> None:
    """The second teaser: fewest lines wins, and it must have lines at all."""
    from mlos_review.blocks import aj_outcome_codes, requires_aj
    from mlos_review.deck import aj_teaser_stratifier
    from mlos_review.settings import Settings

    candidates = []
    for stratifier in bundle.stratifiers():
        if stratifier == "all" or not requires_aj(bundle, stratifier):
            continue
        codes = aj_outcome_codes(bundle, stratifier)
        if any(bundle.figure("aj_cif", stratifier, outcome=c) for c in codes):
            candidates.append(stratifier)

    pick = aj_teaser_stratifier(bundle, Settings())
    if not candidates:
        expect(f"{case}: no second teaser without a plotted stratifier",
               pick is None, f"got {pick}")
        return

    fewest = min(len(bundle.levels(s)) for s in candidates)
    expect_equal(f"{case}: the second teaser has the fewest lines",
                 len(bundle.levels(pick)), fewest)
    expect_equal(f"{case}: ties go to the bundle's stratifier order", pick,
                 next(s for s in candidates if len(bundle.levels(s)) == fewest))


def check_schema_version() -> None:
    """The version this package expects is the version R writes.

    The one constant that has to be maintained in two languages, so it is the
    one worth a test: mlos_results.R is read for its MLOS_RESULTS_SCHEMA_VERSION
    and compared against bundle.EXPECTED_SCHEMA_VERSION. Bumping the schema
    without telling this side is then a test failure rather than a warning
    printed on every run from then on.

    The goldens are checked with it, which is what proves the warning stays
    quiet in the ordinary case: they are written by the same R the constant came
    from, so any of them tripping the mismatch means the field itself moved.
    """
    from mlos_review import bundle as bundle_mod

    section("schema version")
    source = (REPO_ROOT / "mlos_results.R").read_text()
    match = re.search(r"^MLOS_RESULTS_SCHEMA_VERSION\s*<-\s*(\d+)L?\s*$",
                      source, re.MULTILINE)
    expect("mlos_results.R declares a schema version", match is not None)
    if match:
        expect_equal("mlos_review expects the version R writes",
                     bundle_mod.EXPECTED_SCHEMA_VERSION, int(match.group(1)))

    for path in sorted((REPO_ROOT / "tests" / "golden").glob("*/results.json")):
        data = json.loads(path.read_text())
        expect_equal(f"{path.parent.name}: golden carries the expected version",
                     data.get("schema_version"), bundle_mod.EXPECTED_SCHEMA_VERSION)


def check_dependency_declaration() -> None:
    """pyproject.toml and the Colab notebook agree on what has to be installed.

    The notebook keeps its own explicit `pip install` rather than running
    `pip install -e .`, and deliberately: a Colab session holds whatever files
    were uploaded into it, and a cell that fails because pyproject.toml was not
    among them would be a worse notebook than the one that names its packages.
    The cost of that is a second copy of the dependency list, so the copies are
    compared here. pyproject.toml is the declaration; the notebook is checked
    against it, never the other way round.

    Distribution names only, not versions. The notebook takes whatever Colab's
    image has, and pinning it to lower bounds it cannot honour would turn a
    working session into a failing one for no gain.

    The expected set is spelled out rather than derived, so that adding an
    import does not quietly widen what both copies are allowed to say. Two of
    these were missing from both for a while, because the deck built anyway on
    a machine that happened to carry them: `matplotlib`, which figures.py
    imports at module scope and deck.py reaches through, and `openpyxl`, which
    every build needs for the companion workbook.
    """
    section("dependency declaration")
    pyproject = (REPO_ROOT / "pyproject.toml").read_text()
    block = re.search(r"^dependencies\s*=\s*\[(.*?)\]", pyproject,
                      re.MULTILINE | re.DOTALL)
    expect("pyproject.toml declares dependencies", block is not None)
    if block is None:
        return
    declared = {name.lower() for name in
                re.findall(r'"([A-Za-z0-9_.-]+)', block.group(1))}
    expect("and they are the five the package imports",
           declared == {"python-pptx", "pandas", "pyyaml", "matplotlib",
                        "openpyxl"},
           f"got {sorted(declared)}")

    notebook = (REPO_ROOT / "colab_mlos.ipynb").read_text()
    installs = re.findall(r"pip install --quiet ([a-zA-Z0-9 _.-]+)", notebook)
    expect("the notebook installs the Python packages", len(installs) > 0)
    # Every occurrence, not the first: the notebook has carried two copies of
    # this line before, and a check that stopped at one would have let them
    # drift apart.
    for line in installs:
        expect_equal("the notebook installs exactly what pyproject declares",
                     {name.lower() for name in line.split()}, declared)


def check_notebook_file_listing() -> None:
    """The Colab notebook's upload list names every file the package has.

    A Colab session holds only what was uploaded into it, so that list is the
    install: a module missing from it is missing from the session, and Cell 3
    ends in an ImportError rather than in a deck. Nothing about the list is
    derived, which is how `recommend.py` and `salience.py` came to be absent
    from it while `deck.py` imported both.

    Compared as sets, since the notebook lists the modules in reading order and
    not alphabetically. The data files beside the modules are checked too: the
    diagram is opened at build time and travels with the package as
    package-data, so leaving it behind fails the same way an absent module
    would.
    """
    section("notebook file listing")
    notebook = json.loads((REPO_ROOT / "colab_mlos.ipynb").read_text())
    markdown = "".join("".join(cell["source"]) for cell in notebook["cells"]
                       if cell["cell_type"] == "markdown")
    block = re.search(r"^    mlos_review/\n((?:^        \S+\n)+)", markdown,
                      re.MULTILINE)
    expect("the notebook lists the package's files", block is not None)
    if block is None:
        return
    listed = set(block.group(1).split())
    present = {path.name for path in (REPO_ROOT / "mlos_review").iterdir()
               if path.is_file() and not path.name.startswith(".")}
    expect_equal("and lists exactly the files the package has", listed, present)


def check_stacked_underscore_codes() -> None:
    """An outcome code containing an underscore survives the reshape whole.

    The codes come from the user's palette configuration, so nothing upstream
    of this package forbids an underscore in one. Splitting a column name at
    its last underscore would cut such a code in half and file its numbers
    under the wrong outcome; the split has to go by known measure prefix. No
    fixture carries such a code, so the boundary is pinned synthetically.
    """
    from mlos_review.blocks import Table, stacked_by_measure

    section("stacked reshape (synthetic)")
    frame = pd.DataFrame(
        {
            "aj_final_cif_COMM_LIVE": [0.6, 0.4],
            "aj_final_cif_T": [0.4, 0.6],
            "aj_restricted_mean_COMM_LIVE": [12.0, 8.0],
            "aj_restricted_mean_T": [20.0, 30.0],
        },
        index=["a", "b"],
    )
    table = Table(df=frame, title="synthetic", formats={}, flags=None)
    stacked = stacked_by_measure(table, Vocabulary({}))

    expect_equal("an underscore code heads one column, intact",
                 list(stacked.df.columns), ["COMM_LIVE", "T"])
    share_label = Vocabulary({}).metric("aj_final_cif").text(short=True)
    expect_equal("its share lands under it",
                 stacked.df.loc[("a", share_label), "COMM_LIVE"], 0.6)


def check_capitalize_first() -> None:
    section("capitalization (synthetic)")
    expect_equal("acronyms survive", capitalize_first("median LOS"), "Median LOS")
    expect_equal("interior case untouched", capitalize_first("KM restricted mean"),
                 "KM restricted mean")
    expect_equal("empty string is safe", capitalize_first(""), "")


# ---------------------------------------------------------------------------
# Fixture checks
# ---------------------------------------------------------------------------


def check_vocabulary(case: str, bundle: Bundle) -> None:
    vocab, fallbacks, _ = names.report_fallbacks(bundle.data)
    expect(f"{case}: every name has a curated display name",
           not fallbacks, f"fallbacks={fallbacks}")

    bad_case = []
    for machine in sorted(names.METRICS):
        label = names.METRICS[machine].label
        if label[:1].isupper() and label.split()[0] not in ALLOWED_CAPITALIZED:
            bad_case.append(machine)
    expect(f"{case}: labels begin lowercase unless the first word is an acronym",
           not bad_case, f"offenders={bad_case}")

    # A CI suffix is a role on a base measure, not a measure of its own.
    round_trips = True
    for machine in list(names.METRICS)[:5]:
        for suffix, role in names.ROLE_SUFFIXES.items():
            base, parsed = vocab.split_role(machine + suffix)
            round_trips &= (base == machine and parsed == role)
    expect(f"{case}: CI and SE suffixes parse back to base and role", round_trips)

    # Every configured outcome code, plus the Any pseudo-state.
    unresolved = []
    for code in vocab.outcome_labels:
        for pattern, _, _ in names.OUTCOME_TEMPLATES:
            probe = pattern.pattern.strip("^$").replace(r"(\w+)", code)
            if "\\" in probe:
                continue
            if vocab.metric(probe).label == names._tokenize(probe):
                unresolved.append(probe)
    expect(f"{case}: outcome-templated names resolve for every code",
           not unresolved, f"unresolved={unresolved}")


def check_stratifier_order(case: str, bundle: Bundle) -> None:
    """Presentation order is the bundle's own order, with the baseline lifted.

    The point of the checks is that no list in this package decides the
    sequence. A reorder of the R registry must reach the deck unaided, and a
    stratifier this package has never heard of must still be presented rather
    than dropped, so both are asserted against a synthetic bundle whose keys
    are deliberately not the ones the vocabulary knows.
    """
    from mlos_review.bundle import BASELINE_STRATIFIER

    order = bundle.stratifiers()
    keys = [k for k in bundle.data.get("strata", {})]
    expect_equal(f"{case}: every present stratifier is presented",
                 sorted(order), sorted(keys))
    if BASELINE_STRATIFIER in keys:
        expect_equal(f"{case}: the baseline leads", order[0], BASELINE_STRATIFIER)
    expect_equal(f"{case}: the rest keep the bundle's order",
                 [s for s in order if s != BASELINE_STRATIFIER],
                 [k for k in keys if k != BASELINE_STRATIFIER])

    # Synthetic, and unrelated to the fixture: a reordered registry carrying a
    # dimension no vocabulary entry covers. Reading order from the JSON is what
    # makes both of these work, so a regression to a hardcoded list fails here
    # rather than silently dropping a stratifier from a future deck.
    invented = Bundle(data={"strata": {"group": {}, "all": {}, "intake": {},
                                       "facility": {}, "period": {}}},
                      root=Path("."))
    expect_equal(f"{case}: a reordered registry reaches the deck unaided",
                 invented.stratifiers(),
                 ["all", "group", "intake", "facility", "period"])


def check_bundle_access(case: str, bundle: Bundle) -> None:
    expect(f"{case}: has() rejects a path that is not there",
           not bundle.has("strata", "no_such_stratifier"))
    expect(f"{case}: figure() returns None for an unknown kind",
           bundle.figure("no_such_kind", "group") is None)
    check_stratifier_order(case, bundle)

    for stratifier in bundle.stratifiers():
        levels = bundle.levels(stratifier)
        expect(f"{case}: levels() is a list ({stratifier})", isinstance(levels, list))
        if bundle.has("strata", stratifier, "km"):
            frame = bundle.stratum(stratifier, "km")
            expect_equal(f"{case}: matrix rows match levels ({stratifier})",
                         list(frame.index), levels)


def check_full_table(case: str, bundle: Bundle, stratifier: str) -> blocks.Table:
    table = full_table(bundle, stratifier)
    levels = bundle.levels(stratifier)

    expect_equal(f"{case}: full_table has one row per level ({stratifier})",
                 list(table.df.index), levels)
    expect_equal(f"{case}: full_table columns are the declared measures ({stratifier})",
                 list(table.df.columns), [m for m, _ in FULL_TABLE_MEASURES])

    # The block SELECTS, it does not compute. Every cell must be the number the
    # bundle already holds, so a future refactor cannot quietly start deriving.
    mismatches = []
    for measure, block in FULL_TABLE_MEASURES:
        source = bundle.stratum(stratifier, block)[measure]
        for level in levels:
            got, want = table.df.loc[level, measure], source[level]
            if not (got == want or (got != got and want != want)):
                mismatches.append(f"{measure}/{level}")
    expect(f"{case}: every cell equals the bundle value ({stratifier})",
           not mismatches, f"mismatched={mismatches[:5]}")

    if len(levels) > 1:
        expect(f"{case}: flags present for a multi-level stratifier ({stratifier})",
               table.flags is not None)
    else:
        expect(f"{case}: flags suppressed for a single level ({stratifier})",
               table.flags is None)

    if table.flags is not None:
        wrong = []
        for measure in table.df.columns:
            values, marks = table.df[measure], table.flags[measure]
            if values.isna().all():
                continue
            if values.min() == values.max():
                if set(marks[values.notna()]) != {"F"}:
                    wrong.append(measure)
                continue
            if set(marks[values == values.max()]) != {"H"}:
                wrong.append(measure)
            if set(marks[values == values.min()]) != {"L"}:
                wrong.append(measure)
        expect(f"{case}: flags mark the true extremes ({stratifier})",
               not wrong, f"wrong={sorted(set(wrong))}")

    return table


def check_highlights(case: str, bundle: Bundle, stratifier: str,
                     full: blocks.Table) -> None:
    table = highlights_table(bundle, stratifier, full)
    shown, levels = list(table.df.index), list(full.df.index)

    expect(f"{case}: highlights are a subset of the levels ({stratifier})",
           set(shown) <= set(levels))
    expect(f"{case}: highlights keep the stratifier's level order ({stratifier})",
           shown == [lvl for lvl in levels if lvl in shown])
    expect(f"{case}: highlights are non-empty ({stratifier})", len(shown) > 0)

    if len(levels) <= HIGHLIGHT_SHOW_ALL_MAX:
        expect_equal(f"{case}: four levels or fewer show all ({stratifier})",
                     shown, levels)
    else:
        expect(f"{case}: more than four levels are narrowed ({stratifier})",
               len(shown) < len(levels))
        expect(f"{case}: a narrowed table says so in a footnote ({stratifier})",
               any(str(len(levels)) in note for note in table.footnotes),
               f"footnotes={table.footnotes}")

    # The invariant that keeps "H" honest: flags are judged over every level and
    # then subset, never recomputed on the rows that survived.
    if full.flags is not None:
        expect(f"{case}: highlight flags are the full-table flags ({stratifier})",
               table.flags is not None
               and table.flags.equals(full.flags.loc[shown]))

    if len(levels) > HIGHLIGHT_SHOW_ALL_MAX:
        pinned = (levels[-1],)
        with_pin = highlights_table(bundle, stratifier, full, pinned=pinned)
        expect(f"{case}: a pinned level is shown ({stratifier})",
               levels[-1] in with_pin.df.index, f"shown={list(with_pin.df.index)}")
        expect(f"{case}: pinning adds to the union, it does not replace it ({stratifier})",
               set(shown) <= set(with_pin.df.index),
               f"union={shown} with_pin={list(with_pin.df.index)}")
        expect(f"{case}: a pinned level absent from the data is ignored ({stratifier})",
               "NO_SUCH_LEVEL" not in
               highlights_table(bundle, stratifier, full, pinned=("NO_SUCH_LEVEL",)).df.index)

    expect(f"{case}: units are footnoted, not repeated per column ({stratifier})",
           any("in days" in note for note in table.footnotes),
           f"footnotes={table.footnotes}")


def check_formats(case: str, table: blocks.Table) -> None:
    """Percent is a display choice; the frame keeps the fraction."""
    missing = [m for m in table.df.columns if m not in table.formats]
    expect(f"{case}: every column has a format", not missing, f"missing={missing}")

    fraction = table.df.get("km_still_in_care_at_cap")
    if fraction is not None and fraction.notna().any():
        expect(f"{case}: the at-cap column is stored as a fraction, not a percent",
               bool((fraction.dropna() <= 1.0).all()), f"max={fraction.max()}")
        expect(f"{case}: the at-cap column is rendered as a percent",
               table.formats["km_still_in_care_at_cap"].percent)


def check_at_cap_agrees_with_aj(case: str, bundle: Bundle, stratifier: str,
                                full: blocks.Table) -> None:
    """The slide's at-cap column against the competing-risks route to it.

    The R suite asserts this on the bundle; repeating it here is not
    redundant, because what the slide shows is what the block selected, and a
    block that picked up the wrong row would still satisfy the R check.
    """
    if not bundle.has("strata", stratifier, "aj_final_cif"):
        return
    cif = bundle.stratum(stratifier, "aj_final_cif")
    if "aj_final_cif_Any" not in cif.columns:
        return
    def missing(value):
        # The bundle writes an absent number as JSON null, which arrives as
        # None rather than NaN, so a NaN test alone would not catch it.
        return value is None or value != value

    wrong = []
    for level in full.df.index:
        still = full.df.loc[level, "km_still_in_care_at_cap"]
        any_cif = cif.loc[level, "aj_final_cif_Any"]
        if missing(still) or missing(any_cif):
            continue
        if abs(still - (1 - any_cif)) > 1e-8:
            wrong.append(f"{level}: {still} vs {1 - any_cif}")
    expect(f"{case}: at-cap column equals 1 - AJ cif_Any ({stratifier})",
           not wrong, f"mismatched={wrong}")


def check_cap_sensitivity(case: str, bundle: Bundle, stratifier: str,
                          full: blocks.Table) -> None:
    """The caveat fires on exactly the levels whose curves cross the threshold.

    Judged against the FULL table, not the highlights, which is the property
    worth pinning: a level the cap is holding back belongs in the notes whether
    or not it won a row on the slide.
    """
    from mlos_review.names import Vocabulary as _Vocab

    notes = blocks.cap_sensitivity_notes(bundle, stratifier, _Vocab(bundle.data))
    at_cap = full.df["km_still_in_care_at_cap"].dropna()
    over = [str(level) for level, value in at_cap.items()
            if value > blocks.CAP_SENSITIVITY_THRESHOLD]

    expect(f"{case}: the cap caveat appears only when the cap binds ({stratifier})",
           bool(notes) == bool(over), f"over={over} notes={len(notes)}")
    if not over:
        return

    text = " ".join(notes)
    expect(f"{case}: every level over the threshold is named ({stratifier})",
           all(level in text for level in over), f"over={over} text={text}")
    unnamed = [str(level) for level, value in at_cap.items()
               if value <= blocks.CAP_SENSITIVITY_THRESHOLD
               and str(level) in text]
    expect(f"{case}: no level under the threshold is named ({stratifier})",
           not unnamed, f"named anyway={unnamed}")
    # The worst offender leads, so a reader who stops after the first name has
    # the one that matters most.
    expect(f"{case}: the largest at-cap level is named first ({stratifier})",
           text.index(str(at_cap.idxmax())) == min(text.index(l) for l in over),
           f"text={text}")


def check_order_shifts(case: str, bundle: Bundle, stratifier: str,
                       full: blocks.Table) -> None:
    """The note names every pair the three marks disagree about, and no other.

    Recomputed here from the table's own numbers rather than from the block, so
    an error in the rule cannot agree with itself. What is checked is the set
    of pairs: which one leads is `order_shifts`'s ordering and is pinned in the
    synthetic case, where the severities are chosen rather than whatever the
    shelter's data happens to hold.
    """
    vocab = Vocabulary(bundle.data)
    notes = blocks.order_shift_notes(full, vocab.stratifier(stratifier).label, vocab)

    for figure, _, measures in blocks.ORDER_SHIFT_TRIOS:
        present = [m for m in measures if m in full.df.columns]
        crossing, agreeing = [], []
        for first, second in itertools.combinations(full.df.index, 2):
            directions = set()
            for measure in present:
                ahead, behind = full.df.loc[first, measure], full.df.loc[second, measure]
                if ahead != ahead or behind != behind:
                    continue
                if abs(ahead - behind) >= blocks.ORDER_SHIFT_TOLERANCE:
                    directions.add(ahead > behind)
            pair = f"{first} with {second}"
            (crossing if len(directions) == 2 else agreeing).append(pair)

        said = [note for note in notes
                if note.startswith(f"Order shifts on the {figure}")]
        expect_equal(f"{case}: the {figure} note fires on a shift and not "
                     f"otherwise ({stratifier})", len(said), int(bool(crossing)))
        if not crossing:
            continue
        text = said[0]
        expect(f"{case}: every shifting pair of the {figure} is named ({stratifier})",
               all(pair in text for pair in crossing), f"{crossing} / {text}")
        expect(f"{case}: no agreeing pair of the {figure} is named ({stratifier})",
               not any(pair in text for pair in agreeing), f"{agreeing} / {text}")
        counted = "One pair" if len(crossing) == 1 else f"{len(crossing)} pairs"
        expect(f"{case}: the {figure} note counts the pairs it names ({stratifier})",
               counted in text, text)


def check_salience_statistic() -> None:
    """`stay_days_deviation`, its noise floor, and the threshold it must clear.

    Synthetic because every part of this is a number chosen to make one
    property visible: a fixture's spread is whatever its data happens to hold,
    which shows the code runs and not that it computes the right thing.
    """
    section("salience (synthetic)")

    def bundle_with(stays, errors, intakes, pooled=None) -> Bundle:
        levels = [f"L{i}" for i in range(len(stays))]
        matrix = lambda rows: {  # noqa: E731 - one shape, three uses
            "type": "matrix", "rows": list(rows), "columns": levels,
            "values": [list(values) for values in rows.values()]}
        node = {
            "km": matrix({"km_restricted_mean": stays,
                          "km_restricted_mean_se": errors}),
            "observations": matrix({"total_intakes": intakes}),
        }
        strata = {"group": node}
        if pooled is not None:
            strata["all"] = {"km": {
                "type": "matrix", "rows": ["km_restricted_mean"],
                "columns": ["All"], "values": [[pooled]]}}
        return Bundle(root=Path("."), data={"strata": strata})

    # Two levels, equal weights, ten days apart: the spread is half the gap,
    # which is what a standard deviation of two points is.
    plain = salience.spread(bundle_with([10.0, 20.0], [0.0, 0.0], [100, 100]), "group")
    expect_equal("two levels spread by half their gap", round(plain.observed, 6), 5.0)
    expect_equal("with no standard errors there is no floor", plain.noise, 0.0)
    expect_equal("so the deviation is the observed spread",
                 round(plain.deviation, 6), 5.0)

    # The constants themselves, pinned so that changing the rule is a visible
    # edit here rather than a silent change of which stratifiers earn slides.
    expect_equal("the threshold floor is a day",
                 salience.STAY_DAYS_DEVIATION_FLOOR, 1.0)
    expect_equal("and the share is 15% of mean LOS",
                 salience.STAY_DAYS_DEVIATION_SHARE, 0.15)

    # Both terms, taken on the POOLED mean rather than on the weighted mean of
    # the levels, so a bundle carrying both is judged on the published figure a
    # reader can look up. The pooled value here is the one that puts the
    # threshold exactly on a spread of 5, whatever the constants are, so the
    # boundary is tested rather than a number that happens to sit near it.
    on_the_line = (5.0 - salience.STAY_DAYS_DEVIATION_FLOOR) / salience.STAY_DAYS_DEVIATION_SHARE
    scaled = salience.spread(
        bundle_with([10.0, 20.0], [0.0, 0.0], [100, 100], pooled=on_the_line), "group")
    expect_equal("the threshold is a day plus a share of the pooled mean",
                 round(scaled.threshold, 6), 5.0)
    expect_equal("and the share is taken on the pooled figure",
                 scaled.mean_los, on_the_line)
    expect("a spread level with the threshold is not above it", not scaled.salient)

    # The floor with two equal weights is (SE1^2 + SE2^2) / 4, not the naive
    # average of the two variances: each level helps set the mean it is then
    # measured against, which is the (1 - w/W) factor. Getting that wrong
    # overstates the floor by a factor of two exactly here.
    noisy = salience.spread(bundle_with([10.0, 20.0], [4.0, 4.0], [100, 100]), "group")
    expect_equal("the floor accounts for levels setting their own mean",
                 round(noisy.noise, 6), round(math.sqrt((16 + 16) / 4), 6))
    expect_equal("and comes out of the spread in quadrature",
                 round(noisy.deviation, 6), round(math.sqrt(25 - 8), 6))

    # A spread no bigger than the floor is reported as no spread, rather than
    # as the square root of a negative number.
    swamped = salience.spread(bundle_with([10.0, 11.0], [9.0, 9.0], [100, 100]), "group")
    expect_equal("a spread inside the noise floors at zero", swamped.deviation, 0.0)
    expect("and is not salient", not swamped.salient)

    # Weighting is by intake COUNTS. A tiny level far from the rest moves the
    # spread by its share of arrivals, not by being one level of two.
    lopsided = salience.spread(
        bundle_with([10.0, 40.0], [0.0, 0.0], [9900, 100]), "group")
    expect(f"a rare level pulls the spread by its share of arrivals",
           2.9 < lopsided.observed < 3.0, f"{lopsided.observed}")

    # Nothing to compare: one level, or levels with no arrivals behind them.
    expect("one level has no spread",
           salience.spread(bundle_with([10.0], [0.0], [100]), "group") is None)
    expect("nor do levels nobody arrived into",
           salience.spread(bundle_with([10.0, 20.0], [0.0, 0.0], [0, 0]),
                           "group") is None)
    expect("and an uncomputable spread is not salient",
           not salience.is_salient(bundle_with([10.0], [0.0], [100]), "group"))

    # The settings answer first where they answer at all.
    from mlos_review.settings import Emphasis
    wide = bundle_with([10.0, 40.0], [0.0, 0.0], [100, 100])
    expect("NEVER outranks a wide spread",
           not salience.earns_slides(wide, "group", Emphasis(mode="NEVER")))
    narrow = bundle_with([10.0, 10.2], [0.0, 0.0], [100, 100])
    expect("ALWAYS outranks a narrow one",
           salience.earns_slides(narrow, "group", Emphasis(mode="ALWAYS")))
    expect("and AUTO takes the data's answer",
           salience.earns_slides(wide, "group", Emphasis())
           and not salience.earns_slides(narrow, "group", Emphasis()))


def check_single_stratifier_default() -> None:
    """`for_dataset`: what an unstated stratifier gets, and what it never touches."""
    section("emphasis default against the dataset (synthetic)")
    from mlos_review.settings import Emphasis, Settings, from_mapping

    alone = Settings().for_dataset(["group"])
    expect("one stratifier makes the default ALWAYS",
           alone.emphasis_for("group").forced, alone.emphasis_for("group"))
    several = Settings().for_dataset(["period", "intake", "group"])
    expect_equal("more than one leaves it AUTO",
                 several.emphasis_for("group").mode, "AUTO")
    expect_equal("and no stratifier at all leaves it alone",
                 Settings().for_dataset([]).emphasis_for("group").mode, "AUTO")

    # Only the default moves. A file that names the stratifier has said what it
    # wants, including when what it wants is the comparison AUTO stands for.
    stated = from_mapping({"emphasis": {"animal_group": "AUTO"}})
    expect_equal("an explicit AUTO survives a single-stratifier dataset",
                 stated.for_dataset(["group"]).emphasis_for("group").mode, "AUTO")
    refused = from_mapping({"emphasis": {"animal_group": "NEVER"}})
    expect("an explicit NEVER survives it too",
           refused.for_dataset(["group"]).emphasis_for("group").suppressed)
    pinned = from_mapping({"emphasis": {"animal_group": ["LRG"]}})
    expect_equal("and pinned levels are not disturbed",
                 pinned.for_dataset(["group"]).emphasis_for("group"),
                 Emphasis(mode="ALWAYS", levels=("LRG",)))
    # A stratifier the file does not name still takes the moved default, which
    # is the whole point: the file need not mention it at all.
    expect("an unnamed stratifier takes the dataset's default",
           from_mapping({"emphasis": {"animal_group": "NEVER"}})
           .for_dataset(["period"]).emphasis_for("period").forced)


def check_care_days(case: str, bundle: Bundle) -> None:
    """The three workload slides: what they tally, and what they will not divide.

    The spine of the run is an identity, so it is checked as one: residents
    times tenure per resident is animal-days, on the counted side and on the
    fitted side alike. If that stops holding, one of the three slides is
    reading a column it does not think it is reading.

    Shares need the levels to partition the standing population, which periods
    do not, so a period table carries its days owed and no share of them,
    however tempting the arithmetic is.
    """
    from mlos_review.deck import workload_slide, WORKLOAD_SLIDES
    from mlos_review.settings import Settings

    vocab = Vocabulary(bundle.data)
    for section in WORKLOAD_SLIDES:
        slide = workload_slide(bundle, vocab, section, Settings())
        wanted = [s for s in bundle.stratifiers()
                  if s != "all" and blocks.requires_workload(bundle, s, section)]
        expect_equal(f"{case}: the {section} slide is built where the data is",
                     slide is not None, bool(wanted))
        if slide is None:
            continue
        expect_equal(f"{case}: one {section} table per stratifier that has it",
                     [t.title for t in slide.tables],
                     [f"{blocks.WORKLOAD_SECTIONS[section][1]}, by "
                      f"{vocab.stratifier(s).label}" for s in wanted])

        for table, stratifier in zip(slide.tables, wanted):
            divides = blocks._levels_partition_census(bundle, stratifier)
            for share in ("intake_share", "resident_share", "fitted_share",
                          "owed_share"):
                if share not in blocks.WORKLOAD_SECTIONS[section][0]:
                    continue
                expect_equal(f"{case}: {stratifier} shares its {share} only "
                             f"where its levels divide the population",
                             share in table.df.columns, divides)
                # Shares are of the WHOLE, so they add to one only when the
                # table shows every level; a table that had to select says so
                # in its footnote and its shares fall short by what it left
                # out, which is the honest arithmetic rather than a rescaling.
                shown_all = len(table.df) == len(
                    blocks._workload_frame(bundle, stratifier))
                if divides and shown_all:
                    expect(f"{case}: {stratifier}'s {share} sums to one",
                           abs(table.df[share].sum() - 1.0) < 1e-9,
                           f"{table.df[share].sum()}")
                elif divides:
                    expect(f"{case}: {stratifier} says it selected rows",
                           any("are in the workbook" in note
                               for note in table.footnotes),
                           f"{table.footnotes}")
            if section == "census" and not divides:
                expect(f"{case}: {stratifier} says why it has no share",
                       any("divide the calendar" in note
                           for note in table.footnotes),
                       f"{table.footnotes}")

    # Every column comes from the bundle unchanged, and the two products hold.
    for stratifier in bundle.stratifiers():
        frame = blocks._workload_frame(bundle, stratifier)
        if frame.empty:
            continue
        for name, block in ((blocks.WORKLOAD_INTAKES, "observations"),
                            (blocks.WORKLOAD_RESIDENTS, "observations"),
                            (blocks.WORKLOAD_DAYS, "observations"),
                            (blocks.WORKLOAD_RESIDENTS_FITTED, "census"),
                            (blocks.WORKLOAD_TENURE_FITTED, "census"),
                            (blocks.WORKLOAD_DAYS_FITTED, "census"),
                            (blocks.WORKLOAD_DAYS_OWED, "census")):
            if name not in frame.columns:
                continue
            expect_equal(f"{case}: {stratifier} lists the bundle's own {name}",
                         list(frame[name]),
                         list(bundle.stratum(stratifier, block)[name]))

        # Counted and fitted alike: the animal-days column is the census times
        # the tenure column, which is what makes the third slide the first two
        # multiplied rather than a third measurement.
        for count, tenure, days in (
                (blocks.WORKLOAD_RESIDENTS, blocks.WORKLOAD_TENURE,
                 blocks.WORKLOAD_DAYS),
                (blocks.WORKLOAD_RESIDENTS_FITTED, blocks.WORKLOAD_TENURE_FITTED,
                 blocks.WORKLOAD_DAYS_FITTED)):
            if not {count, tenure, days} <= set(frame.columns):
                continue
            product = frame[count] * frame[tenure]
            for level in frame.index:
                if pd.isna(product[level]) or pd.isna(frame.at[level, days]):
                    continue
                expect(f"{case}: {stratifier} {level} residents x tenure is "
                       f"animal-days ({days})",
                       abs(product[level] - frame.at[level, days])
                       <= 1e-6 * max(1.0, abs(frame.at[level, days])),
                       f"{product[level]} against {frame.at[level, days]}")

        # And what an average day's residents still owe is what they have
        # served plus one day each, which is the census.
        if {blocks.WORKLOAD_DAYS_FITTED, blocks.WORKLOAD_DAYS_OWED,
            blocks.WORKLOAD_RESIDENTS_FITTED} <= set(frame.columns):
            gap = (frame[blocks.WORKLOAD_DAYS_OWED]
                   - frame[blocks.WORKLOAD_DAYS_FITTED])
            for level in frame.index:
                if pd.isna(gap[level]):
                    continue
                census = frame.at[level, blocks.WORKLOAD_RESIDENTS_FITTED]
                expect(f"{case}: {stratifier} {level} owes its elapsed days "
                       f"plus one per resident",
                       abs(gap[level] - census) <= 1e-6 * max(1.0, abs(census)),
                       f"{gap[level]} against {census}")

    # The burden-carrier sentence, which fires only on a level that holds a
    # majority of the days owed AND more of them than of the residents. Checked
    # against the bundle rather than against the block.
    said = blocks.findings_for_care_days(bundle)
    carriers = [line for line in said if "of the care days still owed" in line]
    expected = 0
    for stratifier in blocks.summary_fields(bundle):
        owed = blocks._care_days_owed(bundle, stratifier)
        if owed is None or owed.sum() <= 0 or len(owed) < 2:
            continue
        # The shares below are only meaningful where the levels divide the
        # standing population, which is the condition the rule itself applies.
        if not blocks._levels_partition_census(bundle, stratifier):
            continue
        census = bundle.stratum(stratifier, "census")["expected_census"]
        share = owed / owed.sum()
        top = share.idxmax()
        if (share[top] > blocks.CARE_DAYS_MAJORITY
                and share[top] - (census / census.sum())[top]
                >= blocks.CARE_DAYS_DISPROPORTION):
            expected += 1
    expect_equal(f"{case}: a burden-carrier is named exactly where there is one",
                 len(carriers), expected)


def check_recommendations(case: str, bundle: Bundle) -> None:
    """Each rule fires on exactly the condition it claims, and names its knob.

    Every recommendation here is checked against the bundle rather than against
    the rule, and the rules that should stay quiet on a healthy dataset are
    checked for staying quiet: a section that fires on every deck is the
    failure mode, not an empty one.
    """
    from mlos_review.regression import comparison as _comparison, pooled, stratified
    from mlos_review import recommend

    vocab = Vocabulary(bundle.data)
    comparison = _comparison(pooled(bundle), stratified(bundle))

    # Narrowing follows the burden-carrier test exactly, plus a size floor:
    # the deck must not suggest re-analyzing a level it did not name as the
    # carrier, nor one too small to carry the analysis.
    expected = []
    for stratifier in blocks.summary_fields(bundle):
        found = blocks.burden_carrier(bundle, stratifier)
        if found is None or not bundle.has("strata", stratifier, "observations"):
            continue
        stays = bundle.stratum(stratifier, "observations")["total_intakes"]
        if found[0] in stays.index and stays[found[0]] >= recommend.NARROWING_MINIMUM_STAYS:
            expected.append(found[0])
    said = recommend.narrowing(bundle, vocab)
    expect_equal(f"{case}: narrowing is suggested for each burden carrier",
                 len(said), len(expected))
    for level, line in zip(expected, said):
        expect(f"{case}: it names {level} and the setting that filters it",
               level in line and "filter`" in line, line)

    # The gap remedy travels with the gap finding: both are gated on a gap
    # inside the plotted range, so they fire together or not at all.
    expect_equal(f"{case}: the gap remedy fires exactly when the gap is found",
                 bool(recommend.gap_remedy(bundle, vocab)),
                 bool(blocks.findings_for_gaps(bundle, vocab)))

    for stratifier in bundle.stratifiers():
        if stratifier == "all":
            continue
        found = salience.spread(bundle, stratifier)
        expect_equal(f"{case}: a dimension is questioned only where it did not "
                     f"separate ({stratifier})",
                     bool(recommend.dimension_not_separating(bundle, stratifier, vocab)),
                     found is not None and not found.salient)

        # The cap recommendation is deliberately quieter than the cap caveat:
        # saying a number is a lower bound is worth doing at half a percent,
        # and changing a setting is not.
        km = bundle.stratum(stratifier, "km")
        at_cap = km["km_still_in_care_at_cap"].dropna()
        binding = at_cap[at_cap > recommend.CAP_RECOMMENDATION_THRESHOLD]
        said = recommend.cap_binding(bundle, stratifier, vocab)
        expect_equal(f"{case}: raising the cap is suggested only where it binds "
                     f"hard ({stratifier})", bool(said), not binding.empty)
        if said:
            expect(f"{case}: it names the setting ({stratifier})",
                   "restricted_stay_cap" in said[0], said[0])
        expect(f"{case}: the recommendation is quieter than the caveat "
               f"({stratifier})",
               recommend.CAP_RECOMMENDATION_THRESHOLD
               > blocks.CAP_SENSITIVITY_THRESHOLD)

        # A level is flagged for its tail only when it is more dispersed than
        # the shelter it came from, by the named multiple.
        pooled_km = bundle.stratum("all", "km")
        base = recommend._ratio(pooled_km["km_p90_los"].iloc[0],
                                pooled_km["km_median_los"].iloc[0])
        wide = []
        if base:
            for level in km.index:
                ratio = recommend._ratio(km.loc[level, "km_p90_los"],
                                         km.loc[level, "km_median_los"])
                if ratio and ratio >= recommend.TAIL_SPREAD_MULTIPLE * base:
                    wide.append(str(level))
        said = recommend.tail_spread(bundle, stratifier, vocab)
        expect_equal(f"{case}: a tail is questioned once per dispersed level "
                     f"({stratifier})", len(said), len(wide))
        expect(f"{case}: and it names them ({stratifier})",
               all(any(level in line for line in said) for level in wide),
               f"{wide}")

        # Levels the robustness slide could not establish anything about.
        thin = []
        if (not comparison.empty
                and stratifier in comparison.index.get_level_values("stratifier")):
            margins = comparison.xs(stratifier, level="stratifier")["agreement_margin"]
            thin = [str(i) for i in margins[margins >= blocks.WIDE_MARGIN].dropna().index]
        said = recommend.unestimable_levels(bundle, stratifier, comparison, vocab)
        expect_equal(f"{case}: unestimable levels are raised where there are "
                     f"any ({stratifier})", bool(said), bool(thin))
        if thin:
            expect(f"{case}: and every one is named ({stratifier})",
                   all(level in said[0] for level in thin), said[0])

        # The shape rule holds a level only when every variant that estimated
        # it agrees, which is what stops a design-cell artefact from reading as
        # a property of the level.
        said = recommend.falling_hazard(bundle, stratifier, vocab)
        # A crossed variant's tabulated ratio is a contrast taken inside the
        # reference level of the other shape dimension, and reading the run's
        # remaining variants without it would test agreement against a bench
        # one fit short. One crossed variant therefore silences the run, and
        # silence is the assertion.
        variants = bundle.value("weibull", "shape_variants", default={}) or {}
        if any(isinstance(v, dict) and (v.get("shape_crossing") or {}).get("crossed")
               for v in variants.values()):
            expect(f"{case}: no shape recommendation once a variant crossed "
                   f"({stratifier})", not said, said)
        if said:
            stays = (_unified_stays(bundle, stratifier)
                     if bundle.has("strata", stratifier, "unified_stay_counts")
                     else None)
            for level in bundle.levels(stratifier):
                if str(level) not in said[0]:
                    continue
                bounds = _shape_bounds(bundle, stratifier, str(level))
                # Both tests, as the rule applies them: the level's own k has to
                # clear the absolute ceiling AND the shelter's pooled shape
                # scaled down, and it is the whole interval that has to clear,
                # not the estimate.
                pooled = bundle.value("weibull", "shape")
                ceiling = recommend.SHAPE_CEILING
                if pooled is not None and pooled == pooled:
                    ceiling = min(ceiling,
                                  float(pooled) * recommend.SHAPE_POOLED_CEILING)
                expect(f"{case}: {level}'s own shape interval clears the ceiling "
                       f"in every variant that fitted it",
                       bounds and all(u < ceiling for _, u in bounds),
                       f"{bounds} against {ceiling:.4f}")
                if stays is not None and str(level) in stays.index:
                    share = float(stays[str(level)]) / float(stays.sum())
                    expect(f"{case}: {level} is big enough to be recommended "
                           f"at all", share >= recommend.RECOMMENDATION_MINIMUM_SHARE,
                           f"{share:.4f}")


def _shape_bounds(bundle: Bundle, stratifier: str,
                  level: str) -> list[tuple[float, float]]:
    """Every (own k, upper bound) the Weibull variants published for one level.

    Only reached on a run with no crossed variant, since one of those silences
    `falling_hazard` outright, so every variant here is one the rule read.
    """
    import pandas as _pd

    found = []
    for variant in (bundle.value("weibull", "shape_variants", default={}) or {}).values():
        if not isinstance(variant, dict) or "shape_table" not in variant:
            continue
        table = _pd.DataFrame(variant["shape_table"]["columns"])
        for _, row in table.iterrows():
            term = str(row["variable"])
            upper = row.get("shape_own_upper")
            if (term.endswith(level) and len(term) > len(level)
                    and upper is not None and upper == upper):
                found.append((float(row["shape_own"]), float(upper)))
    return found


def check_outcome_spread(case: str, bundle: Bundle, stratifier: str) -> None:
    """Both axes of the outcome-timing table, and the floor they share.

    The same numbers read two ways: down a column, comparing levels on one
    outcome, and across a row, comparing one level's outcomes with each other.
    Recomputed here from the bundle rather than through `_outcome_timings`, so
    a mistake in the shared frame cannot agree with itself.

    The event floor is the part worth pinning on both. Without it the widest
    ratio either way is usually the thinnest cell's, and the finding would be
    about a sample size rather than about the shelter.
    """
    vocab = Vocabulary(bundle.data)
    label = vocab.stratifier(stratifier).label
    spread_said = blocks.findings_for_outcome_spread(bundle, stratifier, label, vocab)
    contrast_said = blocks.findings_for_outcome_contrast(bundle, stratifier,
                                                         label, vocab)
    if not (bundle.has("strata", stratifier, "aj_restricted_mean")
            and bundle.has("strata", stratifier, "outcomes")):
        expect(f"{case}: no AJ means, no timing findings ({stratifier})",
               not spread_said and not contrast_said,
               f"{spread_said}{contrast_said}")
        return

    means = bundle.stratum(stratifier, "aj_restricted_mean")
    events = bundle.stratum(stratifier, "outcomes")
    floor = blocks.OUTCOME_SPREAD_MINIMUM_EVENTS

    # The frame both findings read, rebuilt: a cell survives only where the
    # level had enough of that outcome to average.
    cells = {}
    for code in blocks.aj_outcome_codes(bundle, stratifier):
        column, counted = f"aj_restricted_mean_{code}", f"outcome_{code}_events"
        if column not in means.columns or counted not in events.columns:
            continue
        for level in means.index:
            # A run that never saw this outcome in this level writes null,
            # which arrives as None rather than as NaN.
            value = means.loc[level, column]
            if value is None or value != value or float(value) <= 0:
                continue
            if float(events[counted].get(level, 0) or 0) >= floor:
                cells[(str(level), code)] = float(value)

    def widest(groups):
        best = (0.0, None)
        for key, values in groups.items():
            if len(values) < 2:
                continue
            ratio = max(values.values()) / min(values.values())
            if ratio > best[0]:
                best = (ratio, (key, values))
        return best

    by_outcome, by_level = {}, {}
    for (level, code), value in cells.items():
        by_outcome.setdefault(code, {})[level] = value
        by_level.setdefault(level, {})[code] = value

    for name, groups, said in (("spread", by_outcome, spread_said),
                               ("contrast", by_level, contrast_said)):
        ratio, chosen = widest(groups)
        fires = ratio >= blocks.OUTCOME_SPREAD_MULTIPLE
        expect_equal(f"{case}: the {name} is stated once, when it is wide "
                     f"({stratifier})", len(said), int(fires))
        if not fires:
            continue
        _, values = chosen
        slow = max(values, key=values.get)
        fast = min(values, key=values.get)
        shown = f"{ratio:.1f}-fold" if name == "spread" else f"{ratio:.1f} times"
        expect(f"{case}: the {name} names both ends and the multiple "
               f"({stratifier})",
               str(slow) in said[0] and str(fast) in said[0] and shown in said[0],
               f"{slow}/{fast}/{ratio:.1f}: {said[0]}")

    # Nothing too thin to average may anchor either finding, which is the whole
    # job of the floor.
    thin = {f"{level} {code}" for code in blocks.aj_outcome_codes(bundle, stratifier)
            for level in means.index
            if (level, code) not in cells}
    expect(f"{case}: the floor blanked every thin cell ({stratifier})",
           all((str(level), code) in cells or f"{level} {code}" in thin
               for level in means.index
               for code in blocks.aj_outcome_codes(bundle, stratifier)),
           f"{sorted(thin)}")


def check_salience_note(case: str, bundle: Bundle, stratifier: str) -> None:
    """The LOS slide reports the spread, and reports the verdict it produced.

    On the LOS slide rather than the outlook slide, which is the property worth
    pinning: a stratifier that failed the test has no outlook slide, and that
    is exactly when a reader wants to know what the number was.
    """
    from mlos_review.deck import los_by_stratifier
    from mlos_review.settings import Emphasis, Settings

    vocab = Vocabulary(bundle.data)
    auto = Settings(emphasis={stratifier: Emphasis(mode="AUTO")})
    slide = los_by_stratifier(bundle, stratifier, vocab, auto)
    if slide is None:
        return
    said = [note for note in slide.notes if note.startswith("Salience.")]
    found = salience.spread(bundle, stratifier)

    expect_equal(f"{case}: the spread is reported where it can be computed "
                 f"({stratifier})", len(said), int(found is not None))
    if found is None:
        return
    text = "\n".join(slide.notes)
    expect(f"{case}: it prints the deviation and the threshold ({stratifier})",
           f"{found.deviation:.1f} days" in text
           and f"{found.threshold:.1f} days" in text, text)
    expect_equal(f"{case}: the verdict it states is the one it acted on "
                 f"({stratifier})",
                 "above the threshold" in text, found.salient)


def check_stratified_outlook(case: str, bundle: Bundle, stratifier: str) -> None:
    """The outlook slide per level: who gets one, and what is on it.

    Built against a staged copy carrying the figures, because the rule requires
    at least one and the golden bundles ship no PNGs, so against the bare
    fixture every assertion here would pass on a slide that was never built.
    """
    try:
        from PIL import Image
    except ImportError:  # pragma: no cover - Pillow ships with python-pptx
        return
    from mlos_review.deck import los_by_stratifier, resident_outlook_by_stratifier
    from mlos_review.settings import Emphasis, Settings

    plots = [e["plot"] for e in bundle.data.get("outputs", []) if e.get("plot")]
    if not plots or not requires_full_table(bundle, stratifier):
        return

    with tempfile.TemporaryDirectory() as tmp:
        staged = Path(tmp) / "fixture"
        shutil.copytree(REPO_ROOT / "tests" / "golden" / case, staged)
        for name in plots:
            Image.new("RGB", (300, 200), "white").save(staged / name)
        staged_bundle = Bundle.load(staged)
        vocab = Vocabulary(staged_bundle.data)

        both = [staged_bundle.figure(kind, stratifier)
                for kind in ("km_census_by_tenure", "km_remaining_los")]
        drawn = any(f is not None for f in both)

        # Under AUTO the spread decides, and the two answers have to agree:
        # the slide exists exactly where the stratifier is salient and there is
        # something to draw on it.
        auto = Settings(emphasis={stratifier: Emphasis(mode="AUTO")})
        expect_equal(f"{case}: AUTO leaves the outlook slide to the spread "
                     f"({stratifier})",
                     resident_outlook_by_stratifier(staged_bundle, stratifier,
                                                    vocab, auto) is not None,
                     drawn and salience.is_salient(staged_bundle, stratifier))
        never = Settings(emphasis={stratifier: Emphasis(mode="NEVER")})
        expect(f"{case}: NEVER declines however far apart the levels are "
               f"({stratifier})",
               resident_outlook_by_stratifier(staged_bundle, stratifier, vocab,
                                              never) is None)

        forced = Settings(emphasis={stratifier: Emphasis(mode="ALWAYS")})
        slide = resident_outlook_by_stratifier(staged_bundle, stratifier, vocab,
                                               forced)
        if not drawn:
            expect(f"{case}: no figures, no outlook slide ({stratifier})",
                   slide is None)
            return

        expect(f"{case}: a salient stratifier gets one ({stratifier})",
               slide is not None)
        expect_equal(f"{case}: it carries the census and remaining curves "
                     f"({stratifier})",
                     slide.figures, [f for f in both if f is not None])
        expect_equal(f"{case}: its columns are the outlook measures ({stratifier})",
                     list(slide.table.df.columns),
                     [m for m in blocks.OUTLOOK_SLIDE_MEASURES
                      if m in full_table(staged_bundle, stratifier).df.columns])
        # The pair reads as one statement about one set of levels, so its rows
        # are the rows of the slide it follows, not a fresh selection.
        los = los_by_stratifier(staged_bundle, stratifier, vocab, forced)
        expect_equal(f"{case}: it shows the levels the LOS slide showed "
                     f"({stratifier})",
                     list(slide.table.df.index), list(los.table.df.index))
        expect(f"{case}: it names the resident mean apart from the readings "
               f"({stratifier})",
               slide.table.headers.get("per_resident_future_days")
               == "Remaining per resident")


def check_order_shift_thresholds() -> None:
    """Which differences count, which are noise, and which pair leads.

    The tolerance is the whole rule, and no fixture pins it: a shelter's levels
    differ by whatever they differ by, so a fixture can show that the rule runs
    but not that it draws the line in the right place. These do.
    """
    section("order shifts (synthetic)")

    measures = ["km_median_los", "km_restricted_mean", "km_p90_los"]

    def shifts(**levels) -> list[str]:
        frame = pd.DataFrame.from_dict(levels, orient="index", columns=measures)
        found = blocks.order_shifts(frame, measures)
        return [f"{first} with {second}" for first, second, _ in found]

    # An exact tie on the median is a tie, not a direction: the pair agrees as
    # far as the marks that separate it go.
    expect_equal("an exact tie is not a shift",
                 shifts(A=[10, 20, 30], B=[10, 25, 40]), [])
    # Just inside the margin on the mean, against a real two-day median gap.
    expect_equal("a difference under a day is noise",
                 shifts(A=[10, 20.9, 30], B=[12, 20.0, 40]), [])
    # And just outside it. One day is a difference; the threshold is where it
    # starts counting, not where it stops.
    expect_equal("a difference of exactly a day counts",
                 shifts(A=[10, 21.0, 30], B=[12, 20.0, 40]), ["A with B"])
    # A mark one level has no value for drops out of that pair. An unreached
    # P90 is common and leaves the median and the mean perfectly comparable.
    expect_equal("an absent mark is dropped, not the pair",
                 shifts(A=[10, 30, float("nan")], B=[12, 25, 40]), ["A with B"])
    expect_equal("and it cannot make a shift on its own",
                 shifts(A=[10, 20, float("nan")], B=[12, 25, 40]), [])

    # Most substantial first, by the SMALLER side of each reversal: B and C
    # cross by nine days against A and B's five, though A and B are further
    # apart on the mark they agree about.
    ordered = shifts(A=[10, 10, 10], B=[20, 5, 10], C=[11, 11, 40])
    expect_equal("the deepest reversal leads", ordered, ["B with C", "A with B"])

    vocab = Vocabulary()
    frame = pd.DataFrame.from_dict(
        {"A": [10, 30, float("nan")], "B": [12, 25, 40]},
        orient="index", columns=measures)
    said = blocks.order_shift_notes(
        blocks.Table(df=frame, title="synthetic"), "animal group", vocab)[0]
    expect("the worked example takes the leading pair", "Take A with B" in said, said)
    expect("it names the marks that split the pair",
           "higher mean LOS (30 against 25)" in said
           and "lower median LOS (10 against 12)" in said, said)
    expect("and says nothing about the mark it could not read",
           "P90" not in said, said)


def check_observation_gaps(case: str, bundle: Bundle) -> None:
    """Every gap R recorded reaches exactly one of the two opening slides.

    A gap does not qualify a curve, it invalidates it from that day onward, so
    losing one is the failure to catch: the check is run against the bundle's
    own gaps table rather than against the split the blocks chose. The periods
    report on the title slide, which is where they are listed, and everything
    else on the summary slide, which carries the pooled row.

    The reassurance is checked too, and by the same rule that produces it: a
    slide speaks about the analyses R actually scanned, so the title slide is
    silent on a single-period run rather than claiming a clean bill it has no
    scan behind.
    """
    from mlos_review.deck import opening_slides

    vocab = Vocabulary(bundle.data)
    opening = opening_slides(bundle, vocab)
    slides = {"window": "\n".join(opening[0].notes),
              "sample": "\n".join(opening[1].notes) if len(opening) > 1 else ""}

    fields = set(blocks.summary_fields(bundle))
    plot_cap = bundle.value("settings", "presentation", "plot_stay_cap")
    gaps = blocks.observation_gaps(bundle, bundle.stratifiers())

    recorded = (len(bundle.table("unified", "gaps"))
                if bundle.has("unified", "gaps") else 0)
    expect_equal(f"{case}: every recorded gap resolves to a stratifier the deck shows",
                 len(gaps), recorded)

    def slide_for(stratifier: str) -> str:
        return ("window" if stratifier != "all" and stratifier not in fields
                else "sample")

    for _, row in gaps.iterrows():
        where = slide_for(row["stratifier"])
        other = "sample" if where == "window" else "window"
        stratum = str(row["stratum"])
        subject = "the pooled data" if row["stratifier"] == "all" else stratum
        detail = f"day {row['gap_start_day']:g} to day {row['gap_end_day']:g}"
        expect(f"{case}: the {stratum} gap is reported on the {where} slide",
               subject in slides[where], slides[where])
        # Below the plot cap it is a hole in a curve the audience is looking
        # at, so it is named with its days. Above it, briefly: the deck's
        # figures stop before it, and what a reader can still act on is where
        # the first one starts, not how far each one runs.
        in_view = plot_cap is None or row["gap_start_day"] < plot_cap
        expect_equal(f"{case}: the {stratum} gap is detailed only inside the plot",
                     detail in slides[where], in_view)
        expect(f"{case}: the {stratum} gap is not also on the {other} slide",
               detail not in slides[other], slides[other])

    if plot_cap is not None:
        past = gaps[gaps["gap_start_day"] >= plot_cap]
        for where in set(past["stratifier"].map(slide_for)):
            starts = past["gap_start_day"][past["stratifier"].map(slide_for) == where]
            expect(f"{case}: the {where} slide dates the first gap past the plot",
                   f"day {starts.min():g}" in slides[where], slides[where])

    # The title slide reports on the periods exactly when R scanned them, and
    # the summary slide always does: the pooled scan runs on every analysis.
    scanned = any(bundle.value("settings", "coverage", stratifier, "included",
                               default=False)
                  for stratifier in bundle.stratifiers()
                  if stratifier != "all" and stratifier not in fields)
    expect_equal(f"{case}: the window slide speaks only for a scan that ran",
                 "Observation gaps" in slides["window"], scanned)
    if len(opening) > 1:
        expect(f"{case}: the summary slide always reports the pooled scan",
               "Observation gaps" in slides["sample"], slides["sample"])
    if gaps.empty:
        expect(f"{case}: a clean run says so rather than staying silent",
               "Observation gaps: none" in slides["sample"] + slides["window"],
               slides["sample"])

    # A gap inside the plotted range invalidates a curve the audience is
    # looking at, so it is a finding and not only a note. One past the cap
    # is not: no figure here crosses it, and the notes already carry it.
    in_view = gaps if plot_cap is None else gaps[gaps["gap_start_day"] < plot_cap]
    found = blocks.findings_for_gaps(bundle, vocab)
    expect_equal(f"{case}: a gap inside the plot is found, one past it is not",
                 len(found), int(not in_view.empty))
    expect(f"{case}: the opening carries that finding into the deck",
           found == opening[0].findings, f"{opening[0].findings}")


def check_falling_hazard_gates() -> None:
    """The three things that silence the shape recommendation.

    No fixture reaches any of them. `falling_hazard` fires on none of the cases
    (their shape ratios sit above the ceiling or their intervals are wide), and
    none has three qualifying predictors, so no fixture carries a crossed
    variant at all. Each gate is therefore built here, against a bundle holding
    exactly what the rule reads.
    """
    from mlos_review import recommend

    section("falling_hazard gates (synthetic)")

    def bundle_with(*, crossed: bool, upper: float, big_share: bool,
                    fallback: str | None = None, pooled: float = 0.90,
                    applicable: bool | None = None,
                    enabled: bool | None = None) -> Bundle:
        stays = [9000, 1000] if big_share else [9900, 100]
        return Bundle(root=Path("."), data={
            "strata": {
                "group": {
                    "labels": ["MED", "LARGE"],
                    "unified_stay_counts": {
                        "type": "table",
                        "columns": {"period_label": ["MED", "LARGE"],
                                    "total_stays": stays},
                    },
                },
            },
            "weibull": {"shape": pooled, "shape_variants": {"period": {
                "shape_crossing": {"crossed": crossed,
                                   "applicable": applicable,
                                   "enabled": enabled,
                                   "fallback_reason": fallback},
                "shape_table": {"type": "table", "columns": {
                    "variable": ["animal_groupMED", "animal_groupLARGE"],
                    "shape_ratio": [1.0, 0.60],
                    "ci_upper": [None, upper],
                    # What the rule actually reads: the level's own k and the
                    # top of its interval. MED is the reference, whose own k is
                    # the baseline and so is an estimate like any other.
                    "shape_own": [0.90, 0.60],
                    "shape_own_upper": [0.95, upper],
                }},
            }}},
        })

    vocab = names.Vocabulary()
    fired = recommend.falling_hazard(
        bundle_with(crossed=False, upper=0.70, big_share=True), "group", vocab)
    expect("falling_hazard: fires on an uncrossed variant clear of the ceiling",
           bool(fired) and "LARGE" in fired[0], fired)

    # The same numbers, from a shape crossed over two other dimensions: the
    # ratio is then a contrast inside the reference cell of the other one, and
    # the sentence could not say so.
    silent = recommend.falling_hazard(
        bundle_with(crossed=True, upper=0.70, big_share=True), "group", vocab)
    expect("falling_hazard: silent when the variant's shape was crossed",
           not silent, silent)

    # An additive fit is read whichever road it arrived by. The default, with
    # crossing available and switched off:
    fired = recommend.falling_hazard(
        bundle_with(crossed=False, upper=0.70, big_share=True,
                    applicable=True, enabled=False), "group", vocab)
    expect("falling_hazard: fires on the additive fit the default produces",
           bool(fired) and "LARGE" in fired[0], fired)

    # And a crossing asked for and refused, which produces the same fit. Reading
    # one but not the other would drop a recommendation for turning an
    # experimental setting on.
    fired = recommend.falling_hazard(
        bundle_with(crossed=False, upper=0.70, big_share=True,
                    applicable=True, enabled=True,
                    fallback="1 of 4 combinations hold fewer than 5 outcomes"),
        "group", vocab)
    expect("falling_hazard: fires on an additive fit reached by fallback",
           bool(fired) and "LARGE" in fired[0], fired)

    # Mixed: one variant crossed, one additive. The additive one alone would
    # fire, and that is the reading the gate refuses -- the crossed variant is
    # the fit that would have checked it.
    mixed = bundle_with(crossed=False, upper=0.70, big_share=True,
                        applicable=True, enabled=True)
    mixed.data["weibull"]["shape_variants"]["intake"] = {
        "shape_crossing": {"crossed": True, "applicable": True, "enabled": True,
                           "fallback_reason": None},
        "shape_table": {"type": "table", "columns": {
            "variable": ["animal_groupMED", "animal_groupLARGE"],
            "shape_ratio": [1.0, 0.60], "ci_upper": [None, 0.70],
            "shape_own": [0.90, 0.60], "shape_own_upper": [0.95, 0.70],
        }},
    }
    silent = recommend.falling_hazard(mixed, "group", vocab)
    expect("falling_hazard: silent when any variant in the run crossed",
           not silent, silent)

    # Estimate 0.60, well under the ceiling, but the interval reaches past it.
    silent = recommend.falling_hazard(
        bundle_with(crossed=False, upper=0.90, big_share=True), "group", vocab)
    expect("falling_hazard: silent when the interval reaches the ceiling",
           not silent, silent)

    # 1% of stays. Nothing about the estimate changed; the level is simply too
    # small to send anyone to, which is what keeps the _UNKNOWN_ fill out
    # without naming it.
    silent = recommend.falling_hazard(
        bundle_with(crossed=False, upper=0.70, big_share=False), "group", vocab)
    expect("falling_hazard: silent for a level under the minimum share",
           not silent, silent)

    # Absolute test passed, pooled test failed. Own k of 0.60 with an interval
    # reaching 0.70 is comfortably under SHAPE_CEILING, but on a shelter whose
    # own shape is 0.60 the level is not standing out from anything: the ceiling
    # becomes 0.60 x 0.95 = 0.57 and the interval does not clear it.
    silent = recommend.falling_hazard(
        bundle_with(crossed=False, upper=0.70, big_share=True, pooled=0.60),
        "group", vocab)
    expect("falling_hazard: silent when the level does not stand out from the "
           "shelter's own shape", not silent, silent)
    # The mirror: the same level on a shelter that mostly does not fall.
    fired = recommend.falling_hazard(
        bundle_with(crossed=False, upper=0.70, big_share=True, pooled=1.10),
        "group", vocab)
    expect("falling_hazard: fires where the shelter's shape leaves room",
           bool(fired) and "LARGE" in fired[0], fired)


def check_gaps_past_the_plot_cap() -> None:
    """The brief form, and the two forms side by side.

    `sim_geometric_period_effect` reaches the brief form on the title slide,
    with two period gaps past its plot cap. What no fixture has is a gap past
    the cap on the summary slide, or one of each on the same slide, so both are
    built here. The band between the two caps is where a real gap is most
    likely to land: a small stratum's risk set thins far out in the tail, long
    after the figures have stopped.
    """
    section("observation gaps past the plot cap (synthetic)")

    def bundle_with(rows) -> Bundle:
        columns = ["analysis", "stratum", "gap_start_day", "gap_end_day"]
        return Bundle(root=Path("."), data={
            "settings": {
                "restricted_stay_cap": 365,
                "presentation": {"plot_stay_cap": 60},
                "coverage": {"period": {"label": "Period", "included": True},
                             "group": {"label": "Animal Group", "included": True}},
            },
            # Presence is all these decide: which stratifiers exist, and which
            # of them the opening summary tallies (see summary_fields).
            "strata": {"period": {}, "all": {"unified_stay_counts": {}},
                       "group": {"unified_stay_counts": {}}},
            "unified": {"gaps": {
                "type": "table",
                "columns": {name: [row[i] for row in rows]
                            for i, name in enumerate(columns)}}},
        })

    vocab = Vocabulary()
    beyond = blocks.sample_gap_notes(
        bundle_with([["KM by Animal Group", "LRG", 212, 240]]), vocab)
    expect_equal("a gap past the plot cap is one note", len(beyond), 1)
    said = beyond[0]
    expect("it says there is none in the plotted range",
           "none inside the plotted range" in said, said)
    expect("it names the level and where the gap starts",
           "LRG (animal group)" in said and "day 212" in said, said)
    # The whole difference between the two forms: no day range, because the
    # end of a gap no figure reaches is not something a reader can act on.
    expect("it does not spell out the gap's extent", "day 240" not in said, said)
    expect("it says what still runs out that far",
           "365-day stay cap" in said and "60-day plot cap" in said, said)

    mixed = blocks.sample_gap_notes(
        bundle_with([["Unified KM", "(all data)", 12, 14],
                     ["KM by Animal Group", "LRG", 212, 240]]), vocab)
    expect_equal("a gap on each side of the plot cap is two notes", len(mixed), 2)
    expect("the one in the figures leads", "day 12 to day 14" in mixed[0], mixed[0])
    expect("the one past them follows as an aside",
           mixed[1].startswith("Further gaps"), mixed[1])

    # The periods are scanned as their own stratifier and reported on the title
    # slide, in the same two forms.
    window = blocks.window_gap_notes(
        bundle_with([["KM by Period", "2021", 30, 33]]), vocab)
    expect_equal("a period gap is the title slide's to report", len(window), 1)
    expect("and it is detailed, being inside the plot",
           "2021 (period), day 30 to day 33" in window[0], window[0])


def check_findings(case: str, bundle: Bundle, stratifier: str,
                      full: blocks.Table) -> None:
    lines = findings_for_stratifier(bundle, stratifier, full, stratifier)

    expect(f"{case}: at most one finding per slide ({stratifier})",
           len(lines) <= 1, f"got {len(lines)}")

    stays = full.df["km_restricted_mean"].dropna()
    flat = stays.empty or stays.min() == stays.max()
    if len(full.df) < 2 or flat:
        expect(f"{case}: nothing found from a flat or single level ({stratifier})",
               not lines, f"got {lines}")
        return

    expect(f"{case}: a finding names the longest-stay level ({stratifier})",
           bool(lines) and str(stays.idxmax()) in lines[0], f"got {lines}")
    expect(f"{case}: a finding is a finished sentence ({stratifier})",
           bool(lines) and lines[0].endswith("."), f"got {lines}")

    # Shares are only meaningful where the levels partition the population at a
    # moment. Periods partition the calendar instead, so their censuses sum to
    # roughly one census per period and a share would be nonsense.
    census = full.df["expected_census"].dropna()
    if bundle.has("strata", "all", "census") and not census.empty:
        total = float(bundle.stratum("all", "census")["expected_census"].iloc[0])
        partitions = total and abs(census.sum() - total) <= 0.01 * total
        expect(f"{case}: a share of residents appears only when levels "
               f"partition the population ({stratifier})",
               bool(partitions) == ("share of residents" in lines[0]),
               f"sum={census.sum():.1f} total={total:.1f} line={lines[0]}")


def check_overall_slide(case: str, bundle: Bundle) -> None:
    """The whole-sample slide: present, leading, and one row wide.

    Its finding is the length-bias contrast, which needs both medians. The
    median LOS is NA on a fixture whose curve never falls to 0.5, so the
    silence is asserted against whichever case the bundle actually is rather
    than assumed.
    """
    from mlos_review.deck import los_overall
    from mlos_review.names import Vocabulary as _Vocab

    if not requires_full_table(bundle, "all"):
        return
    slide = los_overall(bundle, _Vocab(bundle.data))
    expect(f"{case}: the whole sample gets a slide", slide is not None)
    expect_equal(f"{case}: its table is one row", len(slide.table.df), 1)
    expect(f"{case}: a single row carries no high/low flags",
           slide.table.flags is None)
    expect(f"{case}: it leads the presentation order",
           bundle.stratifiers()[0] == "all", f"order={bundle.stratifiers()}")

    # The standing facts are one function, not prose copied into two rules, so
    # the whole-sample slide and a stratified one must end identically. A rule
    # that re-wrote any of them by hand would drift the moment one was edited.
    from mlos_review.blocks import los_standing_notes
    from mlos_review.deck import los_by_stratifier

    tail = los_standing_notes(bundle, "all", _Vocab(bundle.data))
    expect_equal(f"{case}: the whole-sample slide ends with the shared notes",
                 slide.notes[-len(tail):], tail)
    levelled = next((los_by_stratifier(bundle, st, _Vocab(bundle.data))
                     for st in bundle.stratifiers() if st != "all"), None)
    if levelled is not None:
        expect_equal(f"{case}: and a stratified slide ends with the same ones, "
                     f"bar the data-conditional caveat",
                     [n for n in slide.notes if n in tail
                      and "Cap sensitivity" not in n],
                     [n for n in levelled.notes if n in tail
                      and "Cap sensitivity" not in n])

    row = slide.table.df.iloc[0]
    both = all(row.get(m) == row.get(m) and row.get(m) is not None
               for m in ("km_median_los",
                         "per_resident_past_days_restricted_median"))
    expect_equal(f"{case}: the length-bias finding appears only with both "
                 f"medians", bool(slide.findings), both)
    if both:
        expect(f"{case}: the finding cites both medians",
               f"{row['km_median_los']:.0f}" in slide.findings[0]
               and f"{row['per_resident_past_days_restricted_median']:.0f}"
               in slide.findings[0],
               f"got {slide.findings}")


def check_opening(case: str, bundle: Bundle) -> None:
    """The title slide and the descriptive summary that follows it.

    Every assertion here is that the summary agrees with the bundle it claims
    to describe. These are the only numbers in the deck a reader has no other
    way to check: everything later is a statistic they can look up in the
    workbook, and these are counts of the sample that produced it.
    """
    from datetime import datetime

    from mlos_review.deck import (
        AUTOMATION_CAVEAT, OPENING_BULLETS, OPENING_DIAGRAM, opening_slides,
        provenance, summary_slide,
    )

    vocab = Vocabulary(bundle.data)
    opening = opening_slides(bundle, vocab)

    title = opening[0]
    expect_equal(f"{case}: the deck opens on the title slide", title.layout, "TITLE")
    expect_equal(f"{case}: the title slide carries both bullets",
                 title.bullets, list(OPENING_BULLETS))
    # The caveat is a bullet here and a footnote on the findings slide, so it
    # is worth saying that the bullet is the same string: two placements of one
    # sentence, not two sentences that happen to agree today.
    expect(f"{case}: one of them is the automation caveat",
           AUTOMATION_CAVEAT in title.bullets)

    # The truncation and censoring diagram is the one figure that comes from
    # the package rather than from the run, so it is on every deck regardless
    # of what this case computed. Asserted against the file on disk, since a
    # constant pointing at nothing is exactly the failure to catch.
    expect(f"{case}: the diagram ships with the package", OPENING_DIAGRAM.exists(),
           str(OPENING_DIAGRAM))
    expect_equal(f"{case}: the title slide carries the diagram",
                 title.figures, [OPENING_DIAGRAM])

    # Two stamps, both at second resolution, and the analysis one taken from
    # the bundle rather than from the clock this deck was built on.
    stamp = provenance(bundle, datetime(2020, 5, 6, 7, 8, 9))
    expect(f"{case}: the footnote dates the deck", "2020-05-06 07:08:09" in stamp)
    generated = bundle.value("run", "generated_at")
    if generated:
        expect(f"{case}: the footnote dates the statistics",
               str(generated)[:19] in stamp, stamp)
        # The version lead is stripped before this one rather than the pattern
        # loosened to tolerate it: "mLOS 0.1.0." is full of dots followed by
        # digits, and a pattern that lets those through would also let a
        # single-digit fraction of a second through.
        expect(f"{case}: the footnote drops fractions of a second",
               re.search(r"\.\d", re.sub(r"^mLOS \S+\. ", "", stamp)) is None,
               stamp)
    version = bundle.value("run", "mlos_version")
    if version:
        expect(f"{case}: the footnote names the version that computed it",
               stamp.startswith(f"mLOS {version}. "), stamp)

    # The study window rides on the title slide, and lists no period the rest
    # of the deck never shows.
    expect_equal(f"{case}: the title slide carries the study window",
                 [table.title for table in title.tables],
                 ["study window"] if blocks.requires_study_window(bundle) else [])
    if title.tables and "period" in bundle.data.get("strata", {}):
        absent = (set(title.tables[0].df.index) - set(bundle.levels("period"))
                  - set(bundle.levels("all")))
        expect(f"{case}: the window lists no period the deck never shows",
               not absent, f"{sorted(absent)}")

    summary = summary_slide(bundle, vocab)
    if summary is None:
        return

    # The outer index level groups the rows the footnote governs. Dropped for
    # the value assertions, which are about the rows themselves, and asserted
    # separately below.
    grouped = summary.tables[0].df
    frame = grouped.droplevel(0)
    stays = bundle.table("strata", "all", "unified_stay_counts")["total_stays"].iloc[0]
    expect_equal(f"{case}: the sample table counts every stay once",
                 frame.loc[capitalize_first(vocab.metric("total_stays").label),
                           "count"], stays)

    groups = list(dict.fromkeys(grouped.index.get_level_values(0)))
    expect_equal(f"{case}: the sample rows come before the outcome rows",
                 groups,
                 [blocks.SAMPLE_GROUP, blocks.OUTCOME_GROUP]
                 if bundle.has("strata", "all", "outcomes") else [blocks.SAMPLE_GROUP])
    expect(f"{case}: the footnote sits under the outcome rows",
           summary.tables[0].footnotes == [blocks.SAMPLE_FOOTNOTE])

    animals = capitalize_first(vocab.metric("n_animals").label)
    supplied = bool(bundle.value("data_preparation", "animal_ids_supplied",
                                 default=False))
    expect_equal(f"{case}: distinct animals appear only where IDs were supplied",
                 animals in frame.index, supplied)
    if supplied:
        expect_equal(f"{case}: the animal count is the bundle's",
                     frame.loc[animals, "count"],
                     bundle.value("data_preparation", "n_animals"))

    # The cap is a duration in a column of tallies, so its row carries the unit
    # the vocabulary holds for it rather than leaving 365 to be read as a count.
    cap = bundle.value("settings", "restricted_stay_cap")
    if cap:
        label = f"{capitalize_first(vocab.metric('restricted_stay_cap').label)} (days)"
        expect_equal(f"{case}: the stay cap is a row, in days", frame.loc[label, "count"], cap)

    if bundle.has("strata", "all", "outcomes"):
        recorded = bundle.stratum("all", "outcomes")
        level = recorded.index[0]
        codes = [c[len("outcome_"):-len("_events")] for c in recorded.columns
                 if c.startswith("outcome_") and c.endswith("_events")]
        outcomes = [row for row in frame.index
                    if row == "Any" or row.split(" ")[0] in codes]
        expect_equal(f"{case}: the pooled total leads the outcome rows",
                     outcomes[0], "Any")
        expect_equal(f"{case}: Any is the analysis's own completed total",
                     frame.loc["Any", "count"],
                     recorded.loc[level, "completed_outcomes_total"])
        # Left blank rather than filled with the 1.0 it would be by
        # construction: every other share on this slide came from the bundle.
        expect(f"{case}: Any claims no share of itself",
               frame.loc["Any", "share"] != frame.loc["Any", "share"])
        expect_equal(f"{case}: one row per recorded outcome", len(outcomes) - 1, len(codes))
        for row in outcomes[1:]:
            code = row.split(" ")[0]
            expect_equal(f"{case}: {code} is counted as the bundle counts it",
                         frame.loc[row, "count"],
                         recorded.loc[level, f"outcome_{code}_events"])

    # Read straight from the bundle rather than through the block, so an error
    # in the block cannot agree with itself.
    def stays_by_level(stratifier: str) -> dict:
        counts = bundle.table("strata", stratifier, "unified_stay_counts")
        return {str(level): count for level, count
                in zip(counts["period_label"], counts["total_stays"])}

    # Levels are tallied only for fields few enough to plot. A field past the
    # limit is left off rather than listed, which is the only thing the slide
    # branches on.
    wanted = blocks.summary_fields(bundle)
    fits = [s for s in wanted
            if len(bundle.levels(s)) <= blocks.max_listed_levels(bundle)]
    expect_equal(f"{case}: levels are tallied only where they fit",
                 [table.title for table in summary.tables[1:]],
                 [f"stays by {vocab.stratifier(s).label}" for s in fits])
    for table, stratifier in zip(summary.tables[1:], fits):
        expect_equal(f"{case}: {stratifier} lists every level",
                     list(table.df.index), bundle.levels(stratifier))
        expect_equal(f"{case}: {stratifier} lists the bundle's own counts",
                     list(table.df["total_stays"]),
                     [stays_by_level(stratifier)[lvl]
                      for lvl in bundle.levels(stratifier)])

    check_field_summary(case, bundle, vocab, wanted, stays_by_level)


def check_field_summary(case: str, bundle: Bundle, vocab, wanted, stays_by_level) -> None:
    """`field_summary_table`, which no rule builds today.

    Kept for the dataset whose fields have more levels than a slide can list,
    and exercised here for the reason QUADRANTS is exercised in the renderer
    test: kept code that nothing calls and nothing tests is kept code that has
    quietly stopped working.
    """
    if not blocks.requires_field_summary(bundle):
        expect(f"{case}: no field summary without fields", not wanted)
        return

    summary = blocks.field_summary_table(bundle, vocab).df
    expect_equal(f"{case}: every field is summarized", len(summary), len(wanted))
    for row, stratifier in zip(summary.index, wanted):
        levels = bundle.levels(stratifier)
        expect_equal(f"{case}: {stratifier} reports its level count",
                     summary.loc[row, "categories"], len(levels))
        stays = stays_by_level(stratifier)
        if len(set(stays.values())) > 1:
            for end, want in (("smallest", min(stays.values())),
                              ("largest", max(stays.values()))):
                named = summary.loc[row, end]
                expect(f"{case}: {stratifier}'s {end} level is one of its own",
                       named in levels, f"{named!r}")
                expect_equal(f"{case}: {stratifier}'s {end} holds that many stays",
                             summary.loc[row, f"{end}_stays"], want)
                expect_equal(f"{case}: {stratifier}'s {end} is the level with "
                             f"that many", stays[named], want)
        else:
            expect(f"{case}: {stratifier} names no extreme where its levels tie",
                   summary.loc[row, "smallest"] != summary.loc[row, "smallest"])


def check_deck(case: str, bundle: Bundle, directory: Path) -> list:
    from pptx import Presentation

    from mlos_review.deck import (AUTOMATION_CAVEAT, findings_section,
                                  los_by_stratifier,
                                  los_overall, recommendations_section,
                                  resident_outlook,
                                  resident_outlook_by_stratifier)
    from mlos_review.names import Vocabulary as _Vocab
    from mlos_review.settings import Settings as _Settings

    vocab = _Vocab(bundle.data)
    # Mirrors build()'s dispatch, settings resolution included: the whole sample
    # takes its own two rules, in the position the bundle gives it, and every
    # other stratifier the by-level pair, the second of which declines unless
    # the stratifier is salient. A single-stratifier bundle makes it salient by
    # default, which is why the settings are read against the dataset here too.
    deck_settings = _Settings().for_dataset(
        [st for st in bundle.stratifiers() if st != "all"])

    def los_section(stratifier: str) -> list:
        if stratifier == "all":
            return [los_overall(bundle, vocab), resident_outlook(bundle, vocab)]
        return [los_by_stratifier(bundle, stratifier, vocab, deck_settings),
                resident_outlook_by_stratifier(bundle, stratifier, vocab,
                                               deck_settings)]

    section_slides = [
        s for st in bundle.stratifiers() for s in los_section(st) if s is not None
    ]
    # Then the workload slide, which reads across every stratifier and so can
    # only follow all of them.
    from mlos_review.deck import workload_slide, WORKLOAD_SLIDES
    section_slides += [s for s in
                       (workload_slide(bundle, vocab, section, deck_settings)
                        for section in WORKLOAD_SLIDES)
                       if s is not None]
    # Then the two ratio slides per stratifier, which close the LOS section:
    # hazard first and days second, none for `all`, which is a single level
    # with nothing to hold a ratio against.
    from mlos_review.deck import (cox_comparison_by_stratifier,
                                  hazard_ratios_by_stratifier,
                                  los_ratios_by_stratifier, reserve_section)
    from mlos_review.blocks import (findings_for_cox_comparison,
                                    findings_for_parametric_agreement,
                                    requires_cox_comparison)
    from mlos_review.figures import FigureSet
    from mlos_review.recommend import unestimable_levels
    from mlos_review.regression import (comparison as _comparison, pooled,
                                        stratified, hazard_ratio_panel,
                                        los_ratio_panel, panel_stratifiers,
                                        HR_SERIES)
    cox = _comparison(pooled(bundle), stratified(bundle))
    hr_panel = hazard_ratio_panel(bundle)
    los_panel = los_ratio_panel(bundle)
    with tempfile.TemporaryDirectory() as figure_tmp:
        figures = FigureSet(directory=Path(figure_tmp))
        wanted = [st for st in panel_stratifiers(hr_panel, HR_SERIES)
                  if st != "all"]
        # Every hazard ratio, then every length-of-stay ratio, which is the
        # order build() assembles them in and the order the notes argue for.
        for st in wanted:
            # The retired robustness slide's sentences ride the hazard-ratio
            # slide now, and so does the parametric-agreement one. They have to
            # ride the mirror's too: they are what the closing sections are
            # counted from.
            label = vocab.stratifier(st).label
            slide = hazard_ratios_by_stratifier(
                bundle, st, hr_panel, vocab, figures, deck_settings,
                findings=(findings_for_parametric_agreement(hr_panel, st, label)
                          + (findings_for_cox_comparison(cox, st, label)
                             if requires_cox_comparison(cox, st) else [])),
                recommendations=unestimable_levels(bundle, st, cox, vocab))
            if slide is not None:
                section_slides.append(slide)
        for st in wanted:
            slide = los_ratios_by_stratifier(bundle, st, los_panel, vocab,
                                             figures, deck_settings)
            if slide is not None:
                section_slides.append(slide)

    # Then the competing-risks section in the same order build() assembles it:
    # the teaser, then the one stratified slide the default TEASER coverage
    # carries. Both contribute to the closing slide, so this list has to match
    # or the count below counts the wrong findings.
    from mlos_review.blocks import requires_aj_teaser
    from mlos_review.deck import (_aj_figures, aj_by_stratifier, aj_teaser,
                                  aj_teaser_stratifier, resident_destination)

    pick = aj_teaser_stratifier(bundle, deck_settings)
    aj_slides = [aj_teaser(bundle, vocab)] if requires_aj_teaser(bundle) else []
    # Behind the teaser and ahead of any breakdown, the way build() places it.
    aj_slides.append(resident_destination(bundle, vocab))
    if pick is not None:
        aj_slides.extend(aj_by_stratifier(bundle, pick, vocab))
    aj_slides = [s for s in aj_slides if s is not None]

    section_slides = section_slides + aj_slides
    # Plus the opening, which is not a section: a title slide always, and the
    # descriptive slides for whichever of them this bundle has the numbers for.
    from mlos_review.deck import opening_slides
    opening = opening_slides(bundle, vocab)

    # BOTH closing sections gather over the opening too, not only the sections:
    # the observation-gap finding and its remedy ride the title slide, and a
    # mirror that gathered over the sections alone would lose them and then
    # disagree with the real deck about how many pages it takes.
    built = opening + section_slides
    gathered = [line for s in built for line in s.findings]
    suggested = [line for s in built for line in s.recommendations]
    finding_pages = findings_section(built)
    recommendation_pages = recommendations_section(built)
    closing = finding_pages + recommendation_pages
    # The reserve pages come after the closing sections and contribute nothing
    # to either, which is what lets them carry the retired robustness slides
    # without their sentences being counted a second time.
    with tempfile.TemporaryDirectory() as figure_tmp:
        reserve = reserve_section(bundle, cox, vocab,
                                  FigureSet(directory=Path(figure_tmp)),
                                  deck_settings)
    expect_equal(f"{case}: the reserve slides gather nothing",
                 [line for s in reserve for line in s.findings + s.recommendations],
                 [])
    # And the educational section, last of all and gathering nothing either.
    # Counted here rather than left out, so a slide added to it has to be
    # accounted for in the deck's length like every other.
    from mlos_review.deck import educational_section
    educational = educational_section(bundle, vocab)
    expect_equal(f"{case}: the educational slides gather nothing",
                 [line for s in educational
                  for line in s.findings + s.recommendations],
                 [])
    expected = (len(opening) + len(section_slides) + len(closing)
                + len(reserve) + len(educational))

    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "deck.pptx"
        written, archived = build(directory, out)
        expect(f"{case}: deck written", out.exists())
        expect_equal(f"{case}: build reports where it wrote", written, out)
        expect(f"{case}: nothing was archived over an empty directory",
               archived is None, f"got {archived}")
        deck = Presentation(str(out))
        expect_equal(f"{case}: one slide per usable stratifier",
                     len(deck.slides), expected)
        rendered = []
        for index, slide in enumerate(deck.slides):
            texts = [
                shape.text_frame.text for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            rendered.append(texts)
            expect(f"{case}: slide {index + 1} has a title", bool(texts))
            expect(f"{case}: slide {index + 1} has speaker notes",
                   slide.has_notes_slide
                   and bool(slide.notes_slide.notes_text_frame.text.strip()))

        # The automation caveat, checked on the rendered file because where it
        # lands is the requirement: the title slide, and the head of the
        # Findings section, which is the page that says what the deck concluded
        # and the one whose selection a reader cannot check against anything in
        # front of them. Twice and no more. Not the continuation pages, and not
        # every slide, which would be a caveat nobody reads.
        carrying = [i + 1 for i, texts in enumerate(rendered)
                    if any(AUTOMATION_CAVEAT in text for text in texts)]
        expect(f"{case}: the title slide carries the automation caveat",
               1 in carrying, f"carried by slides {carrying}")
        if finding_pages:
            findings_at = len(opening) + len(section_slides) + 1
            expect_equal(f"{case}: and the Findings section's first page, "
                         f"and no other", carrying, [1, findings_at])
            expect_equal(f"{case}: which is the page it heads",
                         _slide_title(deck.slides[findings_at - 1]), "Findings")
        else:
            expect_equal(f"{case}: a deck that found nothing caveats once",
                         carrying, [1])


    # Findings are carried by the slides but never printed on them; they
    # reach the reader only through the closing section the deck assembles.
    # Recommendations follow the findings, in their own section: a reader
    # has to be able to tell a finding from a suggestion without a tag on the
    # line, which is the whole reason they are two fields and two sections.
    expect_equal(f"{case}: recommendations are gathered separately, and last",
                 [line for page in recommendation_pages for line in page.bullets],
                 suggested)
    expect(f"{case}: a section slide never displays its own recommendations",
           not any(line in s.bullets for s in section_slides
                   for line in s.recommendations))
    if gathered:
        expect(f"{case}: the closing section collects every finding",
               [line for page in finding_pages for line in page.bullets]
               == gathered)
        expect(f"{case}: a section slide never displays its own findings",
               all(not s.bullets for s in section_slides))
        expect(f"{case}: every closing page is titled",
               all(page.title for page in closing))
    else:
        expect(f"{case}: no closing pages when nothing was found",
               closing == [])

    # The mirror, handed on. It is the deck's slides as objects, which is what
    # a variant deck composes from, and what check_slide_titles_addressable
    # needs in order to ask its question of the titles rather than of a
    # rendered file.
    return opening + section_slides + closing + reserve + educational


def check_slide_titles_addressable(case: str, slides: list) -> None:
    """Every slide the deck builds can be named, once, from an outline.

    A variant deck addresses a base slide by its exact title, so a title used
    twice is a slide nobody can ask for. The one intended exception is a
    continuation run, which `_gathered_section` names `T` then `T, continued`
    repeated and unnumbered: the run is addressed as a set, through its head,
    because which sentence lands on which page is decided by height.

    This is the contract a variant relies on, checked where the slides are
    built rather than trusted.
    """
    from mlos_review.variant import CONTINUATION, compose, parse_outline, runs

    orphans = [index + 1 for index, slide in enumerate(slides)
               if slide.title.endswith(CONTINUATION)
               and (index == 0
                    or slides[index - 1].title not in
                    (slide.title, slide.title[:-len(CONTINUATION)]))]
    expect_equal(f"{case}: no continuation page without its head",
                 orphans, [])

    seen, repeated = set(), []
    for slide in slides:
        if slide.title.endswith(CONTINUATION):
            continue
        if slide.title in seen:
            repeated.append(slide.title)
        seen.add(slide.title)
    expect_equal(f"{case}: every slide title is addressable, once",
                 sorted(repeated), [])
    if repeated:
        return

    # The round trip, which is the assertion that actually pins the contract:
    # what an outline gets back for a title is the whole run, in order.
    grouped = runs(slides)
    for head, run in grouped.items():
        borrowed = compose(parse_outline(f"@insert {head}", "test"), slides)
        expect_equal(f"{case}: @insert {head!r} borrows its whole run",
                     [s.title for s in borrowed], [s.title for s in run])
    expect_equal(f"{case}: the runs account for every slide",
                 sum(len(run) for run in grouped.values()), len(slides))


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------


def check_margin_takes_the_further_bound() -> None:
    """The margin is set by the bound further from 1, not the larger bound.

    No fixture pins this, because whether an interval sits off center is a
    property of one dataset's arithmetic rather than of any shape a fixture
    was built for, and the two readings agree whenever it does not. They come
    apart in two ways, and both are silent: on an interval whose LOWER bound is
    the further, the signed maximum understates the margin, and on one lying
    entirely below 1 it returns a negative margin and claims agreement within
    less than nothing.

    Built by hand from two hazard ratios rather than from a bundle, so the
    expected value can be worked out on paper: HRs of 1.0 and 1.1 with equal
    standard errors of 0.05 give a log difference of 0.0953 and a half-width of
    1.96 * sqrt(2) * 0.05 = 0.1386, so the interval runs from exp(-0.0433) to
    exp(0.2339) and the margin is exp(0.0953 + 0.1386) - 1 = 26.4%, set by the
    upper bound. Flip the two hazard ratios and the interval mirrors, the lower
    bound becomes the further one, and the margin has to stay 26.4%.
    """
    import math

    from mlos_review import regression as reg

    def frame(hr, se):
        index = pd.MultiIndex.from_tuples([("group", "A")], names=reg.INDEX_NAMES)
        return pd.DataFrame({"hr": [hr],
                             "hr_ci_lower": [hr * math.exp(-reg.Z_95 * se)],
                             "hr_ci_upper": [hr * math.exp(reg.Z_95 * se)],
                             "p_value": [0.5]}, index=index)

    expected = math.exp(math.log(1.1) + reg.Z_95 * math.hypot(0.05, 0.05)) - 1

    rising = reg.comparison(frame(1.0, 0.05), frame(1.1, 0.05))
    expect(f"margin from the upper bound: {rising.iloc[0]['agreement_margin']:.4f}",
           math.isclose(rising.iloc[0]["agreement_margin"], expected, rel_tol=1e-9))

    # The mirror image: same two fits, swapped. A signed maximum would report
    # a smaller margin here, and on a wide enough interval a negative one.
    falling = reg.comparison(frame(1.1, 0.05), frame(1.0, 0.05))
    expect(f"margin from the lower bound: {falling.iloc[0]['agreement_margin']:.4f}",
           math.isclose(falling.iloc[0]["agreement_margin"], expected, rel_tol=1e-9))

    expect("the mirrored pair earns the same margin",
           math.isclose(rising.iloc[0]["agreement_margin"],
                        falling.iloc[0]["agreement_margin"], rel_tol=1e-12))


def check_comparison_declines_without_both_fits() -> None:
    """A stratifier with only the pooled fit yields no comparison and no slide.

    The ordinary case for a run with one usable stratifier, and for one whose
    variant failed to converge. What must not happen is a slide of blanks: the
    pooled column would fill, the other two would not, and the table would look
    like a result that had come out empty rather than one never attempted.
    """
    from mlos_review.blocks import requires_cox_comparison
    from mlos_review import regression as reg

    index = pd.MultiIndex.from_tuples([("group", "A"), ("group", "B")],
                                      names=reg.INDEX_NAMES)
    pooled = pd.DataFrame({"hr": [1.0, 1.4], "hr_ci_lower": [None, 1.2],
                           "hr_ci_upper": [None, 1.6], "p_value": [None, 0.01]},
                          index=index)
    empty = pd.DataFrame(columns=reg.HR_COLUMNS,
                         index=pd.MultiIndex.from_arrays([[], []], names=reg.INDEX_NAMES),
                         dtype=float)

    comparison = reg.comparison(pooled, empty)
    expect_equal("no comparison without a stratified fit",
                 reg.comparable_stratifiers(comparison), [])
    expect("the pooled column still fills",
           comparison["hr_pooled"].notna().all())
    expect("no margin is invented for it",
           comparison["agreement_margin"].isna().all())
    expect("and no slide is built",
           not requires_cox_comparison(comparison, "group"))


def check_workbook_and_manifest(case: str, bundle: Bundle, directory: Path) -> None:
    """The workbook and the figure manifest one build writes.

    Runs a real build into a temporary directory, because the point is what
    reaches disk: sheets that are correct in memory and a workbook Excel will
    not open are the same failure to a reader.
    """
    import json as _json

    from openpyxl import load_workbook

    from mlos_review import workbook
    from mlos_review.deck import build, figure_directory, workbook_path

    vocab = Vocabulary(bundle.data)
    sheets = workbook.sheets(bundle, vocab)

    names = [sheet.name for sheet in sheets]
    expect_equal(f"{case}: sheet names are unique", len(set(names)), len(names))
    for name in names:
        expect(f"{case}: {name} is a legal sheet name",
               name == workbook.sheet_name(name), f"would become {workbook.sheet_name(name)}")
    for sheet in sheets:
        expect(f"{case}: {sheet.name} carries at least one table", bool(sheet.tables))
        # The stacking rule: tables share a sheet only when their columns match,
        # which is what lets a reader scan down one set of headings.
        columns = {tuple(table.df.columns) for table in sheet.tables}
        expect_equal(f"{case}: {sheet.name} stacks only matching columns", len(columns), 1)

    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "deck.pptx"
        build(directory, out)

        tables = workbook_path(out)
        expect_equal(f"{case}: a workbook is written exactly when there are sheets",
                     tables.exists(), bool(sheets))
        if tables.exists():
            book = load_workbook(tables)
            expect_equal(f"{case}: every sheet reaches the file",
                         book.sheetnames, names)
            for sheet in sheets:
                worksheet = book[sheet.name]
                # Heading rows, one blank line between tables, then the rows.
                # Counted rather than parsed: the arithmetic is what a stacked
                # sheet gets wrong, by dropping a gap or a heading.
                heading = 1 if len(sheet.tables) > 1 else 0
                wanted = sum(heading + 1 + len(table.df) + len(table.footnotes)
                             for table in sheet.tables) + (len(sheet.tables) - 1)
                expect_equal(f"{case}: {sheet.name} has the rows it should",
                             worksheet.max_row, wanted)
                expect(f"{case}: {sheet.name} says something",
                       worksheet.max_column >= 2)

        # Every figure this package drew is in the manifest, and every entry
        # names a file that is really there. A manifest is worth nothing if it
        # can drift from the directory it describes.
        manifest = figure_directory(out) / "manifest.json"
        pngs = sorted(p.name for p in figure_directory(out).glob("*.png")) \
            if figure_directory(out).exists() else []
        expect_equal(f"{case}: a manifest is written exactly when figures are",
                     manifest.exists(), bool(pngs))
        if manifest.exists():
            entries = _json.loads(manifest.read_text())["figures"]
            expect_equal(f"{case}: the manifest lists every PNG",
                         sorted(entry["plot"] for entry in entries), pngs)
            for entry in entries:
                expect(f"{case}: {entry['plot']} is described",
                       bool(entry["description"]) and bool(entry["kind"]))
                expect(f"{case}: {entry['plot']} names a real stratifier",
                       entry["stratifier"] in bundle.stratifiers())


def check_ratio_panel_sheets(case: str, bundle: Bundle) -> None:
    """The two ratio panels reach the workbook, and the slides are views of them.

    The check the workbook needed and did not have. Both panels were built for
    the slides and written nowhere, so the Weibull fits and the Kaplan-Meier
    ratio, which no other sheet carries, were missing from the companion
    entirely while every generic workbook assertion passed. What pins that shut
    is not "a sheet named HR_Panel exists" but the relation the workbook exists
    to keep: every number a slide shows is in the sheet, at the same value.
    """
    from mlos_review import regression as reg
    from mlos_review import workbook

    vocab = Vocabulary(bundle.data)
    by_name = {sheet.name: sheet for sheet in workbook.sheets(bundle, vocab)}

    for name, panel, series in (
            ("HR_Panel", reg.hazard_ratio_panel(bundle), reg.HR_SERIES),
            ("LOS_Ratio_Panel", reg.los_ratio_panel(bundle), reg.LOS_SERIES)):
        wanted = not panel.empty and bool(panel.notna().any().any())
        expect_equal(f"{case}: {name} is written exactly when the panel holds anything",
                     name in by_name, wanted)
        if not wanted:
            continue

        sheet = by_name[name]
        expect_equal(f"{case}: {name} carries one table", len(sheet.tables), 1)
        table = sheet.tables[0]
        # Estimate and both bounds for each reading, in the panel's own order.
        expect_equal(f"{case}: {name} carries every reading with its interval",
                     list(table.df.columns),
                     [column for key, _, _ in series
                      for column in (key, f"{key}_ci_lower", f"{key}_ci_upper")])
        check_formats(f"{case}: {name}", table)

        for stratifier in reg.panel_stratifiers(panel, series):
            slide = blocks.ratio_panel_table(panel, series, stratifier,
                                             title="", footnotes=[])
            for level in slide.df.index:
                for measure in slide.df.columns:
                    expect_equal(
                        f"{case}: {name} holds {stratifier}/{level} {measure}",
                        _same_number(table.df.at[(stratifier, level), measure],
                                     slide.df.at[level, measure]),
                        True)

    # The Weibull specifically, since its absence is what went unnoticed: a run
    # that fitted one must be able to show a reader its numbers.
    if bundle.value("weibull", "has_analysis"):
        columns = {column for sheet in by_name.values() for table in sheet.tables
                   for column in table.df.columns}
        for measure in ("hr_weibull_pooled", "los_weibull_pooled"):
            expect(f"{case}: the workbook carries {measure}", measure in columns)


def _same_number(left, right) -> bool:
    """Equal, counting two missing values as equal, which `==` does not."""
    if pd.isna(left) or pd.isna(right):
        return bool(pd.isna(left) and pd.isna(right))
    return bool(left == right)


def check_cox_regression_frames(case: str, bundle: Bundle) -> None:
    """The two regression frames, and the margin derived from them.

    Fixture-driven, so it runs over every bundle: ones with both fits, ones
    with only the pooled fit (fewer than two usable stratifiers), and ones with
    no Cox regression at all. The first assertion is the one the whole design
    rests on, that the two frames line up without anyone arranging it.
    """
    import math

    from mlos_review import regression as reg

    pooled = reg.pooled(bundle)
    stratified = reg.stratified(bundle)

    expect(f"{case}: regression frames share their columns",
           list(pooled.columns) == list(stratified.columns) == reg.HR_COLUMNS)
    if not stratified.empty:
        expect(f"{case}: pooled and stratified frames are the same shape",
               pooled.index.equals(stratified.index),
               f"{list(pooled.index)} vs {list(stratified.index)}")

    # Every level of every usable stratifier is present, reference included.
    # A frame that quietly dropped the reference would still look right in a
    # table and would leave the figure with no 1.0 bar to anchor on.
    if not pooled.empty:
        for stratifier in pooled.index.get_level_values("stratifier").unique():
            expect_equal(f"{case}: {stratifier} keeps every level",
                         list(pooled.loc[stratifier].index), bundle.levels(stratifier))

    comparison = reg.comparison(pooled, stratified)
    for key in comparison.index:
        row = comparison.loc[key]
        name = f"{case}: {key[0]}/{key[1]}"

        if pd.isna(row["hr_stratified"]) or pd.isna(row["agreement_margin"]):
            # Either no stratified fit, or the reference level. Both are
            # ordinary; what must not happen is a margin without a ratio, or
            # one comparison column filled while the other is not.
            expect(f"{name}: no margin without both fits",
                   pd.isna(row["agreement_margin"]) or pd.notna(row["hr_stratified"]))
            expect(f"{name}: the two comparison columns stand or fall together",
                   pd.isna(row["p_difference"]) == pd.isna(row["agreement_margin"]))
            continue

        lower, upper = row["hr_ratio_ci_lower"], row["hr_ratio_ci_upper"]
        expect(f"{name}: the interval brackets the ratio",
               lower <= row["hr_ratio"] <= upper,
               f"{lower} .. {row['hr_ratio']} .. {upper}")

        # THE defining property of the achieved equivalence margin: the
        # interval for the ratio lies inside [1/(1+m), 1+m], and no smaller
        # symmetric margin would contain it. Both halves are checked, because
        # the first alone is satisfied by any margin large enough.
        bound = 1 + row["agreement_margin"]
        expect(f"{name}: the margin contains the interval",
               lower >= 1 / bound - 1e-12 and upper <= bound + 1e-12,
               f"[{lower}, {upper}] outside [{1 / bound}, {bound}]")
        expect(f"{name}: the margin is the tightest one that does",
               math.isclose(max(abs(math.log(lower)), abs(math.log(upper))),
                            math.log(bound), rel_tol=1e-9),
               f"margin {row['agreement_margin']} on [{lower}, {upper}]")

        expect(f"{name}: the margin covers the observed gap",
               row["agreement_margin"] >= abs(row["hr_ratio"] - 1) - 1e-12)

        # The p-value and the interval are two readings of one z statistic at
        # one confidence level, so they cannot disagree: p < 0.05 exactly when
        # the interval for the ratio excludes 1. Worth pinning because they are
        # computed on separate lines and a future edit could change the
        # constant in one without the other, which no eyeball would catch.
        p = row["p_difference"]
        expect(f"{name}: the p-value is a probability", 0.0 <= p <= 1.0, f"got {p}")
        if abs(p - 0.05) > 1e-9:
            expect(f"{name}: the p-value and the interval agree",
                   (p < 0.05) == (lower > 1.0 or upper < 1.0),
                   f"p={p} against [{lower}, {upper}]")

    usable = reg.comparable_stratifiers(comparison)
    for stratifier in bundle.stratifiers():
        from mlos_review.blocks import requires_cox_comparison
        expect_equal(f"{case}: {stratifier} comparison availability",
                     requires_cox_comparison(comparison, stratifier),
                     stratifier in usable)
        if stratifier in usable:
            table = blocks.cox_comparison_table(comparison, stratifier)
            expect_equal(f"{case}: {stratifier} comparison table rows",
                         list(table.df.index), bundle.levels(stratifier))
            check_formats(case, table)


def run_fixture(case: str, directory: Path) -> None:
    section(f"fixture: {case}")
    bundle = Bundle.load(directory)
    run_check(check_cox_regression_frames, case, bundle)
    run_check(check_ratio_panel_sheets, case, bundle)
    run_check(check_workbook_and_manifest, case, bundle, directory)
    run_check(check_vocabulary, case, bundle)
    run_check(check_bundle_access, case, bundle)
    run_check(check_opening, case, bundle)
    run_check(check_observation_gaps, case, bundle)
    run_check(check_care_days, case, bundle)
    run_check(check_recommendations, case, bundle)
    run_check(check_aj_teaser, case, bundle)
    run_check(check_aj_teaser_pick, case, bundle)
    run_check(check_overall_slide, case, bundle)
    for stratifier in bundle.stratifiers():
        if stratifier != "all":
            run_check(check_aj_by_stratifier, case, bundle, stratifier)
            run_check(check_stacked_table, case, bundle, stratifier)
        # Presence-based dispatch: a bundle without the matrices a table needs
        # is not a failure, it is a run that legitimately has less in it.
        if not requires_full_table(bundle, stratifier):
            continue
        # Required: the checks below take its table, so under --only it runs
        # muted rather than handing them nothing.
        full = run_check(check_full_table, case, bundle, stratifier, required=True)
        if full is None:
            continue
        run_check(check_formats, case, full)
        run_check(check_at_cap_agrees_with_aj, case, bundle, stratifier, full)
        run_check(check_cap_sensitivity, case, bundle, stratifier, full)
        run_check(check_order_shifts, case, bundle, stratifier, full)
        if stratifier != "all":
            run_check(check_salience_note, case, bundle, stratifier)
            run_check(check_outcome_spread, case, bundle, stratifier)
            run_check(check_stratified_outlook, case, bundle, stratifier)
        run_check(check_highlights, case, bundle, stratifier, full)
        run_check(check_findings, case, bundle, stratifier, full)
    slides = run_check(check_deck, case, bundle, directory, required=True)
    if slides is not None:
        run_check(check_slide_titles_addressable, case, slides)
    run_check(check_deck_with_figures, case, bundle, directory)


# ---------------------------------------------------------------------------
# The documents: what two renderers make of them
# ---------------------------------------------------------------------------

# make_docx.sh is the declaration of which documents get a Word export; this
# list is checked against it below rather than trusted, the same way the Colab
# notebook is checked against pyproject.toml. README.md is deliberately absent:
# it is read on GitHub only, has no cross-references and no .docx, so the half
# of this that matters is the half it cannot fail.
DOCUMENTS = [
    "mlos_math_methods.md",
    "mlos_user_guide.md",
    "presentation_guide.md",
    "documentation_rules.md",
]

# Pandoc, searched the way make_docx.sh searches: an explicit override first,
# then the pandoc RStudio bundles through Quarto, then whatever is on PATH.
PANDOC_CANDIDATES = [
    "/Applications/RStudio.app/Contents/Resources/app/quarto/bin/tools/aarch64/pandoc",
    "/Applications/RStudio.app/Contents/Resources/app/quarto/bin/tools/x86_64/pandoc",
    "/Applications/RStudio.app/Contents/Resources/app/quarto/bin/tools/pandoc",
]


def _find_pandoc() -> str | None:
    override = os.environ.get("PANDOC")
    if override and Path(override).is_file() and os.access(override, os.X_OK):
        return override
    for candidate in PANDOC_CANDIDATES:
        if Path(candidate).is_file() and os.access(candidate, os.X_OK):
            return candidate
    return shutil.which("pandoc")


def _anchor(text: str) -> str:
    """GitHub's slug rule. Underscores survive it; other punctuation does not.

    Getting this wrong in the lenient direction is harmless and in the strict
    direction is not: a slug that drops underscores calls every settings
    heading in the User Guide (`animal_group_columns` and its forty siblings) a
    broken link, which is a wall of false failures nobody will read twice.
    """
    text = re.sub(r"`|\*|\(|\)|\.|,|:|;|'|\"", "", text.strip().lower())
    return re.sub(r"[^a-z0-9\-_ ]", "", text).replace(" ", "-")


def _heading_text(text: str) -> str:
    """Heading text as the reader sees it, with the markup and typography off.

    Pandoc pretty-prints a long heading across lines, drops code-span backticks
    and emphasis markers, and turns quotes into curly ones. None of those is a
    difference in the document, and comparing without normalizing them first
    makes every heading over about sixty characters look like a disagreement.
    """
    text = text.replace("`", "").replace("*", "")
    text = text.translate(str.maketrans({"\u2019": "'", "\u2018": "'",
                                         "\u201c": '"', "\u201d": '"'}))
    return re.sub(r"\s+", " ", text).strip()


def _blank_line_rules(name: str, source: str) -> None:
    """The three things pandoc needs a blank line above and GitHub does not.

    These need no renderer, which is why they are separate from the comparison
    below: they are the checks that still run on a machine with neither tool,
    and each of the three has actually shipped.

    A list or a heading without one is folded into the paragraph above, taking
    its bookmark with it, so every cross-reference to that section dies in
    Word. A `---` without one is worse, because it does not look like damage: a
    setext underline turns the line above into a heading, which then takes the
    anchor the real section wanted and pushes the real one to `<slug>-1`. The
    contents list of the presentation guide shipped that way in August 2026,
    from a script that regenerated the list and joined it straight onto the
    rule below. CommonMark reads the same three characters as a thematic break,
    so GitHub showed nothing wrong at all.
    """
    lines = source.split("\n")
    offenders = {"rule": [], "heading": [], "list": []}
    in_fence = False
    for i, line in enumerate(lines):
        if line.lstrip().startswith("```"):
            in_fence = not in_fence
            continue
        if in_fence or i == 0 or not lines[i - 1].strip():
            continue
        previous = lines[i - 1]
        if re.match(r"^-{3,}\s*$", line):
            offenders["rule"].append(i + 1)
        elif re.match(r"^#{1,6} ", line):
            offenders["heading"].append(i + 1)
        elif (re.match(r"^\s*[-*+] ", line)
              and not re.match(r"^\s*[-*+] ", previous)
              and not previous.startswith("  ")):
            offenders["list"].append(i + 1)

    for kind, lineno in offenders.items():
        expect(f"{name}: no {kind} without a blank line above it",
               not lineno, f"at line(s) {lineno}")


def _math_underscores(name: str, source: str) -> None:
    """The subscript spelling that eats itself, and only that one.

    GitHub runs its markdown pass over math content, where `}_` can open
    emphasis. ONE of those in a paragraph is safe, since nothing closes it; two
    open and close it and both subscripts vanish. Screening on underscore
    parity instead flags every legal subscript in the file, which on the math
    document is 62 paragraphs of 190, none of them wrong.
    """
    paragraphs = [p for p in source.split("\n\n") if "$" in p]
    if not paragraphs:
        return
    doubled = [p for p in paragraphs if p.count("}_") >= 2]
    expect(f"{name}: no paragraph carries two `}}_` subscripts",
           not doubled,
           f"{len(doubled)} of {len(paragraphs)} math paragraphs")


def _structure_agrees(name: str, source: str, pandoc: str | None) -> None:
    """What markdown-it and pandoc each made of the file, held against each other.

    markdown-it-py is CommonMark, which is the closest local stand-in for
    GitHub's renderer, and pandoc is the Word path. A construct both agree on
    survives the trip to both places. A construct only one of them sees is
    exactly where the two copies of a document come apart, and it is invisible
    from either side alone: the setext bug above looked perfect on GitHub and
    broke every link to the last section in Word.

    Counts, not content. This is not a golden of the rendered output, which
    would churn on every edit; it is the block structure, which only changes
    when the markdown changes meaning.
    """
    try:
        from markdown_it import MarkdownIt
    except ImportError:
        skipped(f"{name}: renderer comparison",
                "markdown-it-py is not installed (pip3 install markdown-it-py)")
        return
    if pandoc is None:
        skipped(f"{name}: renderer comparison",
                "no pandoc found; set PANDOC=/path/to/pandoc")
        return

    tokens = MarkdownIt("commonmark", {"html": True}).enable("table").parse(source)
    mdit_heads = [(int(t.tag[1]), _heading_text(tokens[i + 1].content))
                  for i, t in enumerate(tokens) if t.type == "heading_open"]
    # Fenced blocks are skipped, the way the sibling check above skips them. A
    # `#` line inside a fence is a heading nobody wrote: the variant outline
    # the guide shows is markdown ABOUT markdown, and a scan that read its
    # example as document headings would report three that markdown-it was
    # right not to parse.
    written, in_fence = [], False
    for line in source.split("\n"):
        if line.lstrip().startswith("```"):
            in_fence = not in_fence
            continue
        match = None if in_fence else re.match(r"^(#{1,6}) (.*)$", line)
        if match:
            written.append((len(match.group(1)), _heading_text(match.group(2))))

    # A heading in the file but not in the parse was swallowed by the block
    # above it, which is the blank-line rule failing where the naive scan
    # above cannot see it.
    expect_equal(f"{name}: every heading written is parsed as one",
                 [h for h in written if h not in mdit_heads], [])

    anchors = [_anchor(text) for _, text in mdit_heads]
    expect_equal(f"{name}: no two headings claim one anchor",
                 sorted({a for a in anchors if anchors.count(a) > 1}), [])
    targets = set(anchors)
    expect_equal(f"{name}: every internal link resolves",
                 sorted({link for link in re.findall(r"\]\(#([^)]+)\)", source)
                         if link not in targets}), [])

    # Fed on stdin rather than by path, so both renderers are demonstrably
    # reading the same bytes and this function can be exercised on a source
    # that is not on disk.
    result = subprocess.run([pandoc, "-f", "markdown", "-t", "html", "--mathjax"],
                            input=source, capture_output=True, text=True)
    expect(f"{name}: pandoc parses it without warnings", not result.stderr.strip(),
           result.stderr.strip()[:200])
    html = result.stdout

    pandoc_heads = [(int(level), _heading_text(re.sub(r"<[^>]+>", "", text)))
                    for level, text in
                    re.findall(r"<h([1-6])[^>]*>(.*?)</h\1>", html, re.S)]
    expect_equal(f"{name}: both renderers see the same headings",
                 ([h for h in mdit_heads if h not in pandoc_heads],
                  [h for h in pandoc_heads if h not in mdit_heads]),
                 ([], []))

    # An id pandoc had to disambiguate is a second element claiming a slug some
    # heading already owns, and the link goes to whichever came first.
    stolen = sorted({m.group(1) for m in
                     (re.match(r"^(.*)-\d+$", i) for i in
                      re.findall(r'id="([^"]*)"', html)) if m}
                    & targets)
    expect_equal(f"{name}: no anchor is taken by something that is not its section",
                 stolen, [])

    for block, opener, closer in (("bullet lists", "bullet_list_open", "<ul>"),
                                  ("tables", "table_open", "<table"),
                                  ("code blocks", "fence", "<pre")):
        expect_equal(f"{name}: both renderers see the same number of {block}",
                     sum(1 for t in tokens if t.type == opener),
                     html.count(closer))


# Two documents number their sections, and a § names neither: it is written
# bare, so it can only mean the one the writing cites. That is the math
# methods, and documentation_rules.md is numbered but never cited this way.
# Both halves are checked rather than assumed, so a third numbered document, or
# a § aimed at the rules, fails here and asks to be taught which is meant.
NUMBERED_DOCUMENTS = ["mlos_math_methods.md", "documentation_rules.md"]
SECTION_SYMBOL_MEANS = "mlos_math_methods.md"

# Where a §-reference can be written: everything in the repository a person
# writes prose into. The golden .txt and .csv under tests/cases are output
# rather than writing, and results/ is regenerable, so neither can hold a
# citation for a reader to follow.
REFERENCE_SUFFIXES = {".md", ".R", ".py", ".yaml", ".sh", ".toml", ".ipynb"}
REFERENCE_SKIP = {".git", "results", "mlos_review.egg-info", "__pycache__"}

# A reference, and the range it may open: §5.3, §§7.1–7.5, §6.5–6.7. Only the
# endpoints of a range are named; the interior follows from the sequence check.
REFERENCE = re.compile(r"§§?\s*(\d+(?:\.\d+)*(?:\s*[–\-]\s*\d+(?:\.\d+)*)*)")
NUMBERED_HEADING = re.compile(r"^(#{1,6}) (\d+(?:\.\d+)*)\.?\s", re.M)


def _section_numbers(name: str, source: str) -> set[str]:
    """The numbers a § may cite in one document, checked as they are read.

    Sequence is the half nothing else can check. A reference names a number and
    no renderer recomputes it, so a document that runs 5.6, 5.8 leaves every
    citation of 5.7 indistinguishable from a typo, and one that repeats 5.6
    sends two citations to two different places.

    Which heading level carries the sections is read from the document rather
    than fixed, since the math methods number `#` and the rules number `##`.
    """
    top = min((len(m.group(1)) for m in NUMBERED_HEADING.finditer(source)),
              default=None)
    if top is None:
        return set()

    numbers: set[str] = set()
    faults: list[str] = []
    current, next_section, next_subsection = None, 1, 1

    for match in re.finditer(r"^(#{1,6}) (.+)$", source, re.M):
        level, text = len(match.group(1)), match.group(2).strip()
        written = re.match(r"(\d+(?:\.\d+)*)\.?\s", text)
        number = written.group(1) if written else None
        parts = number.split(".") if number else []

        if level < top or (level == top and number is None):
            current = None          # the document title, and the appendix
        elif level == top:
            if len(parts) != 1:
                faults.append(f"section {number} carries a subsection number")
            elif int(parts[0]) != next_section:
                faults.append(f"section {number} follows section "
                              f"{next_section - 1}")
            current = int(parts[0])
            next_section, next_subsection = current + 1, 1
            numbers.add(number)
        elif level > top + 1:
            # Signposts inside a subsection carry no number. A dotted one here
            # is a third level of citation, which §N.M cannot reach and the
            # contents list does not map.
            if len(parts) > 1:
                faults.append(f"{number} is numbered below subsection level")
        elif current is not None:   # a subsection of the appendix is exempt
            if number is None:
                faults.append(f"subsection {text!r} of section {current} "
                              "carries no number")
            elif len(parts) != 2:
                faults.append(f"subsection {number} is not numbered N.M")
            else:
                if int(parts[0]) != current:
                    faults.append(f"subsection {number} sits in section {current}")
                elif int(parts[1]) != next_subsection:
                    faults.append(f"subsection {number} follows "
                                  f"{current}.{next_subsection - 1}")
                next_subsection = int(parts[1]) + 1
                numbers.add(number)

    expect_equal(f"{name}: sections are numbered in sequence", faults, [])
    return numbers


def _dangling_references(text: str, numbers: set[str]) -> list[tuple[int, str]]:
    """Every §-reference in `text` naming no section, with the line it sits on."""
    dangling = []
    for match in REFERENCE.finditer(text):
        line = text.count("\n", 0, match.start()) + 1
        for cited in re.split(r"\s*[–\-]\s*", match.group(1)):
            if cited not in numbers:
                dangling.append((line, cited))
    return dangling


def check_section_references() -> None:
    """Every §-reference in the repository names a section that exists.

    Nothing renders a §, so a reference to a section that was renumbered away,
    or was never there, reads exactly like one that works, in the markdown and
    in the Word export alike. This is the check that fires when a section is
    inserted or deleted and the citations of everything below it go stale.

    A § naming the wrong existing section is a different failure and survives
    this: README_TESTS.md cited §5.7 for the census-by-tenure metrics from the
    initial public release until August 2026, and §5.7 is a real section about
    something else, so nothing here would have objected. Only a reader catches
    that one.
    """
    section("section references")

    numbered = [name for name in DOCUMENTS
                if NUMBERED_HEADING.search((REPO_ROOT / name).read_text())]
    expect_equal("the numbered documents are the ones a § can mean",
                 sorted(numbered), sorted(NUMBERED_DOCUMENTS))

    numbers = set()
    for name in numbered:
        found = _section_numbers(name, (REPO_ROOT / name).read_text())
        if name == SECTION_SYMBOL_MEANS:
            numbers = found

    dangling = []
    for path in sorted(REPO_ROOT.rglob("*")):
        if (path.suffix not in REFERENCE_SUFFIXES or not path.is_file()
                or REFERENCE_SKIP & set(path.relative_to(REPO_ROOT).parts)):
            continue
        for line, cited in _dangling_references(path.read_text(), numbers):
            dangling.append(f"{path.relative_to(REPO_ROOT)}:{line} §{cited}")
    expect_equal(f"and every § in the repository names a section of "
                 f"{SECTION_SYMBOL_MEANS}", dangling, [])


def check_documents() -> None:
    """The guides render the same way on GitHub and in Word.

    They are the deliverable the audience actually reads, they are maintained
    by hand, and the ways they break are silent on the side you happen to be
    looking at. Nothing here reads the analysis: it is the markdown itself.
    """
    section("documents")

    # make_docx.sh decides which documents get a Word export, and every one of
    # those has to be checked here, since a rendering break is exactly what the
    # export makes visible. It is the declaration and this list is checked
    # against it, never the other way round, so adding a guide there fails here
    # until it is covered.
    #
    # Containment rather than equality: the reverse does not hold.
    # documentation_rules.md is maintainer-facing and gets no export, and is
    # still checked here, because it is markdown that people read.
    script = (REPO_ROOT / "make_docx.sh").read_text()
    declared = re.search(r"^DOCS=\((.*?)^\)", script, re.M | re.S)
    expect("make_docx.sh declares the documents it exports", declared is not None)
    if declared is not None:
        expect_equal("and this check covers every document it exports",
                     sorted(set(declared.group(1).split()) - set(DOCUMENTS)), [])

    pandoc = _find_pandoc()
    for name in DOCUMENTS:
        source = (REPO_ROOT / name).read_text()
        _blank_line_rules(name, source)
        _math_underscores(name, source)
        _structure_agrees(name, source, pandoc)


def _flag_value(argv: list[str], flag: str) -> str | None:
    """The value after `flag`, or None if it was not given.

    Reads the value rather than indexing past the end when the flag is last,
    so a mistyped invocation says what is wrong instead of raising IndexError.
    """
    if flag not in argv:
        return None
    index = argv.index(flag) + 1
    if index >= len(argv):
        raise SystemExit(f"{flag} requires a value (e.g. {flag} care_days)")
    return argv[index]


def main(argv: list[str]) -> int:
    prefix = _flag_value(argv, "--prefix")
    _state["only"] = _flag_value(argv, "--only")
    _state["quiet"] = "--quiet" in argv

    for check in (
        check_flag_extremes,
        check_selector_abstains,
        check_tied_extremes,
        check_cell_formatting,
        check_settings,
        check_schema_version,
        check_dependency_declaration,
        check_notebook_file_listing,
        check_documents,
        check_section_references,
        check_report_archiving,
        check_stacked_underscore_codes,
        check_capitalize_first,
        check_margin_takes_the_further_bound,
        check_comparison_declines_without_both_fits,
        check_bullet_pagination,
        check_gaps_past_the_plot_cap,
        check_falling_hazard_gates,
        check_order_shift_thresholds,
        check_salience_statistic,
        check_single_stratifier_default,
        check_capital_widths,
        check_layouts_render,
        check_educational_section,
        check_slide_citations,
    ):
        run_check(check)

    golden = sorted((REPO_ROOT / "tests" / "golden").glob("*/results.json"))
    if not golden:
        print("\nNo golden bundles found; run tests/run_tests.R --update-golden first.")
        return 1

    for path in golden:
        case = path.parent.name
        if prefix and not case.startswith(prefix):
            continue
        # The outer half of the isolation: run_check already contains anything
        # a check itself raises, so what reaches here is the scaffolding
        # around them -- loading the bundle, listing its stratifiers, asking
        # whether a table is available. One fixture whose bundle will not open
        # costs that fixture, not the twenty-seven after it.
        try:
            run_fixture(case, path.parent)
        except Exception as exc:  # noqa: BLE001 - any escape is this fixture failing
            errored(f"{case}: fixture aborted", exc)

    if _state["only"] is not None and not _state["matched"]:
        print(f"\n--only {_state['only']!r} matched no check.")
        return 1

    print("\n" + "=" * 71)
    summary = f"SUMMARY: {_state['pass']} passed, {_state['fail']} failed"
    if _state["error"]:
        summary += f" ({_state['error']} of them raised)"
    if _state["skip"]:
        summary += f", {_state['skip']} skipped"
    print(summary)
    print("=" * 71)
    return 1 if _state["fail"] else 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
