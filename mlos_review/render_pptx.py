"""Rendering slides to .pptx.

Deliberately thin. Everything about what a slide says lives in the blocks and
the rules; this file knows only how to put a title, some figures, a table and
some notes onto a page, and how to turn a machine name into words by asking the
vocabulary. Swapping in a markdown or xlsx renderer means writing a sibling of
this file, not touching anything upstream of it.

Geometry is computed from the slide size rather than hardcoded, because the
figure count varies: a rule that finds only one of its two figures should get
one wide figure, not a gap where the other would have been.
"""

from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from mlos_review.blocks import Format, Table, cell_format, header_text
from mlos_review.names import Vocabulary, capitalize_first

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.4)
TITLE_HEIGHT = Inches(0.8)
GUTTER = Inches(0.2)

# Preserving a figure's ratio matters more than filling the box it is given: a
# stretched survival curve misreads. The ratio is READ from each file rather
# than assumed, because it is R's decision and R has changed it once already;
# an assumption here would have gone silently stale the day it did. The
# fallback is only for a file this cannot parse, which on a manifest entry that
# resolved to an existing PNG should not happen.
FIGURE_ASPECT_FALLBACK = 3 / 2

FLAG_COLOR = RGBColor(0x99, 0x99, 0x99)

# COLOR mode paints the number itself instead of appending a letter. Blue for
# high and orange for low rather than the obvious green/red, which is the one
# pairing a red-green color blind reader cannot separate. "F" for flat has no
# color: a column where every level ties has nothing to point at, and an
# unmarked column is already the honest picture of that. The footnote still
# spells out what the marks mean, in either mode.
HIGH_COLOR = RGBColor(0x1F, 0x5C, 0xA8)
LOW_COLOR = RGBColor(0xC0, 0x5A, 0x0E)

# One size for everything inside a table: the row labels, the column headers,
# and the numbers. An earlier draft set headers smaller so a long label would
# not compete with the numbers, with an exemption for headers short enough not
# to (a bare outcome code is an identifier, not a label, and set small it read
# as an afterthought). The exemption was the tell: a table headed "From", "To",
# "Days", "Count", "Share" got three of those at one size and two at another,
# on the same row. A table reads as a grid, and a grid whose type size varies
# by cell content does not.
VALUE_PT = Pt(15)
FLAG_PT = Pt(11)
FOOTNOTE_PT = Pt(10)
BULLET_PT = Pt(18)

# The slide title, and the larger one the opening slide gets. A title slide
# that is set like every other slide does not announce anything, and the
# opening slide has the room: it carries no figure, and its one table is small.
TITLE_PT = Pt(28)
OPENING_TITLE_PT = Pt(40)

# A table's own title, drawn where the slide's title does not name it: a slide
# carrying several tables, or the opening slide, where the table sits under the
# deck's name rather than under a heading of its own. Set at the table's own
# size, so the whole block is one size of type.
TABLE_TITLE_PT = VALUE_PT
TABLE_TITLE_HEIGHT = Inches(0.35)

# What one bullet costs vertically at BULLET_PT: a line with ordinary leading,
# the gap after the paragraph, and a per-character width for working out how
# many lines it wraps to. The bullet glyph and its two spaces count as
# characters like any other. All three scale with the size the bullets are
# actually set in, so a slide set a notch smaller is measured a notch smaller
# rather than being charged for type it does not use.
BULLET_LINE_HEIGHT = Inches(0.3)
BULLET_SPACING = Pt(10)
BULLET_CHAR_WIDTH = Inches(0.13)
BULLET_PREFIX_CHARS = 3

# The size a slide's bullets are tried at, in order. A slide that would take a
# template's artwork and does not fit the band at the first size is measured
# again at the next, and takes the artwork at whichever size fits. Falling one
# notch is cheaper than the alternative, which is a slide that loses the
# branding its neighbours have over a quarter of an inch.
BULLET_SIZES = (Pt(18), Pt(16), Pt(14))

# How far a sub-bullet sits in, and the glyph it takes. The indent is written
# onto the paragraph rather than left to `level` alone: a plain text box has no
# list master behind it, so the outline level says what the line is without any
# renderer being obliged to draw it as anything.
BULLET_SUB_INDENT = Inches(0.42)
BULLET_GLYPHS = ("\u2022", "\u2013")

# The gap under a lead line. Wider than the gap between bullets, because it
# separates two kinds of text rather than two items of one kind.
LEAD_SPACING = Pt(16)

# Widest a single table column is allowed to get before the table stops
# stretching and starts centering instead.
MAX_COLUMN_WIDTH = Inches(1.8)

# How wide a string is guessed to be: a fraction of an em per character,
# applied at whatever point size that string is actually SET in. Doing it per
# size matters, because a flag mark at 11pt is narrower than the 15pt number it
# rides beside, and charging it at the value size was padding every flagged
# column by a tenth of an inch it never used.
#
# CAPITALS are charged half again as much, which is not a refinement but a
# repair. Level names are the one thing in these tables that is written in
# capitals, and at the ordinary rate "OWNER" was budgeted 0.72in against the
# 0.92in it needs, so it wrapped a letter onto a second line, and it was not
# the only level over the line. Measured against Arial, whose capitals are
# wider than the Calibri these decks ask for, so the estimate errs the safe
# way: 0.52 em for digits, lowercase and punctuation, 0.78 for capitals.
#
# Still crude. pptx offers no text metrics, and these numbers only have to
# decide how wide a table wants to be.
CHAR_WIDTH_PER_PT = Inches(0.0072)
CAPITAL_WIDTH_PER_PT = Inches(0.0108)

# The cell's own left and right insets, set below to something tighter than
# pptx's default tenth of an inch a side, plus a little slack.
CELL_MARGIN = Inches(0.06)
CELL_VERTICAL_MARGIN = Inches(0.02)
CELL_PADDING = 2 * CELL_MARGIN + Inches(0.06)

# What a plain textbox costs beyond its text: pptx insets a tenth of an inch
# either side and this file does not override it, so a box sized to its text
# alone would wrap the last word.
TEXT_BOX_INSETS = Inches(0.2)

# "No Style, No Grid": white cells, no rules, no banding of its own. The
# banding this file draws is its own, see _shade_alternate_rows.
TABLE_STYLE_ID = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"

# Every other data row gets a wash of gray instead of a rule between rows.
# Light enough to sit under a number without competing with the flag colors.
BAND_COLOR = RGBColor(0xF1, 0xF1, 0xF1)

# A shade darker, for the header of the one column a slide is about.
HIGHLIGHT_COLOR = RGBColor(0xD6, 0xD6, 0xD6)

# What the flag tokens are DRAWN as. The tokens stay "H", "L" and "F" in the
# frame, where they are data a spreadsheet renderer might express as a
# conditional format instead; only this file turns them into ink.
#
# Not letters, which is the point. The competing-risk tables head their columns
# with outcome codes, and on the OC data those are L, T and N: an "L" beside a
# number meant "lowest" while the "L" above it meant community live, on the
# same table. Any letter can collide with a code, since the codes come from the
# user's data. Arrows cannot.
FLAG_MARKS = {"H": "\u2191", "L": "\u2193", "F": "="}


# Height one table row is given before pptx grows it to fit the text.
TABLE_ROW_HEIGHT = Inches(0.4)
# One line of footnote, and the floor a footnote box is given. The height was a
# flat 0.35in, which was two things at once and wrong as both: a third of an
# inch too tall under a one-line footnote, on a slide where the tables were
# already at the foot of the page, and too short for a two-line one, which then
# ran out of its own box. It is counted now, like everything else that wraps.
FOOTNOTE_LINE_HEIGHT = Inches(0.18)
FOOTNOTE_HEIGHT = Inches(0.35)

# The typeface everything this file draws is set in, and everything it
# measures. The column widths, the row heights and the bullet pagination are
# all estimated against Calibri; a template brings its own theme fonts, and a
# wider one wraps a table header mid-word and a slide title onto a second line
# over the lead beneath it. So a template contributes its artwork and its
# colors, and the type stays what the geometry was measured in.
BODY_FONT = "Calibri"

# How close a decorated slide's content comes to the artwork above and below
# it. Narrower than MARGIN, because artwork ruled off from the body is already
# a boundary and a second one the width of a page margin reads as a gap.
DECORATION_GUTTER = Inches(0.15)

# The name a template's empty layout goes by, before the fallback of taking
# whichever layout carries the fewest placeholders.
BLANK_LAYOUT = "Blank"

# How far a template's page may be from this renderer's before it is refused.
# Not an equality test: SLIDE_WIDTH is Inches(13.333) and PowerPoint writes the
# same page as 12192000 EMU, which is 13 and a third, so the two differ by four
# ten-thousandths of an inch. A tenth of an inch is the width of a hairline at
# the edge of a projected slide and is not a different page.
SIZE_TOLERANCE = Inches(0.1)


@dataclass(frozen=True)
class Bullet:
    """One bullet line, and how far in it sits.

    A plain string is a level 0 bullet, which is what every rule with no
    sub-list passes and why `bullets` still takes strings. Depth is a field
    rather than a marker inside the text, for the reason flags travel beside a
    number rather than inside it: a depth written into prose has to be parsed
    back out by every consumer, and the one that forgets shows the marker to
    the audience.

    One level of nesting is what the renderer draws. A deeper list on a slide
    is a document.
    """

    text: str
    level: int = 0


def _as_bullet(entry: Bullet | str) -> Bullet:
    return entry if isinstance(entry, Bullet) else Bullet(str(entry))


@dataclass
class Slide:
    """One page: a title, zero or more figures, an optional table, notes.

    `notes` is a list rather than one string because notes accumulate from
    several places: what the rule wants to say about the figures, caveats a
    block attaches, and standing facts like the stay cap. Joining them here
    would mean every contributor had to know the separator and the order.
    Rendering writes one paragraph per entry.

    `findings` are finished sentences that do NOT appear on this slide. The
    deck collects them for a closing summary, so the summary says what the
    sections found rather than being written independently of them.

    `recommendations` are gathered the same way and displayed in their own
    closing section. A separate FIELD rather than findings carrying a tag
    inside their text, for the reason flags travel beside a number instead of
    being pasted into it: a kind written into prose has to be parsed back out
    by every consumer, and the one that forgets shows the tag to the audience.
    Two lists rather than one list of typed findings because kind is the only
    axis that exists today; when a third kind arrives, that is the moment to
    make it a type rather than a field.

    `table` is the one table a slide holds beside its figures. `tables` is for
    tables read side by side: a slide made OF them, which the TABLES layout
    draws, or a row of them under a STACKED slide's figures. They are separate
    fields rather than one list because SPLIT and QUADRANTS place exactly one,
    and a rule that handed them three would otherwise be silently drawing two
    of them off the edge.

    `footnote` is a line at the foot of the whole slide, as opposed to a
    table's own footnotes, which belong under that table and travel with it.
    It is for what is true of the page rather than of anything on it: where
    these numbers came from and when.

    `lead` is the same idea at the top of the body and at a size an audience
    reads: a standing qualification on everything the page goes on to say. Not
    a bullet, because it is not one of the page's items and a reader scanning a
    list should not have to work out that the first entry is a different kind
    of thing. Not the footnote either, because a qualification set at footnote
    size beneath a list is one the reader has already finished the list without
    seeing. Both are placed by `render`, which then tells the layout how much
    room is left, so neither depends on which layout the slide uses.

    `layout` names one of LAYOUT_FUNCTIONS:

        STACKED    figures in one row, table full width beneath them, or
                   `tables` side by side there instead. The default, and right
                   whenever the figures are two readings of one thing.
        SPLIT      figures left, table right, both vertically centered. For one
                   figure held against a table, and for a run of slides whose
                   figure changes while the table does not.
        QUADRANTS  figures and table in a two-column grid. Written for a slide
                   that has since moved to SPLIT, and kept for the four-figure
                   case; no rule uses it today.
        TABLES     `tables` in one row, each headed by its own title. For a
                   slide whose subject is several small tables that answer one
                   question between them.
        TITLE      bullets and tables in a column, centered in the body, with
                   a figure taking the half beside them if there is one. The
                   opening slide.

    `schematic` says the figures on this slide are diagrams rather than
    readings: they carry no value anyone reads off them, so room taken from
    them costs nothing, and a slide that would otherwise go unbranded can take
    a template's artwork by giving its diagram the smaller half. A rule sets it
    about its own figures, since what a figure is showing is the rule's to say.

    `layout` and `schematic` are what a rule declares about presentation,
    because how many figures a slide carries, whether they are peers and
    whether they hold values are things only the rule knows; the geometry that
    follows from them is decided here.
    """

    title: str
    figures: list[Path] = field(default_factory=list)
    table: Table | None = None
    tables: list[Table] = field(default_factory=list)
    bullets: list[Bullet | str] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)
    findings: list[str] = field(default_factory=list)
    recommendations: list[str] = field(default_factory=list)
    footnote: str = ""
    lead: str = ""
    layout: str = "STACKED"
    schematic: bool = False


# --- Templates -----------------------------------------------------------
#
# A template is an ordinary one-slide .pptx whose slide carries decoration and
# nothing else: a header band, a footer strip of logos. Its artwork is copied
# onto the slides that have room for it rather than being put on the master,
# because the figures are opaque white PNGs and a master would put the artwork
# behind a rectangle. Which slides have room is `takes_decoration`.
#
# What a template brings beyond its artwork is its theme, so the deck is set in
# the template's fonts and reads its color scheme, and that applies to every
# slide whether decorated or not.


@dataclass(frozen=True)
class Decoration:
    """A template's artwork, and the band of the slide it leaves free.

    `top` and `bottom` bound the tallest horizontal band no template shape sits
    in. A decorated slide draws everything between them: its title, its body
    and its footnote, so the shrinking is charged once and every layout
    inherits it without knowing a template is in play.
    """

    elements: tuple = ()
    images: dict = field(default_factory=dict)
    top: int = int(MARGIN)
    bottom: int = int(SLIDE_HEIGHT - MARGIN)


def takes_decoration(spec: "Slide") -> bool:
    """Whether a template's artwork is worth putting on this slide.

    A slide carrying a figure is not. The figures are opaque white PNGs sized
    to whatever the body leaves them, so artwork behind one is artwork nobody
    sees, and taking an inch off the top and half an inch off the foot of a
    survival curve costs more than a band is worth.

    A schematic is the exception its `Slide` field describes, and only where it
    sits BESIDE the text rather than under it. On TITLE it does, so what it
    gives up to the band is width it can spare; under a list it would be giving
    up the height it is read in.

    Whether the artwork FITS is a second question, asked in `render` once the
    slide's content has been measured.
    """
    if not spec.figures:
        return True
    return spec.schematic and spec.layout == "TITLE"


def free_band(shapes) -> tuple[int, int]:
    """The tallest horizontal band of a slide that no shape sits in.

    Measured rather than assumed, so a template with its artwork all at the
    foot, or split above and below, is read as what it is. Overlapping shapes
    are merged as the scan goes, since a logo strip is usually a group with a
    rule across it.
    """
    spans = sorted((int(shape.top), int(shape.top + shape.height))
                   for shape in shapes
                   if shape.top is not None and shape.height is not None)
    best, floor = (0, 0), 0
    for top, bottom in spans + [(int(SLIDE_HEIGHT), int(SLIDE_HEIGHT))]:
        if top - floor > best[1] - best[0]:
            best = (floor, top)
        floor = max(floor, bottom)
    return best


def template_band(path: str | Path | None) -> Decoration:
    """The band a template leaves free, without its shapes.

    For a caller that has to know how much room a page has before the renderer
    is reached: the closing sections paginate their bullets, and a page broken
    against the full slide would run its last lines under the footer strip.
    A template of None answers with the undecorated slide, so the caller can
    ask unconditionally.
    """
    if path is None:
        return Decoration()
    deck = Presentation(str(path))
    _require_slide_size(deck, path)
    if not deck.slides:
        return Decoration()
    top, bottom = free_band(deck.slides[0].shapes)
    return Decoration(top=top + int(DECORATION_GUTTER),
                      bottom=bottom - int(DECORATION_GUTTER))


def _require_slide_size(deck, path) -> None:
    """Refuse a template of another size rather than stretching its artwork.

    Every measurement in this file is taken from SLIDE_WIDTH and SLIDE_HEIGHT,
    so a template of another shape would have to move all of them. Resizing
    the template instead is the caller's one line, and it is the line that
    keeps the artwork the proportions it was drawn at.
    """
    if (abs(deck.slide_width - SLIDE_WIDTH) <= SIZE_TOLERANCE
            and abs(deck.slide_height - SLIDE_HEIGHT) <= SIZE_TOLERANCE):
        return
    raise ValueError(
        f"Template '{path}' is {Emu(deck.slide_width).inches:g} by "
        f"{Emu(deck.slide_height).inches:g} inches; this renderer draws "
        f"{Emu(SLIDE_WIDTH).inches:g} by {Emu(SLIDE_HEIGHT).inches:g}. "
        "Resize the template to match.")


def _read_decoration(deck, path) -> Decoration:
    """Lift the artwork off the template's own slide, leaving the slide there.

    The slide goes at the end of the render, not here. It is what keeps the
    template's images reachable while the deck is built, and pptx names a new
    image part after the ones the package already holds: with the template
    slide gone, the first figure added is offered `image1.png`, which the
    template is still holding, and both are written into the file under that
    one name.
    """
    _require_slide_size(deck, path)
    if not deck.slides:
        return Decoration()
    source = deck.slides[0]
    images = {rid: rel.target_part for rid, rel in source.part.rels.items()
              if rel.reltype == RT.IMAGE and not rel.is_external}
    top, bottom = free_band(source.shapes)
    return Decoration(
        elements=tuple(deepcopy(shape._element) for shape in source.shapes),
        images=images,
        top=top + int(DECORATION_GUTTER),
        bottom=bottom - int(DECORATION_GUTTER))


def _drop_template_slide(deck) -> None:
    """Take the template's own slide back out, once the deck is built.

    Dropped rather than kept and skipped over: what was wanted from it is its
    shapes, and a template that printed its own page would put a blank one at
    the front of the deck. Its images stay, every slide stamped with them
    having related to them by now.
    """
    entries = deck.slides._sldIdLst
    if len(entries):
        deck.part.drop_rel(entries[0].rId)
        entries.remove(entries[0])


# Where a copied shape names the image it draws. The second is a linked rather
# than embedded picture, and the SVG beside a fallback PNG is reached the same
# way, from an extension element the scan below walks into like any other.
IMAGE_REFERENCES = (qn("r:embed"), qn("r:link"))


def _stamp(slide, decoration: Decoration) -> None:
    """Put the artwork on this slide, behind whatever is drawn next.

    Called before the content, which settles two things at once: shapes early
    in the tree are at the back of the z-order, and the ids pptx hands the
    content are numbered above the ones copied here.
    """
    tree = slide.shapes._spTree
    for element in decoration.elements:
        copy = deepcopy(element)
        fresh: dict[str, str] = {}
        for node in copy.iter():
            for name in IMAGE_REFERENCES:
                rid = node.get(name)
                if rid not in decoration.images:
                    continue
                if rid not in fresh:
                    fresh[rid] = slide.part.relate_to(
                        decoration.images[rid], RT.IMAGE)
                node.set(name, fresh[rid])
        tree.append(copy)


def _set_font(slide, skip: int = 0) -> None:
    """Set the deck's typeface on everything drawn on this slide.

    Applied to the whole slide at the end rather than at each of the dozen
    places a run is written, so a drawing routine added later inherits it
    without having to remember. `skip` is how many shapes were already there,
    which is the template's artwork: that is set in the template's own type and
    is not this file's to restyle.
    """
    for shape in list(slide.shapes)[skip:]:
        frames = []
        if shape.has_text_frame:
            frames.append(shape.text_frame)
        if shape.has_table:
            frames.extend(cell.text_frame for row in shape.table.rows
                          for cell in row.cells)
        for frame in frames:
            for paragraph in frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = BODY_FONT


def _blank_layout(deck):
    """The layout to build every slide on: the one with nothing already on it.

    Found by name, then by placeholder count, rather than by the index that is
    Blank in the default template. A template is someone else's file, and a
    layout carrying a title box would print that box, empty, on every slide.
    """
    for layout in deck.slide_layouts:
        if layout.name == BLANK_LAYOUT:
            return layout
    return min(deck.slide_layouts, key=lambda layout: len(layout.placeholders))



def _format_cell(value, fmt: Format) -> str:
    """Write one cell out. Percent scaling happens HERE, not in the frame.

    Text passes through untouched. A column may hold a date or the name of a
    level, and there is nothing for a decimal count to do to either; see the
    note on text columns in `Table`.
    """
    if isinstance(value, str):
        return value
    if value is None or value != value:  # NaN
        return "" if fmt.blank else "n/a"
    if fmt.percent:
        return f"{value * 100:,.{fmt.decimals}f}%"
    return f"{value:,.{fmt.decimals}f}"


def _text_width(text: str, size: Pt) -> int:
    """How wide a string will be, set at one point size."""
    capitals = sum(1 for character in text if character.isupper())
    return int(((len(text) - capitals) * CHAR_WIDTH_PER_PT
                + capitals * CAPITAL_WIDTH_PER_PT) * size.pt)


def _width_for(content: int) -> int:
    return min(MAX_COLUMN_WIDTH, int(content + CELL_PADDING))


def _word_width(text: str, size: Pt) -> int:
    """The widest unbreakable run in a string: its longest word.

    What a column may never be squeezed below. PowerPoint wraps at spaces, so a
    column narrower than its longest WORD has nowhere legal to break and breaks
    mid-word instead, one character to a line. "Pct" in a column of a third of
    an inch is the case that found this.
    """
    return max([_text_width(word, size) for word in text.split()] or [0])


def _wrapped_lines(text: str, size: Pt, width: int) -> int:
    """How many lines a string takes in a cell of this width.

    Greedy, breaking at spaces, which is what PowerPoint does. Crude like
    everything else here, and it only has to decide how tall a row wants to be:
    a header that wraps is fine, a header that wraps into a row sized for one
    line is what puts a footnote over the last row of a table.
    """
    room = width - CELL_PADDING
    if room <= 0:
        return len(text.split()) or 1
    lines, current = 1, 0
    for word in text.split():
        word_width = _text_width(word, size)
        step = word_width + (_text_width(" ", size) if current else 0)
        if current and current + step > room:
            lines, current = lines + 1, word_width
        else:
            current += step
        # A word too wide for the cell breaks INSIDE itself, as many times as
        # it has to. Counting that is the difference between a height estimate
        # that is merely rough and one that is wrong in the one case that
        # matters: "Pct" in a third of an inch is three lines, not one, and a
        # footnote placed under one of them lands on the table.
        if current > room:
            extra = -(-current // room) - 1
            lines, current = lines + extra, current - extra * room
    return lines


def _column_widths(table: Table, vocab: Vocabulary,
                   flag_style: str = "MARK") -> list[int]:
    """A width per rendered column, each sized to its OWN longest cell.

    One width for every column would be set by the widest cell anywhere in the
    table, which on a table of level names and short numbers means a column of
    five-character labels paid for by an eleven-character header two columns
    over. Sizing per column is what lets a table be genuinely narrow, and a
    narrow table on a split slide is width handed to the figure beside it.

    The index columns come first, in index order, then one per measure.
    """
    widths = []
    for depth in range(table.df.index.nlevels):
        labels = table.df.index.get_level_values(depth)
        widths.append(_width_for(
            max([_text_width(str(v), VALUE_PT) for v in labels] or [0])))

    # A flag rides in the same cell as its number but is set smaller, so it is
    # charged at its own size rather than at three value characters. COLOR
    # mode paints the number instead of appending a mark, so it is charged
    # nothing: the condition mirrors the branch in _add_table that decides
    # whether a mark run is added at all.
    flag = (_text_width("  \u2191", FLAG_PT)
            if table.flags is not None and flag_style != "COLOR" else 0)
    for measure in table.df.columns:
        header = header_text(table, vocab, measure)
        widest = _text_width(header, VALUE_PT)
        for row in table.df.index:
            text = _format_cell(table.df.loc[row, measure],
                                cell_format(table, row, measure))
            widest = max(widest, _text_width(text, VALUE_PT) + flag)
        widths.append(_width_for(widest))
    return widths


def _column_minimums(table: Table, vocab: Vocabulary,
                     flag_style: str = "MARK") -> list[int]:
    """The narrowest each column may be drawn: its longest word, plus padding.

    A header may wrap, and on a slide of three tables it usually must. What it
    may not do is wrap inside a word, which is what a proportional squeeze with
    no floor produces and what it produced here: three tables asking for 21
    inches of a 12.5-inch slide, every column scaled to 59%, and a column
    headed "Pct" drawn at a third of an inch.

    Index columns are floored at their longest LABEL rather than their longest
    word, since a level name is a name and breaking "_UNKNOWN_" across two
    lines is not a wrap, it is a typo.
    """
    minimums = []
    for depth in range(table.df.index.nlevels):
        labels = table.df.index.get_level_values(depth)
        minimums.append(_width_for(
            max([_text_width(str(v), VALUE_PT) for v in labels] or [0])))

    flag = (_text_width("  \u2191", FLAG_PT)
            if table.flags is not None and flag_style != "COLOR" else 0)
    for measure in table.df.columns:
        widest = _word_width(header_text(table, vocab, measure), VALUE_PT)
        for row in table.df.index:
            text = _format_cell(table.df.loc[row, measure],
                                cell_format(table, row, measure))
            # A formatted number never contains a space, so its own width IS
            # its minimum: "10,228" cannot be broken anywhere.
            widest = max(widest, _text_width(text, VALUE_PT) + flag)
        minimums.append(_width_for(widest))
    return minimums


def _squeezed(widths: list[int], minimums: list[int], box_width: int) -> list[int]:
    """Column widths brought within `box_width` without crushing any of them.

    The deficit comes out of the columns that have room to give it, in
    proportion to how much each HAS to give, so a wide column of long headers
    yields and a three-character column of percentages does not.

    Where even the minimums do not fit, the minimums are returned and the table
    comes out WIDER than the box it was given. That is deliberate. The
    alternative, scaling them down together, is what produced a percentage
    column a third of an inch wide with its heading broken into one letter per
    line, and it did it silently: the table was inside its box, so nothing
    downstream could tell. A table hanging off the edge of the slide is caught
    by the geometry check on every fixture, and by anyone looking at the slide.
    The fix for it is never here; it is fewer columns.
    """
    if sum(widths) <= box_width:
        return widths
    if sum(minimums) >= box_width:
        return minimums

    slack = [max(0, w - m) for w, m in zip(widths, minimums)]
    deficit = sum(widths) - box_width
    total_slack = sum(slack)
    return [int(w - deficit * (s / total_slack)) if total_slack else w
            for w, s in zip(widths, slack)]


def _fitted_widths(table: Table, vocab: Vocabulary, flag_style: str,
                   box_width: int) -> list[int]:
    """The widths a table is actually drawn at, in the box it was given."""
    return _squeezed(_column_widths(table, vocab, flag_style),
                     _column_minimums(table, vocab, flag_style), box_width)


def _plain_table_style(grid) -> None:
    """Strip the blue banded style pptx gives a new table, and tighten the cells.

    The default is Medium Style 2 Accent 1: a dark blue header and light blue
    bands. It costs the deck twice. A high value painted blue by the COLOR flag
    style barely separates from a blue-tinted row, which is the mark failing at
    the one job it has; and a banded table pulls attention away from figures
    drawn on white.

    The cell insets go with it. pptx leaves a tenth of an inch either side of
    every cell, which across seven columns is most of an inch of table that
    could have been figure.
    """
    # By tag, not by position: the style id is the element named tableStyleId,
    # wherever it sits among tblPr's children, and a table born without one
    # gets one rather than a crash.
    properties = grid._tbl.tblPr
    style = properties.find(qn("a:tableStyleId"))
    if style is None:
        style = properties.makeelement(qn("a:tableStyleId"), {})
        properties.append(style)
    style.text = TABLE_STYLE_ID
    grid.first_row = False
    grid.horz_banding = False
    for row in grid.rows:
        for cell in row.cells:
            cell.margin_left = CELL_MARGIN
            cell.margin_right = CELL_MARGIN
            cell.margin_top = CELL_VERTICAL_MARGIN
            cell.margin_bottom = CELL_VERTICAL_MARGIN


def _shade_alternate_rows(grid) -> None:
    """Wash every other data row, instead of ruling a line between all of them.

    Bands separate rows with less ink than rules do, and on the stacked
    competing-risk table they separate something else for free: its rows come in
    share-then-days pairs, so banding every other one shades every share row and
    leaves every days row white. What was a footnote telling the reader which
    half was which becomes the shape of the table.
    """
    # Row 0 is the header and is even, so the one test covers it too.
    for index, row in enumerate(grid.rows):
        if index % 2 == 0:
            continue
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND_COLOR


def _add_title(slide, text: str, size: Pt = TITLE_PT,
               top: int = int(MARGIN)) -> None:
    box = slide.shapes.add_textbox(MARGIN, top, SLIDE_WIDTH - 2 * MARGIN,
                                   TITLE_HEIGHT)
    frame = box.text_frame
    frame.text = capitalize_first(text)
    frame.paragraphs[0].runs[0].font.size = size
    frame.paragraphs[0].runs[0].font.bold = True


def figure_aspect(path: Path) -> float:
    """Width over height of a PNG, from its header. No decoding, no library.

    Sixteen bytes settle it, so there is no reason to pull an image library
    into a package that otherwise needs pandas and python-pptx. A file that is
    not a PNG, or is truncated, falls back rather than raising: a figure that
    lands in the wrong-shaped box is a blemish, an exception is a lost deck.
    """
    try:
        with open(path, "rb") as handle:
            head = handle.read(24)
    except OSError:
        return FIGURE_ASPECT_FALLBACK
    if len(head) < 24 or head[:8] != b"\x89PNG\r\n\x1a\n" or head[12:16] != b"IHDR":
        return FIGURE_ASPECT_FALLBACK
    width = int.from_bytes(head[16:20], "big")
    height = int.from_bytes(head[20:24], "big")
    return width / height if width and height else FIGURE_ASPECT_FALLBACK


def _add_figure_row(slide, figures: list[Path], left: Emu, top: Emu,
                    width: Emu, height: Emu, valign: str = "center") -> None:
    """Figures in one row inside a box, centered as a group horizontally.

    Each figure is as wide as its share of the box and its own aspect ratio
    allow. `valign` decides the vertical position of the leftovers: "center"
    for a figure held beside a table it should sit level with, "top" for the
    stacked layout, where the table follows beneath and the slack belongs
    between them rather than above the figures.
    """
    if not figures:
        return
    cell_width = int((width - GUTTER * (len(figures) - 1)) / len(figures))
    drawn = [(path, min(cell_width, int(height * figure_aspect(path))))
             for path in figures]
    # Center the row when the aspect ratio leaves the cells underfilled.
    used = sum(w for _, w in drawn) + GUTTER * (len(figures) - 1)
    x = int(left + (width - used) / 2)
    for path, drawn_width in drawn:
        drawn_height = int(drawn_width / figure_aspect(path))
        y = top if valign == "top" else int(top + (height - drawn_height) / 2)
        slide.shapes.add_picture(str(path), x, y, width=drawn_width)
        x += drawn_width + GUTTER


def _add_figure_in(slide, path: Path, left: Emu, top: Emu,
                   width: Emu, height: Emu) -> None:
    """One figure, centered inside a box, at its own aspect ratio."""
    aspect = figure_aspect(path)
    drawn_width = min(int(width), int(height * aspect))
    drawn_height = int(drawn_width / aspect)
    slide.shapes.add_picture(
        str(path),
        int(left + (width - drawn_width) / 2),
        int(top + (height - drawn_height) / 2),
        width=drawn_width)


def _add_table(slide, table: Table, vocab: Vocabulary, top: Emu, height: Emu,
               flag_style: str = "MARK",
               left: Emu | None = None, available: Emu | None = None,
               widths: list[int] | None = None) -> None:
    """One header row of short labels, then one row per level.

    Flags ride alongside their number in the same cell but as a separate,
    grayed run, so the table reads as numbers with marks rather than as
    strings. A renderer that can do better, like xlsx with conditional
    formatting, has the flag frame available and need not follow this.

    `left` and `available` bound the space the table may use, defaulting to the
    full text width. A quadrant layout passes its cell instead, so the same
    sizing rule applies within a narrower box rather than the table having to
    know which layout it is in.

    `widths` lets a caller that has already worked out the columns hand them
    over. A caller that needs the table's HEIGHT needs its widths first, since
    a header wraps or does not depending on them, and computing them twice is
    how the height and the drawing come to disagree.
    """
    index = table.df.index
    measures = list(table.df.columns)
    depth = index.nlevels
    rows, cols = len(index) + 1, len(measures) + depth

    # A narrow table is centered in the space it was given rather than stretched
    # across it. The nine-measure LOS table fills the slide and should; a table
    # of short columns looks like a mistake when it does.
    box_left = MARGIN if left is None else left
    box_width = (SLIDE_WIDTH - 2 * MARGIN) if available is None else available
    if widths is None:
        widths = _fitted_widths(table, vocab, flag_style, box_width)
    width = sum(widths)
    table_left = int(box_left + (box_width - width) / 2)

    shape = slide.shapes.add_table(rows, cols, table_left, top, int(width), height)
    grid = shape.table
    _plain_table_style(grid)
    _shade_alternate_rows(grid)
    for column, column_width in zip(grid.columns, widths):
        column.width = column_width

    # Units are not repeated per column; they go in the table's footnote. Most
    # of the measures on an LOS table are in days, so a units line under every
    # header was spending a row of the reader's attention to say "days" over
    # and over. The index columns are unheaded for the same reason: a column of
    # level names headed "level" tells the reader nothing they cannot see.
    for j in range(depth):
        grid.cell(0, j).text = ""
    for j, measure in enumerate(measures, start=depth):
        cell = grid.cell(0, j)
        cell.text = header_text(table, vocab, measure)
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.size = VALUE_PT
            run.font.bold = True
        if measure == table.highlight:
            cell.fill.solid()
            cell.fill.fore_color.rgb = HIGHLIGHT_COLOR

    previous: tuple = ()
    for i, row in enumerate(index, start=1):
        parts = row if isinstance(row, tuple) else (row,)
        for j, part in enumerate(parts):
            # An outer label that repeats is written once. Six levels each
            # spanning a share row and a days row should read as six groups,
            # not as twelve rows that happen to pair up.
            if previous[:j + 1] == parts[:j + 1]:
                continue
            cell = grid.cell(i, j)
            cell.text = str(part)
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = VALUE_PT
        previous = parts

        for j, measure in enumerate(measures, start=depth):
            paragraph = grid.cell(i, j).text_frame.paragraphs[0]
            value = paragraph.add_run()
            value.text = _format_cell(table.df.loc[row, measure],
                                      cell_format(table, row, measure))
            value.font.size = VALUE_PT
            flag = "" if table.flags is None else table.flags.loc[row, measure]
            if not flag:
                continue
            if flag_style == "COLOR":
                if flag in ("H", "L"):
                    value.font.color.rgb = HIGH_COLOR if flag == "H" else LOW_COLOR
                    value.font.bold = True
            else:
                mark = paragraph.add_run()
                mark.text = f"  {FLAG_MARKS.get(flag, flag)}"
                mark.font.color.rgb = FLAG_COLOR
                mark.font.size = FLAG_PT


def bullet_height(line: Bullet | str, width: Emu | None = None,
                  size: Pt = BULLET_PT) -> int:
    """Height one bullet needs once it wraps, including the space beneath it.

    Estimated from the character count, like the table column widths and for
    the same reason: pptx cannot measure text, and a caller only has to decide
    how many bullets fit on a page. Deliberately generous, since the cost of
    over-estimating is an early page break and the cost of under-estimating is
    text running off the bottom of a slide.

    `width` is the column the text is set in, defaulting to the full text
    width. A layout that puts bullets beside something passes its own column,
    since the same sentence wraps to twice the lines in half the width.

    `size` is the type the bullets are set in. Both the line and the character
    it holds scale with it, so a smaller size buys height twice over: shorter
    lines, and fewer of them.
    """
    bullet = _as_bullet(line)
    scale = size.pt / BULLET_PT.pt
    full = SLIDE_WIDTH - 2 * MARGIN
    # A sub-bullet is set in a narrower column, so the same sentence wraps to
    # more lines there than it would at the margin.
    column = (full if width is None else width) - bullet.level * BULLET_SUB_INDENT
    per_line = max(1, int(column / (BULLET_CHAR_WIDTH * scale)))
    lines = max(1, -(-(len(bullet.text) + BULLET_PREFIX_CHARS) // per_line))
    return int(lines * BULLET_LINE_HEIGHT * scale + int(BULLET_SPACING) * scale)


def lead_height(text: str) -> int:
    """Height a lead line needs, wrapping included.

    Measured like a bullet minus the glyph, since it is set at the same size in
    the same column: what it does not have is the bullet and its two spaces.
    A blank lead costs nothing, so a slide without one loses no room to it.
    """
    if not text:
        return 0
    per_line = max(1, int((SLIDE_WIDTH - 2 * MARGIN) / BULLET_CHAR_WIDTH))
    lines = max(1, -(-len(text) // per_line))
    return lines * BULLET_LINE_HEIGHT + int(LEAD_SPACING)


def text_budget(decoration: Decoration | None = None) -> int:
    """Height a slide of nothing but bullets has for its list.

    Takes the band a template leaves free when there is one, since a bullets
    slide carries no figure and is therefore decorated.
    """
    decoration = decoration or Decoration()
    return decoration.bottom - decoration.top - int(TITLE_HEIGHT)


def bullet_pages(bullets: list[Bullet | str], reserved: int = 0,
                 budget: int | None = None) -> list[list[Bullet | str]]:
    """Split bullets into pages that fit a slide, keeping their order.

    Geometry lives here rather than in the rule that gathers the bullets: how
    much fits on a page is this file's business, and the rule's business is
    only that everything gathered has to appear somewhere. A single bullet
    taller than a page still gets its own page rather than being dropped or
    split mid-sentence.

    `reserved` is height the FIRST page has already spent, which is what a lead
    line costs. Only the first, because that is where a lead goes;
    charging every page for it would break the section a bullet early from the
    second page onward.

    `budget` is how tall a page is, defaulting to the undecorated slide. A
    caller building against a template passes the smaller band, because a page
    broken against the full slide runs its last lines under the artwork.
    """
    budget = text_budget() if budget is None else budget
    pages: list[list[Bullet | str]] = []
    page: list[Bullet | str] = []
    used = reserved
    for line in bullets:
        height = bullet_height(line)
        if page and used + height > budget:
            pages.append(page)
            page, used = [], 0
        page.append(line)
        used += height
    if page:
        pages.append(page)
    return pages


def _add_bullets(slide, bullets: list[Bullet | str], top: Emu, height: Emu,
                 width: Emu | None = None, size: Pt = BULLET_PT) -> None:
    box = slide.shapes.add_textbox(
        MARGIN, top, (SLIDE_WIDTH - 2 * MARGIN) if width is None else width,
        height)
    frame = box.text_frame
    frame.word_wrap = True
    for index, line in enumerate(bullets):
        bullet = _as_bullet(line)
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        glyph = BULLET_GLYPHS[min(bullet.level, len(BULLET_GLYPHS) - 1)]
        run.text = f"{glyph}  {bullet.text}"
        run.font.size = size
        paragraph.space_after = Pt(BULLET_SPACING.pt * size.pt / BULLET_PT.pt)
        if bullet.level:
            paragraph.level = min(bullet.level, 8)
            # The indent itself, since a bare text box carries no list master
            # to turn the level into one. marL moves the whole line in and
            # indent 0 keeps the glyph with its text rather than hanging it.
            properties = paragraph._p.get_or_add_pPr()
            properties.set("marL", str(int(bullet.level * BULLET_SUB_INDENT)))
            properties.set("indent", "0")


def _add_lead(slide, text: str, top: Emu, height: Emu) -> None:
    """The standing qualification above a page's body.

    Bullet size, so it is read rather than skipped, and italic, so a reader can
    see at a glance that it is not one of the lines below it. No bullet glyph,
    for the same reason.
    """
    box = slide.shapes.add_textbox(
        MARGIN, top, SLIDE_WIDTH - 2 * MARGIN, height)
    frame = box.text_frame
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = BULLET_PT
    run.font.italic = True


FLAG_LEGEND = {
    "MARK": "\u2191 highest, \u2193 lowest, = flat (all levels tied).",
    "COLOR": "Blue highest, orange lowest; unmarked where all levels tie.",
}


def _footnote_text(table: Table, flag_style: str) -> str:
    """One line of footnote from the table's own notes and its flag legend."""
    notes = list(table.footnotes)
    if table.flags is not None:
        notes.insert(0, FLAG_LEGEND[flag_style])
    return "  ".join(notes)


def _footnote_height(table: Table | None, flag_style: str = "MARK",
                     width: Emu | None = None) -> int:
    """How deep a table's footnote runs, wrapping counted.

    A table at the foot of a slide is where this matters: half an inch of
    unused box under a one-line footnote is half an inch the tables above it
    could have had, and it was enough to push the last of them off the page.
    """
    if table is None:
        return 0
    text = _footnote_text(table, flag_style)
    if not text:
        return 0
    room = (SLIDE_WIDTH - 2 * MARGIN) if width is None else width
    lines = _wrapped_lines(text, FOOTNOTE_PT, int(room) - int(TEXT_BOX_INSETS)
                           + int(CELL_PADDING))
    return lines * int(FOOTNOTE_LINE_HEIGHT)


def _add_footnotes(slide, table: Table, top: Emu, flag_style: str,
                   left: Emu | None = None, width: Emu | None = None) -> None:
    notes = list(table.footnotes)
    if table.flags is not None:
        notes.insert(0, FLAG_LEGEND[flag_style])
    if not notes:
        return
    box_width = (SLIDE_WIDTH - 2 * MARGIN) if width is None else width
    box = slide.shapes.add_textbox(
        MARGIN if left is None else left, top, box_width,
        _footnote_height(table, flag_style, box_width))
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = "  ".join(notes)
    frame.paragraphs[0].runs[0].font.size = FOOTNOTE_PT
    frame.paragraphs[0].runs[0].font.italic = True


def _add_slide_footnote(slide, text: str,
                        bottom: int = int(SLIDE_HEIGHT - MARGIN)) -> None:
    """A line at the foot of the page, in the same voice as a table's footnote.

    Same size and italics deliberately: it is the same kind of remark, a
    qualification the reader can take or leave, and giving provenance its own
    typography would make it look like a finding.
    """
    box = slide.shapes.add_textbox(
        MARGIN, int(bottom - FOOTNOTE_HEIGHT),
        SLIDE_WIDTH - 2 * MARGIN, FOOTNOTE_HEIGHT)
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = text
    frame.paragraphs[0].runs[0].font.size = FOOTNOTE_PT
    frame.paragraphs[0].runs[0].font.italic = True


def _header_lines(table: Table, vocab: Vocabulary, widths: list[int]) -> int:
    """How many lines the deepest header in this table wraps to.

    One when the columns got what they asked for, which is the ordinary case;
    two or three on a slide of three tables, where the row is squeezed and the
    headers wrap at their spaces. It has to be counted rather than assumed,
    because pptx grows a row to fit its text whatever height the table was
    given, so a header wrapping into a row sized for one line does not overflow
    the row: it pushes everything below it down, on top of the footnote.
    """
    header_widths = widths[table.df.index.nlevels:]
    return max([_wrapped_lines(header_text(table, vocab, measure), VALUE_PT,
                               width)
                for measure, width in zip(table.df.columns, header_widths)]
               or [1])


def _table_heights(table: Table | None, vocab: Vocabulary | None = None,
                   widths: list[int] | None = None,
                   flag_style: str = "MARK") -> tuple[int, int]:
    """The height a table wants, and the height its footnote needs.

    `widths` and `vocab` are optional only because a caller that has not sized
    the table yet still needs a figure to lay out around; without them the
    header is assumed to fit on one line, which is what it does whenever the
    table is not being squeezed.
    """
    if table is None:
        return 0, 0
    lines = (_header_lines(table, vocab, widths)
             if vocab is not None and widths is not None else 1)
    return (int(TABLE_ROW_HEIGHT * (len(table.df) + lines)),
            _footnote_height(table, flag_style,
                             sum(widths) if widths else None))


def _layout_stacked(slide, spec: Slide, vocab: Vocabulary, top: Emu, height: Emu,
                    flag_style: str,
                    bullet_pt: Pt = BULLET_PT) -> None:
    """Bullets, then a row of figures, then the table across the full width.

    `tables` puts several of them side by side beneath the figures instead,
    each headed by its own title, for a slide whose figures need more than one
    table read with them. The row is measured and drawn by the code the TABLES
    layout uses, so a pair of tables under a pair of figures sits at the widths
    it would have on a slide of its own.
    """
    row = list(spec.tables)
    available = int(SLIDE_WIDTH - 2 * MARGIN)
    if row:
        # A titled table's depth already counts its footnotes, so the row asks
        # for one height and the single-table path keeps its two.
        table_height, footnote_height = _row_depth(row, vocab, flag_style,
                                                   available), Emu(0)
        widths = None
    else:
        widths = (_fitted_widths(spec.table, vocab, flag_style, available)
                  if spec.table is not None else None)
        table_height, footnote_height = _table_heights(spec.table, vocab, widths)

    bullets_height = Emu(0)
    if spec.bullets:
        bullets_height = sum(bullet_height(line, size=bullet_pt)
                             for line in spec.bullets)
        _add_bullets(slide, spec.bullets, top, int(bullets_height),
                     size=bullet_pt)

    figure_height = height - table_height - footnote_height - bullets_height
    if spec.figures:
        figure_height -= GUTTER
        _add_figure_row(slide, spec.figures, MARGIN, int(top + bullets_height),
                        SLIDE_WIDTH - 2 * MARGIN, int(figure_height),
                        valign="top")

    table_top = int(top + bullets_height + figure_height
                    + (GUTTER if spec.figures else 0))
    if row:
        _draw_table_row(slide, row, vocab, MARGIN, table_top, available,
                        int(table_height), flag_style)
    elif spec.table is not None:
        _add_table(slide, spec.table, vocab, table_top, table_height,
                   flag_style=flag_style, widths=widths)
        _add_footnotes(slide, spec.table, int(table_top + table_height), flag_style)


def _layout_split(slide, spec: Slide, vocab: Vocabulary, top: Emu, height: Emu,
                  flag_style: str,
                  bullet_pt: Pt = BULLET_PT) -> None:
    """Figures down the left, table down the right, both vertically centered.

    For a slide that shows ONE figure against a table the reader is meant to
    hold in view while the figure changes. Stacking them would cap the figure
    at whatever height the table left over, which on a six-row table is under
    three inches; side by side, the figure runs to the full body height unless
    the table's width stops it first. On the OC competing-risk slides that is
    7.8 by 5.2 inches against the 4.3 by 2.9 a quadrant cell gave.

    The table takes its natural width but never more than half, so a narrow
    table hands the surplus to the figure rather than sitting in a column of
    its own making.
    """
    body_width = SLIDE_WIDTH - 2 * MARGIN
    table_width = 0
    if spec.table is not None:
        table_width = min(int(body_width / 2),
                          sum(_column_widths(spec.table, vocab, flag_style)))

    figure_width = int(body_width - table_width - (GUTTER if table_width else 0))
    if spec.figures:
        _add_figure_row(slide, spec.figures, MARGIN, top, figure_width, height)

    if spec.table is None:
        return
    widths = _fitted_widths(spec.table, vocab, flag_style, table_width)
    table_height, footnote_height = _table_heights(spec.table, vocab, widths)
    left = int(SLIDE_WIDTH - MARGIN - table_width)
    table_top = int(top + (height - table_height - footnote_height) / 2)
    _add_table(slide, spec.table, vocab, table_top, table_height,
               flag_style=flag_style, left=left, available=table_width,
               widths=widths)
    _add_footnotes(slide, spec.table, int(table_top + table_height), flag_style,
                   left=left, width=table_width)


def _layout_quadrants(slide, spec: Slide, vocab: Vocabulary, top: Emu, height: Emu,
                      flag_style: str,
                      bullet_pt: Pt = BULLET_PT) -> None:
    """Figures and table in a two-column grid, filled in reading order.

    NOT USED BY ANY RULE TODAY, and kept deliberately. The slide it was written
    for, three competing-risk panels over a table, moved to SPLIT: a 4:3-ish
    figure in a half-width, half-height cell cannot fill it on a 16:9 slide, so
    three panels and a table of a different width read as four things scattered
    rather than one exhibit. That is an argument about THAT slide, not about
    the grid. Four figures of equal standing and no table is the shape this
    serves, and it is the shape a future rule is most likely to want; the
    synthetic renderer test keeps it from rotting in the meantime.

    The table takes the cell after the last figure, which on the slide this was
    built for is the fourth of four. Nothing here assumes three figures: a
    dataset with a different number of competing outcomes gets as many rows as
    it needs, and the table still follows the figures.

    Each figure is centered in its cell at its own aspect ratio, and the table
    with its footnote is centered vertically in its own, so a table shorter than
    its cell sits beside the figures rather than hanging from the top of the
    grid.

    There is no bullet area: a grid this full has nowhere to put one, and the
    slides that carry bullets are the findings pages, which are STACKED. A
    rule that set both would lose its bullets here, so it should not.
    """
    cells = len(spec.figures) + (1 if spec.table is not None else 0)
    rows = max(1, -(-cells // 2))
    cell_width = int((SLIDE_WIDTH - 2 * MARGIN - GUTTER) / 2)
    cell_height = int((height - GUTTER * (rows - 1)) / rows)

    def cell_origin(index: int) -> tuple[int, int]:
        row, column = divmod(index, 2)
        return (int(MARGIN + column * (cell_width + GUTTER)),
                int(top + row * (cell_height + GUTTER)))

    for index, path in enumerate(spec.figures):
        left, cell_top = cell_origin(index)
        _add_figure_in(slide, path, left, cell_top, cell_width, cell_height)

    if spec.table is None:
        return

    left, cell_top = cell_origin(len(spec.figures))
    widths = _fitted_widths(spec.table, vocab, flag_style, cell_width)
    wanted, footnote_height = _table_heights(spec.table, vocab, widths)
    # The table gets the height it wants or the height there is, whichever is
    # smaller. pptx grows a row that cannot fit its text, so squeezing here
    # costs a little of the space between rows and never the numbers.
    table_height = min(wanted, cell_height - footnote_height)
    cell_top += int((cell_height - table_height - footnote_height) / 2)
    _add_table(slide, spec.table, vocab, cell_top, table_height,
               flag_style=flag_style, left=left, available=cell_width,
               widths=widths)
    _add_footnotes(slide, spec.table, int(cell_top + table_height), flag_style,
                   left=left, width=cell_width)


def _titled_table_height(table: Table, vocab: Vocabulary | None = None,
                         widths: list[int] | None = None) -> int:
    """How deep one titled table runs: its title, its rows, its footnotes."""
    rows, footnote = _table_heights(table, vocab, widths)
    return rows + footnote + (int(TABLE_TITLE_HEIGHT) if table.title else 0)


def _titled_table_width(table: Table, vocab: Vocabulary, flag_style: str) -> int:
    """How wide a titled table wants to be: its columns, or its title.

    The title is allowed to set the width, because it is given one line and
    wrapping it would run it into the header row underneath. A narrow tally
    under a long title is the case: three digits of column against six words of
    heading. The table then centers inside the column it was given, which
    `_add_table` already does for any table narrower than its box.
    """
    columns = sum(_column_widths(table, vocab, flag_style))
    if not table.title:
        return columns
    heading = _text_width(capitalize_first(table.title), TABLE_TITLE_PT)
    return max(columns, heading + int(TEXT_BOX_INSETS))


def _add_titled_table(slide, table: Table, vocab: Vocabulary, left: Emu,
                      top: Emu, width: Emu, height: Emu, flag_style: str) -> None:
    """One table in a column of its own: its title, itself, then its footnotes.

    The title is drawn here. Elsewhere a slide holds one table
    and the slide's own title names it; a slide holding three needs each of
    them to say what it is, and `Table.title` is already the sentence that says
    so, written where the table was built.
    """
    if table.title:
        box = slide.shapes.add_textbox(left, top, width, TABLE_TITLE_HEIGHT)
        frame = box.text_frame
        frame.word_wrap = True
        frame.text = capitalize_first(table.title)
        run = frame.paragraphs[0].runs[0]
        run.font.size = TABLE_TITLE_PT
        run.font.bold = True
        top = int(top + TABLE_TITLE_HEIGHT)
        height = int(height - TABLE_TITLE_HEIGHT)

    # Sized once, then used for both the drawing and the footnote's position.
    # These are the columns the table will actually have, so the header's line
    # count is known and the footnote goes under the table rather than into it.
    columns = _fitted_widths(table, vocab, flag_style, width)
    wanted, footnote_height = _table_heights(table, vocab, columns)
    # What the table NEEDS, never what is left over. pptx grows a row to fit
    # its text whatever height the shape claims, so a table clamped to the
    # space available does not get smaller, it just has its footnote drawn
    # inside it. If that runs past the foot of the slide, the geometry check
    # says so, which is the honest failure.
    table_height = wanted
    _add_table(slide, table, vocab, top, table_height, flag_style=flag_style,
               left=left, available=width, widths=columns)
    _add_footnotes(slide, table, int(top + table_height), flag_style,
                   left=left, width=width)


def _spec_tables(spec: Slide) -> list[Table]:
    """The tables a table layout draws: `tables`, or the single-table field.

    Falling back to `table` keeps these layouts supersets of the single-table
    ones rather than forks of them, so a rule that only ever has one table need
    not know which field this layout reads.
    """
    return list(spec.tables) or ([spec.table] if spec.table is not None else [])


def _row_widths(tables: list[Table], vocab: Vocabulary, flag_style: str,
                available: Emu) -> list[int]:
    """How wide each table in a row is drawn, after any squeeze.

    Shared by the drawing and by the depth calculation, because the two must
    agree: how deep a table runs depends on whether its headers wrapped, which
    depends on how much width the squeeze left it.
    """
    widths = [_titled_table_width(table, vocab, flag_style) for table in tables]
    gutters = GUTTER * (len(tables) - 1)
    if sum(widths) + gutters > available:
        room = available - gutters
        widths = [int(w * room / sum(widths)) for w in widths]
    return widths


def _row_depth(tables: list[Table], vocab: Vocabulary, flag_style: str,
               available: Emu) -> int:
    """How deep a row of titled tables runs, which is its deepest member."""
    widths = _row_widths(tables, vocab, flag_style, available)
    return max([_titled_table_height(
        table, vocab, _fitted_widths(table, vocab, flag_style, width))
        for table, width in zip(tables, widths)] or [0])


def _draw_table_row(slide, tables: list[Table], vocab: Vocabulary, left: Emu,
                    top: Emu, available: Emu, height: Emu, flag_style: str,
                    center: bool = True) -> None:
    """Titled tables side by side, each at its NATURAL width.

    Natural width rather than an equal share, because these tables are not
    peers in size: a one-column tally beside a five-column summary should look
    like what it is. They are only squeezed when the row genuinely does not
    fit, and then proportionally, so the wide one gives up the most.

    Tops are SHARED, so the titles line up and the row reads as a row. The
    tables then run to whatever depth they need, which differs, and that is
    fine: a reader compares them across, not down.

    `center` decides what happens to the width the row does not use. A slide
    made of tables centers, having nothing else to align to; the opening slide
    does not, because its tables sit under bullets set flush at the margin and
    a centered table under a flush-left bullet reads as a mistake.
    """
    widths = _row_widths(tables, vocab, flag_style, available)
    gutters = GUTTER * (len(tables) - 1)
    if center:
        left = int(left + (available - sum(widths) - gutters) / 2)

    for table, width in zip(tables, widths):
        _add_titled_table(slide, table, vocab, int(left), top, width, height,
                          flag_style)
        left += width + GUTTER


def _layout_tables(slide, spec: Slide, vocab: Vocabulary, top: Emu, height: Emu,
                   flag_style: str,
                   bullet_pt: Pt = BULLET_PT) -> None:
    """A slide made of tables: one row of them, centered on the page.

    One, two or three of them on the slides that use it today, and the row
    adapts by construction rather than by branching on the count: each table
    asks for what it needs, and what is left over is split evenly either side.
    A single table therefore arrives centered, which is what a slide carrying
    one table should look like.

    The row's shared top is placed so that the deepest table is centered
    vertically, which keeps the block off the ceiling without breaking the
    alignment that centering each table separately would.
    """
    tables = _spec_tables(spec)
    if spec.bullets:
        used = sum(bullet_height(line) for line in spec.bullets)
        _add_bullets(slide, spec.bullets, top, int(used))
        top, height = int(top + used + GUTTER), int(height - used - GUTTER)
    if not tables:
        return

    available = SLIDE_WIDTH - 2 * MARGIN
    depth = _row_depth(tables, vocab, flag_style, available)
    _draw_table_row(slide, tables, vocab, MARGIN,
                    int(top + max(0, height - depth) / 2),
                    available, min(int(height), depth), flag_style)


def _layout_title(slide, spec: Slide, vocab: Vocabulary, top: Emu, height: Emu,
                  flag_style: str,
                  bullet_pt: Pt = BULLET_PT) -> None:
    """The opening slide: what the deck is, and what it was computed over.

    Bullets, then any tables beneath them, the whole block centered in the body
    rather than hung from the top. That centering is the only thing that
    distinguishes this from a STACKED slide carrying the same pieces, and it is
    the reason it is worth being its own layout: a title slide sitting flush
    against its title looks like a slide that lost its figure.

    Everything is set flush left at the margin, bullets and tables alike, so
    the slide reads as one column of type under the deck's name.

    A figure moves that column into the left half and takes the right, rather
    than sitting under the text as STACKED would put it. Under the text it
    would get whatever the bullets and the window table left over, which on a
    multi-period window is not enough to read a diagram in; beside them it is
    limited by width, and both halves are centered in their own, so neither
    reads as hanging off the other.
    """
    tables = _spec_tables(spec)
    text_width = int(SLIDE_WIDTH - 2 * MARGIN)
    if spec.figures:
        text_width = int((text_width - GUTTER) / 2)
        _add_figure_row(slide, spec.figures,
                        int(SLIDE_WIDTH - MARGIN - text_width), top,
                        text_width, height)

    bullets = sum(bullet_height(line, text_width, bullet_pt)
                  for line in spec.bullets)
    depth = _row_depth(tables, vocab, flag_style, text_width)
    gap = int(GUTTER) if spec.bullets and tables else 0

    y = int(top + max(0, height - bullets - gap - depth) / 2)
    if spec.bullets:
        _add_bullets(slide, spec.bullets, y, int(bullets), width=text_width,
                     size=bullet_pt)
        y += int(bullets) + gap
    if tables:
        _draw_table_row(slide, tables, vocab, MARGIN, y, text_width, depth,
                        flag_style, center=False)


LAYOUT_FUNCTIONS = {
    "STACKED": _layout_stacked,
    "SPLIT": _layout_split,
    "QUADRANTS": _layout_quadrants,
    "TABLES": _layout_tables,
    "TITLE": _layout_title,
}


# What the two note sections are headed, and what they are headed WITH: the
# titles of the closing slides that gather the same sentences. A presenter who
# reads "Findings" at the foot of slide 9's notes and then arrives at a slide
# called Findings has been told these are the same thing, which they are.
#
# The sentences are repeated rather than moved. They were always meant for the
# closing sections; what was missing is that the presenter had no way to see,
# while talking through a slide, what that slide had contributed, and no way to
# quote the number without reading it off the table and rounding it themselves.
NOTES_SECTIONS = (("findings", "FINDINGS FROM THIS SLIDE:"),
                  ("recommendations", "WHERE TO LOOK NEXT:"))


def _notes_sections(spec: Slide) -> list[str]:
    """The findings and recommendations, as trailing note paragraphs.

    One paragraph per section, its sentences run together, rather than a
    paragraph each: the notes pane of a slide is read in a hurry from a lectern,
    and a heading that owns its sentences is easier to find there than a heading
    followed by an indeterminate number of peers.

    A section with nothing in it prints nothing, heading included. Most slides
    carry findings, several carry no recommendation, and a bare heading over
    nothing would read as a rule that failed rather than as a rule that had
    nothing to say.
    """
    sections = []
    for field_name, heading in NOTES_SECTIONS:
        lines = getattr(spec, field_name)
        if lines:
            sections.append(f"{heading} " + " ".join(lines))
    return sections


def _text_column(spec: Slide) -> int:
    """The width this slide's text is set in, which decides how far it wraps.

    Full width, except where a figure takes the other half: TITLE puts its
    schematic beside the text rather than under it, and the same sentence wraps
    to twice the lines in half the column.
    """
    width = int(SLIDE_WIDTH - 2 * MARGIN)
    if spec.layout == "TITLE" and spec.figures:
        return int((width - GUTTER) / 2)
    return width


def _body_depth(spec: Slide, vocab: Vocabulary, flag_style: str,
                bullet_pt: Pt = BULLET_PT) -> int:
    """How much height this slide's body needs, measured before it is drawn.

    Asked so that a slide whose content is taller than the band a template
    leaves free keeps the whole slide and goes without the artwork, rather than
    running its last table row over a strip of logos. Measured with the same
    routines the layouts use, in the column they will set the text in, so the
    answer is the one they will reach.

    A figure costs no height here. Every slide that reaches this either carries
    none, or carries a schematic beside the text, which takes width rather than
    height and shrinks to whatever is left.
    """
    column = _text_column(spec)
    depth = sum(bullet_height(line, column, bullet_pt) for line in spec.bullets)
    if spec.tables:
        depth += _row_depth(spec.tables, vocab, flag_style, column)
    if spec.table is not None:
        widths = _fitted_widths(spec.table, vocab, flag_style, column)
        depth += sum(_table_heights(spec.table, vocab, widths))
    return int(depth)


def _fitting_size(spec: Slide, vocab: Vocabulary, flag_style: str,
                  room: int) -> Pt | None:
    """The largest bullet size this slide fits the band at, or None.

    Tried in order, so a slide takes the artwork at the size it reads best in
    and steps down only as far as it has to. None is a slide that does not fit
    at any of them, which keeps the whole plain page instead.
    """
    for size in BULLET_SIZES:
        spent = (int(TITLE_HEIGHT) + _body_depth(spec, vocab, flag_style, size)
                 + (lead_height(spec.lead) if spec.lead else 0)
                 + (int(FOOTNOTE_HEIGHT) if spec.footnote else 0))
        if spent <= room:
            return size
    return None


def render(slides: list[Slide], path: str | Path, vocab: Vocabulary,
           flag_style: str = "MARK",
           template: str | Path | None = None) -> Path:
    """Write a deck. Blank layout throughout: this file owns the geometry.

    `template` is a one-slide .pptx whose artwork is stamped onto the slides
    that have room for it and whose theme the whole deck is set in. Without
    one the deck is built on pptx's default template, as it always was.
    """
    if template is None:
        deck = Presentation()
        deck.slide_width = SLIDE_WIDTH
        deck.slide_height = SLIDE_HEIGHT
        decoration = None
    else:
        # The template's own page is left as it was written, having been
        # checked against this renderer's. Setting it here instead would move
        # the artwork by whatever the two rounded differently.
        deck = Presentation(str(template))
        decoration = _read_decoration(deck, template)
    blank = _blank_layout(deck)

    for spec in slides:
        slide = deck.slides.add_slide(blank)
        # The plain slide's own band when there is no template, or when this
        # slide has no room for one, so everything below reads one pair of
        # numbers and no layout has to know a template is in play.
        decorated, bullet_pt = Decoration(), BULLET_PT
        if decoration is not None and takes_decoration(spec):
            size = _fitting_size(spec, vocab, flag_style,
                                 decoration.bottom - decoration.top)
            if size is not None:
                _stamp(slide, decoration)
                decorated, bullet_pt = decoration, size
        artwork = len(slide.shapes)
        _add_title(slide, spec.title,
                   OPENING_TITLE_PT if spec.layout == "TITLE" else TITLE_PT,
                   decorated.top)

        body_top = decorated.top + int(TITLE_HEIGHT)
        body_height = decorated.bottom - body_top
        if spec.lead:
            # Drawn first and then taken out of the body, the same bargain the
            # footnote strikes: the layout is told what room is left rather
            # than the lead being written over what the layout put there.
            lead = lead_height(spec.lead)
            _add_lead(slide, spec.lead, body_top, lead)
            body_top += lead
            body_height -= lead
        if spec.footnote:
            # The layout is told about the space the footnote takes rather than
            # the footnote being drawn over whatever the layout put there.
            body_height -= FOOTNOTE_HEIGHT
            _add_slide_footnote(slide, spec.footnote, decorated.bottom)
        LAYOUT_FUNCTIONS[spec.layout](
            slide, spec, vocab, body_top, int(body_height), flag_style,
            bullet_pt)
        _set_font(slide, artwork)

        paragraphs = list(spec.notes) + _notes_sections(spec)
        if paragraphs:
            frame = slide.notes_slide.notes_text_frame
            frame.text = paragraphs[0]
            for note in paragraphs[1:]:
                frame.add_paragraph().text = note

    if template is not None:
        _drop_template_slide(deck)
    path = Path(path)
    deck.save(str(path))
    return path
